"""Deterministic slide renderer — converts layout-input.json + layout-specs.json to PPTX.

This module is the app's structured PPTX renderer. It reads slide data,
precomputed layout specs, theme/style settings, and slide assets, then writes
the deck directly.

The renderer uses the shared utility functions already defined in
``pptx-python-runner.py`` (``fetch_icon``, ``safe_add_picture``,
``ensure_contrast``, etc.) so rendering, repair, and validation all operate on
the same runtime contract.

Usage (standalone test)::

    python slide_renderer.py \\
        --layout-input  previews/layout-input.json \\
        --layout-specs  previews/layout-specs.json \\
        --output        previews/presentation-preview.pptx \\
        --workspace-dir /path/to/workspace

Usage (from pptx-python-runner.py)::

    from slide_renderer import render_presentation
    render_presentation(layout_input, layout_specs, theme, style, ...)
"""

from __future__ import annotations

import json
import math
import os
import sys
from dataclasses import dataclass, replace
from pathlib import Path

if __package__ in {None, ''}:
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from scripts.layout.layout_specs import (  # noqa: E402
    LayoutSpec,
    RectSpec,
    SLIDE_WIDTH_IN,
    SLIDE_HEIGHT_IN,
    estimate_text_height_in,
)
from scripts.style_config import StyleConfig  # noqa: E402

# ── Constants ────────────────────────────────────────────────────────
EMU_PER_INCH = 914400


# ── Render context ───────────────────────────────────────────────────

@dataclass
class RenderContext:
    """Shared state threaded through all render functions."""

    prs: object  # Presentation
    theme: dict[str, str]
    style: StyleConfig
    font_family: str
    slide_assets: list[dict]
    accent_cycle: list[str]
    template_path: str | None
    workspace_dir: str
    theme_explicit: bool  # True when user provided custom theme colors
    text_box_style: str
    show_slide_icons: bool

    # ── Utilities (injected by caller) ───────────────────────────
    # These are references to the same functions from pptx-python-runner.py
    rgb_color: object  # Callable[[str, str], RGBColor]
    ensure_contrast: object  # Callable[[str, str, float], str]
    set_fill_transparency: object  # Callable[[shape, float], None]
    apply_gradient_fill: object  # Callable[[shape, list[str], int], None]
    fetch_icon: object  # Callable[[str, str, int, str | None], str | None]
    safe_add_picture: object  # Callable[..., object | None]
    safe_add_design_picture: object  # Callable[..., object | None]
    add_design_shape: object  # Callable[..., object]
    add_managed_textbox: object  # Callable[..., object]
    add_managed_shape: object  # Callable[..., object]
    tag_as_design: object  # Callable[[object, str], None]
    resolve_font: object  # Callable[[str, str], str]
    get_blank_layout: object  # Callable[[object], object]
    apply_widescreen: object  # Callable[[object], object]
    slide_image_paths: object  # Callable[[int], list[str]]
    slide_icon_name: object  # Callable[[int], str | None]
    slide_icon_collection: object  # Callable[[int], str | None]
    ensure_parent_dir: object  # Callable[[str], None]
    safe_image_path: object  # Callable[[str], str | None]


# ── Color helpers ────────────────────────────────────────────────────

def _build_accent_cycle(theme: dict[str, str]) -> list[str]:
    keys = ["ACCENT1", "ACCENT3", "ACCENT4", "ACCENT5", "ACCENT6", "ACCENT2"]
    cycle: list[str] = []
    for k in keys:
        v = theme.get(k, "")
        if isinstance(v, str) and v.strip():
            cycle.append(v.strip().lstrip("#").upper())
    if not cycle:
        cycle = ["4472C4", "70AD47", "FFC000", "ED7D31", "5B9BD5", "A5A5A5"]
    return cycle


def _hex_color(theme: dict[str, str], *keys: str) -> str:
    for key in keys:
        val = theme.get(key)
        if isinstance(val, str) and val.strip():
            return val.strip().lstrip("#").upper()
    return "000000"


def _build_colors(theme: dict[str, str]) -> dict[str, str]:
    return {
        "BG": _hex_color(theme, "BG", "LT1", "LIGHT", "WHITE"),
        "TEXT": _hex_color(theme, "TEXT", "DK1", "DARK"),
        "DARK": _hex_color(theme, "DARK", "DK1", "TEXT"),
        "DARK2": _hex_color(theme, "DARK2", "DK2", "BORDER", "TEXT"),
        "LIGHT": _hex_color(theme, "LIGHT", "LT1", "WHITE", "BG"),
        "LIGHT2": _hex_color(theme, "LIGHT2", "LT2", "BORDER", "BG"),
        "SECONDARY": _hex_color(theme, "SECONDARY", "ACCENT2", "LT2"),
        "BORDER": _hex_color(theme, "BORDER", "SECONDARY", "ACCENT2"),
        "ACCENT1": _hex_color(theme, "ACCENT1", "PRIMARY"),
        "ACCENT2": _hex_color(theme, "ACCENT2", "SECONDARY"),
        "ACCENT3": _hex_color(theme, "ACCENT3", "ACCENT1"),
        "ACCENT4": _hex_color(theme, "ACCENT4", "ACCENT2"),
        "ACCENT5": _hex_color(theme, "ACCENT5", "ACCENT3"),
        "ACCENT6": _hex_color(theme, "ACCENT6", "ACCENT4"),
        "PRIMARY": _hex_color(theme, "PRIMARY", "ACCENT1"),
    }


def _mix_hex(a_hex: str, b_hex: str, weight_b: float) -> str:
    """Return a blended hex color with ``weight_b`` contribution from ``b_hex``."""
    weight_b = max(0.0, min(weight_b, 1.0))
    weight_a = 1.0 - weight_b
    a_hex = (a_hex or "000000").strip().lstrip("#")[:6].ljust(6, "0")
    b_hex = (b_hex or "000000").strip().lstrip("#")[:6].ljust(6, "0")
    comps = []
    for idx in (0, 2, 4):
        a_val = int(a_hex[idx:idx + 2], 16)
        b_val = int(b_hex[idx:idx + 2], 16)
        mixed = round(a_val * weight_a + b_val * weight_b)
        comps.append(f"{max(0, min(mixed, 255)):02X}")
    return "".join(comps)


def _luminance(hex_color: str) -> float:
    """Return relative luminance (0..1) from a hex color string."""
    h = (hex_color or "000000").strip().lstrip("#")[:6].ljust(6, "0")
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _effective_panel_bg(ctx: RenderContext, fill_hex: str, colors: dict[str, str]) -> str:
    """Return the colour text will actually sit on after frosted/tinted mixing.

    During rendering, frosted panels mix ``fill_hex`` toward white (or dark),
    and tinted panels use ``fill_hex`` at reduced opacity over the slide BG.
    Contrast must be checked against this effective colour, not ``fill_hex``.
    """
    style = ctx.style
    if style.panel_fill == "frosted":
        if style.dark_mode:
            return _mix_hex(fill_hex, colors.get("DARK", "1B1B1B"), 0.45)
        return _mix_hex(fill_hex, colors.get("LIGHT", "FFFFFF"), 0.55)
    if style.panel_fill == "tinted":
        # Tinted panels are fill_hex at panel_fill_opacity over the slide BG.
        # Approximate the blended visual result.
        opacity = style.panel_fill_opacity
        return _mix_hex(colors["BG"], fill_hex, opacity)
    if style.panel_fill == "transparent":
        return colors["BG"]
    # solid — fill_hex as-is
    return fill_hex


def _apply_slide_bg_gradient(
    slide, color_stops: tuple[str, ...], angle_degrees: float = 90.0,
) -> None:
    """Apply a linear gradient background to a slide via DrawingML XML."""
    from lxml import etree

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # Ensure the background XML structure exists
    slide.background.fill.solid()

    bg_pr = slide.background._element.find(f"{{{ns_p}}}bg/{{{ns_p}}}bgPr")
    if bg_pr is None:
        return

    # Remove existing fill elements
    for tag in ("solidFill", "gradFill", "pattFill", "blipFill", "noFill"):
        for el in list(bg_pr.findall(f"{{{ns_a}}}{tag}")):
            bg_pr.remove(el)

    # Build gradient fill
    grad = etree.Element(f"{{{ns_a}}}gradFill")
    gs_lst = etree.SubElement(grad, f"{{{ns_a}}}gsLst")
    n = max(len(color_stops) - 1, 1)
    for i, c in enumerate(color_stops):
        gs = etree.SubElement(gs_lst, f"{{{ns_a}}}gs")
        gs.set("pos", str(int(round(i * 100000 / n))))
        srgb = etree.SubElement(gs, f"{{{ns_a}}}srgbClr")
        srgb.set("val", c.lstrip("#").upper())
    lin = etree.SubElement(grad, f"{{{ns_a}}}lin")
    lin.set("ang", str(int(round((angle_degrees % 360) * 60000))))
    lin.set("scaled", "1")

    # Insert before effectLst (must come first per OOXML schema)
    bg_pr.insert(0, grad)


def _resolve_slide_colors(
    ctx: RenderContext,
    base_colors: dict[str, str],
    accent_a: str,
    accent_b: str,
    slide_index: int,
) -> dict[str, str]:
    """Derive per-slide colour roles from the user's palette + the active style.

    The palette BG is used as-is for the slide background — no accent mixing.
    Panels, borders, and secondary tones are derived from palette slots with
    accent colouring so they complement the style.
    """
    style = ctx.style
    colors = dict(base_colors)

    # Background: use style fallback when the user has not provided explicit
    # theme colors and the style defines its own characteristic background.
    # First color in bg_colors is the representative solid for contrast math.
    bg_hex = colors["BG"]
    if style.bg_colors:
        bg_hex = style.bg_colors[0]

    if style.dark_mode:
        panel_base = _mix_hex(colors["DARK2"], accent_a, 0.28)
        panel_alt = _mix_hex(colors["DARK2"], accent_b, 0.24)
        border_hex = _mix_hex(colors["LIGHT2"], accent_a, 0.40)
        text_hex = ctx.ensure_contrast(colors["TEXT"], bg_hex)
        secondary_hex = ctx.ensure_contrast(colors["LIGHT2"], bg_hex)
    else:
        panel_base = _mix_hex(colors["LIGHT"], accent_a, 0.28)
        panel_alt = _mix_hex(colors["LIGHT"], accent_b, 0.22)
        border_hex = _mix_hex(colors["BORDER"], accent_a, 0.35)
        text_hex = ctx.ensure_contrast(colors["TEXT"], bg_hex)
        secondary_hex = ctx.ensure_contrast(colors["TEXT"], panel_base)

    # Style-specific panel fill overrides
    if style.panel_fill == "solid":
        panel_base = accent_a
        panel_alt = accent_b if slide_index % 2 == 0 else _mix_hex(accent_b, colors["LIGHT"], 0.20)
    elif style.panel_fill == "frosted":
        if style.dark_mode:
            panel_base = _mix_hex(colors["DARK2"], accent_a, 0.32)
            panel_alt = _mix_hex(colors["DARK2"], accent_b, 0.28)
        else:
            panel_base = _mix_hex(bg_hex, accent_a, 0.25)
            panel_alt = _mix_hex(bg_hex, accent_b, 0.20)
    elif style.panel_fill == "tinted":
        # Tinted panels should be clearly coloured — visible accent wash
        if style.dark_mode:
            panel_base = _mix_hex(colors["DARK2"], accent_a, 0.40)
            panel_alt = _mix_hex(colors["DARK2"], accent_b, 0.35)
        else:
            panel_base = _mix_hex(bg_hex, accent_a, 0.38)
            panel_alt = _mix_hex(bg_hex, accent_b, 0.30)
    elif style.panel_fill == "transparent":
        panel_base = _mix_hex(bg_hex, accent_a, 0.15)
        panel_alt = _mix_hex(bg_hex, accent_b, 0.12)

    colors["BG"] = bg_hex
    colors["TEXT"] = text_hex
    colors["SECONDARY"] = secondary_hex
    colors["BORDER"] = border_hex
    colors["PANEL_BASE"] = panel_base
    colors["PANEL_ALT"] = panel_alt
    return colors


# ── Low-level PPTX writers ──────────────────────────────────────────

def _write_run(paragraph, text: str, font_size_pt: float, color_hex: str,
               bold: bool, font_name: str) -> object:
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color_hex)
    run.font.name = font_name
    return run


def _add_textbox(ctx: RenderContext, slide, rect: RectSpec, text: str,
                 font_size_pt: float, color_hex: str, *,
                 bold: bool = False, name: str = "",
                 align: int = PP_ALIGN.LEFT, line_spacing: float = 0.0) -> object:
    """Add a textbox with NO fill (transparent by default)."""
    tb = ctx.add_managed_textbox(
        slide.shapes,
        Inches(rect.x), Inches(rect.y),
        Inches(rect.w), Inches(rect.h),
        name=name,
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    # NO fill — textboxes are transparent by default

    p = tf.paragraphs[0]
    p.alignment = align
    density = ctx.style.content_density
    if density == "compact":
        spacing_mult = 0.94
    elif density == "spacious":
        spacing_mult = 1.08
    else:
        spacing_mult = 1.0
    if line_spacing:
        p.line_spacing = line_spacing * spacing_mult
    else:
        base_spacing = 1.12 if font_size_pt >= 24 else 1.32
        p.line_spacing = base_spacing * spacing_mult

    _write_run(p, text, font_size_pt, color_hex, bold, ctx.font_family)
    return tb


def _add_panel(ctx: RenderContext, slide, rect: RectSpec,
               fill_hex: str, border_hex: str, *,
               accent_b_hex: str = "", name: str = "") -> object:
    """Add a panel shape with fill controlled by StyleConfig."""
    style = ctx.style
    panel_shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if style.text_box_corner_style == "rounded"
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    colors = _build_colors(ctx.theme)
    if style.panel_shadow != "none" and style.panel_fill != "transparent":
        shadow_dx = 0.10 if style.panel_shadow == "hard" else 0.06
        shadow_dy = 0.11 if style.panel_shadow == "hard" else 0.08
        shadow = ctx.add_design_shape(
            slide.shapes,
            panel_shape_type,
            Inches(rect.x + shadow_dx), Inches(rect.y + shadow_dy),
            Inches(rect.w), Inches(rect.h),
            name=f"{name}_shadow" if name else "panel_shadow",
        )
        if style.panel_shadow == "hard":
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = ctx.rgb_color(colors.get("DARK", "111111"))
        else:
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = ctx.rgb_color(accent_b_hex or colors.get("ACCENT2", fill_hex))
            ctx.set_fill_transparency(shadow, 0.32 if style.dark_mode else 0.22)
        shadow.line.fill.background()
    shape = ctx.add_managed_shape(
        slide.shapes,
        panel_shape_type,
        Inches(rect.x), Inches(rect.y),
        Inches(rect.w), Inches(rect.h),
        name=name,
    )

    if style.panel_fill == "transparent":
        shape.fill.background()
    elif style.panel_fill == "frosted":
        # Frosted glass: semi-transparent accent-tinted fill
        frost_color = _mix_hex(fill_hex, colors.get("LIGHT", "FFFFFF"), 0.55) if not style.dark_mode else _mix_hex(fill_hex, colors.get("DARK", "1B1B1B"), 0.45)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.rgb_color(frost_color)
        ctx.set_fill_transparency(shape, 1.0 - style.panel_fill_opacity)
    elif style.panel_fill == "tinted":
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.rgb_color(fill_hex)
        ctx.set_fill_transparency(shape, 1.0 - style.panel_fill_opacity)
    elif style.panel_fill == "solid":
        if style.color_treatment == "gradient" and accent_b_hex:
            ctx.apply_gradient_fill(shape, [fill_hex, accent_b_hex],
                                    angle_degrees=style.gradient_angle)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = ctx.rgb_color(fill_hex)
        if style.panel_fill_opacity < 1.0:
            ctx.set_fill_transparency(shape, 1.0 - style.panel_fill_opacity)
    else:
        shape.fill.background()

    if style.panel_border:
        shape.line.fill.solid()
        shape.line.color.rgb = ctx.rgb_color(border_hex or fill_hex)
        shape.line.width = Pt(style.panel_border_weight_pt)
    else:
        shape.line.fill.background()

    return shape


def _add_panel_stripe(ctx: RenderContext, slide, rect: RectSpec,
                      color_hex: str, idx: int) -> object | None:
    """Add a thin vertical color stripe on the left edge of a panel."""
    if not ctx.style.panel_stripe:
        return None
    stripe_w = max(min(rect.w * 0.02, 0.08), 0.04)
    stripe = ctx.add_design_shape(
        slide.shapes,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(rect.x), Inches(rect.y),
        Inches(stripe_w), Inches(rect.h),
        name=f"stripe_{idx}",
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ctx.rgb_color(color_hex)
    stripe.line.fill.background()
    return stripe


# ── Design language (decorative accents) ─────────────────────────────

def _add_design_language(ctx: RenderContext, slide, spec: LayoutSpec,
                         accent_a: str, accent_b: str,
                         colors: dict[str, str]) -> None:
    """Add optional decorative elements controlled by StyleConfig."""
    style = ctx.style
    ref = spec.title_rect or spec.content_rect or spec.key_message_rect
    if ref is None:
        return

    if style.background_grid != "none":
        grid_color = _mix_hex(accent_a, colors["BG"], 0.25 if style.dark_mode else 0.45)
        if style.background_grid == "fine":
            x = 0.0
            while x <= SLIDE_WIDTH_IN:
                line = ctx.add_design_shape(
                    slide.shapes,
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(x), Inches(0),
                    Inches(0.01), Inches(SLIDE_HEIGHT_IN),
                    name="bg_grid_v",
                )
                line.fill.solid()
                line.fill.fore_color.rgb = ctx.rgb_color(grid_color)
                ctx.set_fill_transparency(line, 0.86 if style.dark_mode else 0.90)
                line.line.fill.background()
                x += 0.72
            y = 0.0
            while y <= SLIDE_HEIGHT_IN:
                line = ctx.add_design_shape(
                    slide.shapes,
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(0), Inches(y),
                    Inches(SLIDE_WIDTH_IN), Inches(0.01),
                    name="bg_grid_h",
                )
                line.fill.solid()
                line.fill.fore_color.rgb = ctx.rgb_color(grid_color)
                ctx.set_fill_transparency(line, 0.88 if style.dark_mode else 0.92)
                line.line.fill.background()
                y += 0.56
        elif style.background_grid == "perspective":
            horizon_y = SLIDE_HEIGHT_IN * 0.66
            rail_y = horizon_y
            rail_gap = 0.12
            while rail_y < SLIDE_HEIGHT_IN:
                rail = ctx.add_design_shape(
                    slide.shapes,
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(0), Inches(rail_y),
                    Inches(SLIDE_WIDTH_IN), Inches(0.014),
                    name="perspective_rail",
                )
                rail.fill.solid()
                rail.fill.fore_color.rgb = ctx.rgb_color(grid_color)
                ctx.set_fill_transparency(rail, 0.70 if style.dark_mode else 0.82)
                rail.line.fill.background()
                rail_y += rail_gap
                rail_gap *= 1.23
            vanish_x = SLIDE_WIDTH_IN / 2
            bottom_y = SLIDE_HEIGHT_IN
            for idx, base_x in enumerate([0.35, 1.4, 2.7, 4.2, 5.8, 7.2, 8.8, 10.3, 11.6, 12.5]):
                dx = base_x - vanish_x
                line_len = max(bottom_y - horizon_y, 1.6)
                line = ctx.add_design_shape(
                    slide.shapes,
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(base_x), Inches(horizon_y),
                    Inches(0.014), Inches(line_len),
                    name=f"perspective_ray_{idx}",
                )
                line.rotation = max(min(-dx * 7.5, 42), -42)
                line.fill.solid()
                line.fill.fore_color.rgb = ctx.rgb_color(grid_color)
                ctx.set_fill_transparency(line, 0.72 if style.dark_mode else 0.84)
                line.line.fill.background()

    if style.frame_outline != "none":
        frame_color = _mix_hex(accent_a, accent_b, 0.30)
        frame_specs = [(0.22, 0.22, SLIDE_WIDTH_IN - 0.44, SLIDE_HEIGHT_IN - 0.44, 1.2)]
        if style.frame_outline == "double":
            frame_specs.append((0.38, 0.38, SLIDE_WIDTH_IN - 0.76, SLIDE_HEIGHT_IN - 0.76, 0.8))
        for idx, (x, y, w, h, width_pt) in enumerate(frame_specs):
            frame = ctx.add_design_shape(
                slide.shapes,
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(w), Inches(h),
                name=f"frame_outline_{idx}",
            )
            frame.fill.background()
            frame.line.fill.solid()
            frame.line.color.rgb = ctx.rgb_color(frame_color)
            frame.line.width = Pt(width_pt)

    if style.corner_brackets:
        bracket_color = _mix_hex(accent_a, accent_b, 0.45)
        pad = 0.12
        arm = max(min(ref.w * 0.10, 0.42), 0.24)
        thick = 0.03
        corners = [
            (ref.x - pad, ref.y - pad, 1, 1),
            (ref.x + ref.w + pad, ref.y - pad, -1, 1),
            (ref.x - pad, ref.y + ref.h + pad, 1, -1),
            (ref.x + ref.w + pad, ref.y + ref.h + pad, -1, -1),
        ]
        for idx, (cx, cy, sx, sy) in enumerate(corners):
            h_bar = ctx.add_design_shape(
                slide.shapes,
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx if sx > 0 else cx - arm), Inches(cy if sy > 0 else cy - thick),
                Inches(arm), Inches(thick),
                name=f"corner_h_{idx}",
            )
            v_bar = ctx.add_design_shape(
                slide.shapes,
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cx if sx > 0 else cx - thick), Inches(cy if sy > 0 else cy - arm),
                Inches(thick), Inches(arm),
                name=f"corner_v_{idx}",
            )
            for part in (h_bar, v_bar):
                part.fill.solid()
                part.fill.fore_color.rgb = ctx.rgb_color(bracket_color)
                part.line.fill.background()

    if style.accent_rings:
        anchor = spec.icon_rect or spec.hero_rect or spec.sidebar_rect or ref
        ring_color = _mix_hex(accent_b, accent_a, 0.45)
        base_size = max(min(anchor.h * 0.28, 1.2), 0.62)
        ring_x = min(anchor.x + anchor.w - base_size * 0.50, SLIDE_WIDTH_IN - base_size - 0.14)
        ring_y = max(anchor.y - base_size * 0.08, 0.14)
        for idx, scale in enumerate((1.0, 0.72, 0.44)):
            size = base_size * scale
            ring = ctx.add_design_shape(
                slide.shapes,
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(ring_x + (base_size - size) / 2),
                Inches(ring_y + (base_size - size) / 2),
                Inches(size), Inches(size),
                name=f"accent_ring_{idx}",
            )
            ring.fill.background()
            ring.line.fill.solid()
            ring.line.color.rgb = ctx.rgb_color(ring_color)
            ring.line.width = Pt(1.6 if idx == 0 else 1.0)

    # Vertical accent bar — thick enough to be a significant visual feature
    if style.title_accent_bar:
        bottom_limit = (spec.notes_rect.y - 0.08) if spec.notes_rect else min(ref.y + ref.h + 4.8, SLIDE_HEIGHT_IN - 0.12)
        bar_w = max(min(ref.x * 0.28, 0.18), 0.10)
        bar_gap = max(min(ref.x * 0.10, 0.08), 0.04)
        bar_x = max(ref.x - bar_w - bar_gap, 0.02)
        bar_y = max(ref.y - min(ref.h * 0.15, 0.08), 0.0)
        bar_h = max(bottom_limit - bar_y, ref.h + 0.4)
        bar = ctx.add_design_shape(
            slide.shapes,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(bar_x), Inches(bar_y),
            Inches(bar_w), Inches(bar_h),
            name="accent_bar",
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ctx.rgb_color(accent_a)
        bar.line.fill.background()

    if style.decorative_blob:
        anchor = spec.content_rect or ref
        blob_w = min(max(anchor.w * 0.42, 2.8), 4.4)
        blob_h = min(max(anchor.h * 0.62, 1.9), 3.0)
        blob_x = min(anchor.x + anchor.w - blob_w * 0.55, SLIDE_WIDTH_IN - blob_w - 0.10)
        blob_y = max(ref.y - blob_h * 0.18, 0.10)
        blob = ctx.add_design_shape(
            slide.shapes,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(blob_x), Inches(blob_y),
            Inches(blob_w), Inches(blob_h),
            name="bg_blob",
        )
        if style.color_treatment == "gradient":
            ctx.apply_gradient_fill(blob, [_mix_hex(accent_a, "FFFFFF", 0.20), _mix_hex(accent_b, "FFFFFF", 0.40)],
                                    angle_degrees=style.gradient_angle)
        else:
            blob.fill.solid()
            blob.fill.fore_color.rgb = ctx.rgb_color(_mix_hex(accent_a, "FFFFFF", 0.45))
        ctx.set_fill_transparency(blob, 0.65 if style.dark_mode else 0.78)
        blob.line.fill.background()

    # Horizontal accent rule — bold, extends full width of content area
    if style.title_accent_rule and spec.accent_rect is not None:
        rule_h = max(spec.accent_rect.h, 0.06)  # at least 0.06" thick
        rule = ctx.add_design_shape(
            slide.shapes,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(spec.accent_rect.x), Inches(spec.accent_rect.y),
            Inches(spec.accent_rect.w), Inches(rule_h),
            name="accent_rule",
        )
        if style.color_treatment == "gradient":
            ctx.apply_gradient_fill(rule, [accent_a, accent_b],
                                    angle_degrees=style.gradient_angle)
        else:
            rule.fill.solid()
            rule.fill.fore_color.rgb = ctx.rgb_color(accent_a)
        rule.line.fill.background()

    # Decorative circle — large enough to be a visible design element
    if style.decorative_circle:
        anchor = spec.content_rect or ref
        circle_size = max(min(anchor.h * 0.28, 1.0), 0.55)
        circle_x = min(anchor.x + anchor.w - circle_size * 0.55, SLIDE_WIDTH_IN - circle_size - 0.06)
        circle_y = (
            min(spec.notes_rect.y - circle_size - 0.10, SLIDE_HEIGHT_IN - circle_size - 0.06)
            if spec.notes_rect else
            min(anchor.y + anchor.h - circle_size * 0.25, SLIDE_HEIGHT_IN - circle_size - 0.06)
        )
        circle = ctx.add_design_shape(
            slide.shapes,
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(circle_x), Inches(circle_y),
            Inches(circle_size), Inches(circle_size),
            name="decor_circle",
        )
        circle.fill.background()
        circle.line.fill.solid()
        circle.line.color.rgb = ctx.rgb_color(accent_b)
        circle.line.width = Pt(2.0)

    # Rainbow spectrum stripe bars — full-width bands at top and bottom
    if style.rainbow_stripe_bars:
        accent_keys = ["ACCENT1", "ACCENT2", "ACCENT3", "ACCENT4", "ACCENT5", "ACCENT6"]
        stripe_count = len(accent_keys)
        bar_height = 0.08  # inches per stripe
        segment_w = SLIDE_WIDTH_IN / stripe_count
        for pos, y_base in [("top", 0.0), ("bottom", SLIDE_HEIGHT_IN - bar_height)]:
            for i, key in enumerate(accent_keys):
                stripe = ctx.add_design_shape(
                    slide.shapes,
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(i * segment_w), Inches(y_base),
                    Inches(segment_w + 0.01), Inches(bar_height),
                    name=f"rainbow_{pos}_{i}",
                )
                stripe.fill.solid()
                stripe.fill.fore_color.rgb = ctx.rgb_color(colors.get(key, accent_a))
                stripe.line.fill.background()

    # Sparkle stars — small star shapes in corners
    if style.sparkle_stars:
        star_positions = [
            (0.20, 0.18, 0.28),   # top-left
            (SLIDE_WIDTH_IN - 0.52, 0.14, 0.32),  # top-right
            (0.30, SLIDE_HEIGHT_IN - 0.55, 0.24),  # bottom-left
            (SLIDE_WIDTH_IN - 0.44, SLIDE_HEIGHT_IN - 0.48, 0.26),  # bottom-right
        ]
        star_colors = [accent_a, accent_b,
                       colors.get("ACCENT3", accent_a),
                       colors.get("ACCENT4", accent_b)]
        for i, (sx, sy, sz) in enumerate(star_positions):
            star = ctx.add_design_shape(
                slide.shapes,
                MSO_AUTO_SHAPE_TYPE.STAR_4_POINT,
                Inches(sx), Inches(sy),
                Inches(sz), Inches(sz),
                name=f"sparkle_star_{i}",
            )
            star.fill.solid()
            star.fill.fore_color.rgb = ctx.rgb_color(star_colors[i % len(star_colors)])
            ctx.set_fill_transparency(star, 0.35)
            star.line.fill.background()

    # Scan-line overlay — thin horizontal lines across the slide
    if style.scan_lines:
        scan_color = colors.get("WHITE", "FFFFFF")
        y = 0.0
        idx = 0
        while y < SLIDE_HEIGHT_IN:
            line = ctx.add_design_shape(
                slide.shapes,
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0), Inches(y),
                Inches(SLIDE_WIDTH_IN), Inches(0.008),
                name=f"scan_line_{idx}",
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ctx.rgb_color(scan_color)
            ctx.set_fill_transparency(line, 0.92)
            line.line.fill.background()
            y += 0.16
            idx += 1


def _add_key_message_band(ctx: RenderContext, slide, spec: LayoutSpec,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> str | None:
    """Optionally add a semi-transparent band behind the key-message text."""
    if not ctx.style.key_message_band or spec.key_message_rect is None:
        return None
    r = spec.key_message_rect
    # Band is clipped exactly to key_message_rect — no padding into title zone
    band = ctx.add_design_shape(
        slide.shapes,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(r.x), Inches(r.y),
        Inches(r.w), Inches(r.h),
        name="key_band",
    )
    if ctx.style.color_treatment == "gradient":
        ctx.apply_gradient_fill(band, [accent_a, accent_b],
                                angle_degrees=ctx.style.gradient_angle)
    else:
        band.fill.solid()
        band.fill.fore_color.rgb = ctx.rgb_color(accent_a)
    ctx.set_fill_transparency(band, 1.0 - ctx.style.key_message_band_opacity)
    band.line.fill.background()
    return ctx.ensure_contrast(colors["TEXT"], accent_a)


# ── Font sizing helpers ──────────────────────────────────────────────

def _adjust_title_font(
    rect: RectSpec,
    text: str,
    base_pt: float,
    scale: float = 1.0,
    *,
    min_pt: float = 22,
) -> float:
    size = max(base_pt * scale, min_pt)
    required = estimate_text_height_in(text, rect.w, size, line_height=1.08)
    if required <= rect.h * 0.86:
        return size

    while size > min_pt and required > rect.h * 1.02:
        size = max(min_pt, size - 2)
        required = estimate_text_height_in(text, rect.w, size, line_height=1.08)

    if required > rect.h * 0.92 and size > min_pt:
        size = max(min_pt, size - 1)
    return size


def _adjust_body_font(width_in: float, height_in: float, text: str,
                      base_pt: float) -> float:
    required = estimate_text_height_in(text, width_in, base_pt, line_height=1.18)
    if required > height_in * 1.05:
        return max(11, base_pt - 1.2)
    elif required > height_in * 0.85:
        return max(11, base_pt - 0.6)
    return base_pt


def _is_metric_like_title(text: str) -> bool:
    stripped = (text or "").strip()
    if not stripped:
        return False
    words = stripped.split()
    digit_count = sum(ch.isdigit() for ch in stripped)
    alpha_count = sum(ch.isalpha() for ch in stripped)
    if digit_count == 0:
        return len(words) <= 2 and len(stripped) <= 12
    return len(words) <= 4 and digit_count >= max(1, alpha_count // 6)


def _reflow_big_number_spec(
    spec: LayoutSpec,
    title_text: str,
    key_message_text: str,
    *,
    title_pt: float,
    key_pt: float,
) -> LayoutSpec:
    if spec.title_rect is None:
        return spec

    title_height = max(
        min(
            estimate_text_height_in(title_text, spec.title_rect.w, title_pt, line_height=1.08) + 0.08,
            spec.title_rect.h,
        ),
        0.9,
    )
    title_rect = replace(spec.title_rect, h=round(title_height, 4))

    next_y = round(title_rect.y + title_rect.h + 0.12, 4)
    key_rect = spec.key_message_rect
    if key_rect is not None:
        key_height = key_rect.h
        if key_message_text.strip():
            key_height = max(
                estimate_text_height_in(key_message_text, key_rect.w, key_pt, line_height=1.16) + 0.06,
                0.4,
            )
        key_rect = replace(key_rect, y=next_y, h=round(key_height, 4))
        next_y = round(key_rect.y + key_rect.h + 0.18, 4)

    content_rect = spec.content_rect
    if content_rect is not None:
        notes_top = spec.notes_rect.y if spec.notes_rect is not None else (SLIDE_HEIGHT_IN - 0.62)
        content_bottom = max(notes_top - 0.18, next_y + 0.8)
        content_rect = replace(
            content_rect,
            y=next_y,
            h=round(max(content_bottom - next_y, 0.8), 4),
        )

    return replace(spec, title_rect=title_rect, key_message_rect=key_rect, content_rect=content_rect)


def _density_gap_multiplier(density: str) -> float:
    if density == "compact":
        return 0.82
    if density == "spacious":
        return 1.24
    return 1.0


def _whitespace_gap_multiplier(bias: str) -> float:
    """Additional gap scaling from ``StyleLayoutPolicy.whitespace_bias``."""
    if bias == "tight":
        return 0.85
    if bias == "generous":
        return 1.15
    if bias == "editorial":
        return 1.20
    return 1.0


def _density_line_spacing_multiplier(density: str) -> float:
    if density == "compact":
        return 0.94
    if density == "spacious":
        return 1.08
    return 1.0


def _density_padding_adjustment(density: str) -> float:
    if density == "compact":
        return -0.03
    if density == "spacious":
        return 0.04
    return 0.0


def _density_font_adjustment(density: str) -> float:
    if density == "compact":
        return -0.6
    if density == "spacious":
        return 0.5
    return 0.0



# ── Panel text writers ───────────────────────────────────────────────

def _write_panel_text(ctx: RenderContext, shape, title_text: str, body_text: str,
                      fill_hex: str, colors: dict[str, str], *,
                      title_pt: float = 16, body_pt: float = 13,
                      top_reserve: float = 0.0, left_reserve: float = 0.0,
                      align: int = PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    density_name = ctx.style.content_density
    pad = _density_padding_adjustment(density_name)
    spacing_mult = _density_line_spacing_multiplier(density_name)
    title_pt = max(12, title_pt + _density_font_adjustment(density_name))
    body_pt = max(11, body_pt + _density_font_adjustment(density_name))
    left_margin = max(0.10, 0.16 + left_reserve + pad)
    right_margin = max(0.08, 0.14 + pad)
    top_margin = max(0.08, 0.12 + top_reserve + pad)
    bottom_margin = max(0.06, 0.10 + pad)
    tf.margin_left = Inches(left_margin)
    tf.margin_right = Inches(right_margin)
    tf.margin_top = Inches(top_margin)
    tf.margin_bottom = Inches(bottom_margin)

    usable_w = max(shape.width / EMU_PER_INCH - (left_margin + right_margin), 0.5)
    usable_h = max(shape.height / EMU_PER_INCH - (top_margin + bottom_margin), 0.45)
    title_need = estimate_text_height_in(title_text, usable_w, title_pt, line_height=1.10) if title_text else 0.0
    body_need = estimate_text_height_in(body_text, usable_w, body_pt, line_height=1.16) if body_text else 0.0
    density_ratio = (title_need + body_need + (0.06 if title_text and body_text else 0.0)) / usable_h

    line_spacing = 1.22 * spacing_mult
    if density_ratio > 1.05:
        title_pt = max(12, title_pt - 1.6)
        body_pt = max(11, body_pt - 1.2)
        line_spacing = 1.16 * spacing_mult
    elif density_ratio > 0.85:
        title_pt = max(12, title_pt - 0.8)
        body_pt = max(11, body_pt - 0.5)
        line_spacing = 1.18 * spacing_mult

    # Effective background for contrast check — use the actual rendered fill
    bg_for_contrast = _effective_panel_bg(ctx, fill_hex, colors)
    title_color = ctx.ensure_contrast(colors["TEXT"], bg_for_contrast)
    body_color = ctx.ensure_contrast(colors["TEXT"], bg_for_contrast)

    p0 = tf.paragraphs[0]
    p0.alignment = align
    p0.line_spacing = line_spacing
    _write_run(p0, title_text, title_pt, title_color, bold=True, font_name=ctx.font_family)

    if body_text:
        p1 = tf.add_paragraph()
        p1.alignment = align
        p1.line_spacing = line_spacing
        p1.space_before = Pt(4)
        _write_run(p1, body_text, body_pt, body_color, bold=False, font_name=ctx.font_family)


def _write_bullets_panel(ctx: RenderContext, shape, heading: str, items: list[str],
                         fill_hex: str, colors: dict[str, str], *,
                         heading_pt: float = 16, bullet_pt: float = 13,
                         top_reserve: float = 0.0) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    density_name = ctx.style.content_density
    pad = _density_padding_adjustment(density_name)
    spacing_mult = _density_line_spacing_multiplier(density_name)
    heading_pt = max(12, heading_pt + _density_font_adjustment(density_name))
    bullet_pt = max(11, bullet_pt + _density_font_adjustment(density_name))
    left_margin = max(0.10, 0.16 + pad)
    right_margin = max(0.08, 0.14 + pad)
    top_margin = max(0.08, 0.12 + top_reserve + pad)
    bottom_margin = max(0.06, 0.10 + pad)
    tf.margin_left = Inches(left_margin)
    tf.margin_right = Inches(right_margin)
    tf.margin_top = Inches(top_margin)
    tf.margin_bottom = Inches(bottom_margin)

    usable_w = max(shape.width / EMU_PER_INCH - (left_margin + right_margin), 0.5)
    usable_h = max(shape.height / EMU_PER_INCH - (top_margin + bottom_margin), 0.5)
    total_need = estimate_text_height_in(heading, usable_w, heading_pt, line_height=1.10)
    total_need += sum(
        estimate_text_height_in(item, usable_w - 0.18, bullet_pt, line_height=1.16)
        for item in items
    )
    density_ratio = total_need / usable_h
    line_spacing = 1.20 * spacing_mult
    if density_ratio > 1.05:
        heading_pt = max(12, heading_pt - 1.4)
        bullet_pt = max(11, bullet_pt - 1.2)
        line_spacing = 1.16 * spacing_mult
    elif density_ratio > 0.85:
        heading_pt = max(12, heading_pt - 0.7)
        bullet_pt = max(11, bullet_pt - 0.4)
        line_spacing = 1.18 * spacing_mult

    bg_for_contrast = _effective_panel_bg(ctx, fill_hex, colors)
    title_color = ctx.ensure_contrast(colors["TEXT"], bg_for_contrast)
    bullet_color = ctx.ensure_contrast(colors["TEXT"], bg_for_contrast)

    marker = ctx.style.bullet_marker
    marker_bold = ctx.style.bullet_marker_bold

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.line_spacing = line_spacing
    _write_run(p0, heading, heading_pt, title_color, bold=True, font_name=ctx.font_family)

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_before = Pt(3)
        _write_run(p, f"{marker} ", bullet_pt, bullet_color, bold=marker_bold, font_name=ctx.font_family)
        _write_run(p, item, bullet_pt, bullet_color, bold=False, font_name=ctx.font_family)


# ── Geometry helpers ─────────────────────────────────────────────────

def _split_content_for_images(rect: RectSpec) -> tuple[RectSpec, RectSpec]:
    gap = min(rect.w * 0.035, 0.18)
    image_w = min(max(rect.w * 0.36, 2.4), rect.w * 0.42)
    text_w = rect.w - image_w - gap
    text_rect = RectSpec(rect.x, rect.y, text_w, rect.h)
    image_rect = RectSpec(rect.x + text_w + gap, rect.y, image_w, rect.h)
    return text_rect, image_rect


def _grid_rects(rect: RectSpec, count: int, density: str = "normal",
                whitespace_bias: str = "normal",
                corner_style: str = "square") -> list[RectSpec]:
    if count <= 0:
        return []
    use_two_cols = count >= 5 or (rect.h / max(count, 1)) < 0.82
    cols = 2 if use_two_cols and count > 1 else 1
    rows = int(math.ceil(count / cols))
    gap_mult = _density_gap_multiplier(density) * _whitespace_gap_multiplier(whitespace_bias)
    gap_x = min(rect.w * 0.035, 0.18) * gap_mult
    gap_y = min(rect.h * 0.04, 0.18) * gap_mult
    # Rounded panels lose usable area to corner insets — reclaim space from
    # inter-panel gaps so each cell is taller within the same total zone.
    if corner_style == "rounded" and rows > 1:
        gap_y *= 0.5
    cell_w = rect.w if cols == 1 else (rect.w - gap_x * (cols - 1)) / cols
    cell_h = (rect.h - gap_y * (rows - 1)) / rows
    rects: list[RectSpec] = []
    for idx in range(count):
        col = idx % cols
        row = idx // cols
        rects.append(RectSpec(
            rect.x + col * (cell_w + gap_x),
            rect.y + row * (cell_h + gap_y),
            cell_w, cell_h,
        ))
    return rects


def _split_bullet_text(text: str) -> tuple[str, str, bool]:
    cleaned = (text or "").strip()
    is_action = False
    if cleaned.startswith("ACTION:"):
        is_action = True
        cleaned = cleaned[7:].strip()
    for sep in [" — ", " – ", ": ", "："]:
        if sep in cleaned:
            left, right = cleaned.split(sep, 1)
            left, right = left.strip(), right.strip()
            if left and right and len(left) <= 28:
                return left, right, is_action
    return cleaned, "", is_action


def _stats_panel_fonts(title_text: str, body_text: str) -> tuple[float, float]:
    if body_text:
        return 28, 13
    if len(title_text.strip()) > 36:
        return 18, 13
    return 24, 13


def _parse_table_rows(lines: list[str]) -> list[list[str]]:
    """Parse a slide's bullet lines into table rows.

    Supported row syntaxes (auto-detected in priority order):
    - Markdown table rows: ``A | B | C``  (separator rows like ``---|---`` are skipped)
    - TSV rows: ``A\\tB\\tC``
    - CSV-ish rows: ``A, B, C``
    - Single-column fallback when no delimiter is present
    """
    parsed: list[list[str]] = []
    for raw in lines:
        line = (raw or '').strip()
        if not line:
            continue
        # Skip markdown separator rows (e.g. "---|---|---")
        if set(line) <= {'|', '-', ':', ' '}:
            continue

        cells: list[str]
        if '|' in line:
            parts = [part.strip() for part in line.split('|')]
            cells = [part for part in parts if part]
        elif '\t' in line:
            cells = [part.strip() for part in line.split('\t')]
        elif ', ' in line or ',' in line:
            cells = [part.strip() for part in line.split(',')]
        else:
            cells = [line]

        cleaned = [cell for cell in cells if cell]
        if cleaned:
            parsed.append(cleaned)

    if not parsed:
        return []

    col_count = max(len(row) for row in parsed)
    return [row + [''] * (col_count - len(row)) for row in parsed]


def _lighten_hex(hex_color: str, factor: float = 0.92) -> str:
    """Lighten a hex color towards white by the given factor (0=original, 1=white)."""
    hex_color = hex_color.lstrip('#').upper()
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"{r:02X}{g:02X}{b:02X}"


def _is_numeric_cell(text: str) -> bool:
    """Check if a cell value looks numeric (for right-alignment)."""
    cleaned = text.strip().replace(',', '').replace('%', '').replace('$', '').replace('€', '').replace('¥', '').replace('£', '')
    if cleaned.startswith(('+', '-')):
        cleaned = cleaned[1:]
    try:
        float(cleaned)
        return True
    except (ValueError, TypeError):
        return False


def _render_table_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                        data: dict, slide_index: int,
                        accent_a: str, accent_b: str,
                        colors: dict[str, str]) -> None:
    """Structured table slide with themed header, zebra stripes, and auto-sized fonts."""
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    rows = _parse_table_rows(data.get('bullets', []))
    if spec.content_rect and rows:
        content = spec.content_rect
        row_count = len(rows)
        col_count = max(len(row) for row in rows)

        # ── Inset the table slightly for visual breathing room ──
        inset_x = 0.08
        inset_y = 0.06
        tbl_x = content.x + inset_x
        tbl_y = content.y + inset_y
        tbl_w = content.w - inset_x * 2
        tbl_h = content.h - inset_y * 2

        table_shape = slide.shapes.add_table(
            row_count, col_count,
            Inches(tbl_x), Inches(tbl_y),
            Inches(tbl_w), Inches(tbl_h),
        )
        table = table_shape.table
        table.first_row = True

        # ── Column widths — first column slightly wider for labels ──
        if col_count >= 2:
            label_ratio = 1.25
            data_col_share = (tbl_w - tbl_w * label_ratio / (col_count - 1 + label_ratio))
            label_w = tbl_w - data_col_share
            data_w = data_col_share / max(col_count - 1, 1)
            table.columns[0].width = Inches(label_w)
            for col_idx in range(1, col_count):
                table.columns[col_idx].width = Inches(data_w)
        else:
            table.columns[0].width = Inches(tbl_w)

        # ── Row heights — header slightly taller ──
        body_rows = max(row_count - 1, 1)
        header_h = min(tbl_h * 0.16, 0.58)
        body_row_h = (tbl_h - header_h) / body_rows if body_rows > 0 else tbl_h
        if row_count == 1:
            header_h = tbl_h

        # ── Colors ──
        header_fill = accent_a
        bg_base = colors['BG']
        # Zebra stripe: tinted version of accent for alternate rows
        stripe_fill = _lighten_hex(accent_a, 0.92)
        header_text_color = ctx.ensure_contrast(colors['BG'], header_fill)
        body_text_color = ctx.ensure_contrast(colors['TEXT'], bg_base)
        # ── Font sizing — adaptive based on table density ──
        font_name = ctx.font_family
        cell_count = row_count * col_count
        if cell_count <= 12:
            body_font_size = 12
            header_font_size = 13
        elif cell_count <= 24:
            body_font_size = 11
            header_font_size = 12
        elif cell_count <= 40:
            body_font_size = 10
            header_font_size = 11
        else:
            body_font_size = 9
            header_font_size = 10

        # ── Detect numeric columns for right-alignment ──
        numeric_cols: set[int] = set()
        if row_count > 1:
            for col_idx in range(col_count):
                data_cells = [rows[r][col_idx] for r in range(1, row_count) if col_idx < len(rows[r])]
                non_empty = [c for c in data_cells if c.strip()]
                if non_empty and sum(1 for c in non_empty if _is_numeric_cell(c)) / len(non_empty) >= 0.6:
                    numeric_cols.add(col_idx)

        # ── Cell margins (inches → EMU) ──
        cell_margin_lr = Inches(0.10)
        cell_margin_tb = Inches(0.04)

        for row_idx, row_data in enumerate(rows):
            is_header = row_idx == 0
            is_even_body = (row_idx % 2 == 0) and not is_header  # zebra stripe

            table.rows[row_idx].height = Inches(header_h if is_header else body_row_h)

            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = value

                # ── Fill ──
                cell.fill.solid()
                if is_header:
                    cell.fill.fore_color.rgb = ctx.rgb_color(header_fill)
                elif is_even_body:
                    cell.fill.fore_color.rgb = ctx.rgb_color(stripe_fill)
                else:
                    cell.fill.fore_color.rgb = ctx.rgb_color(bg_base)

                # ── Cell margins ──
                cell.margin_left = cell_margin_lr
                cell.margin_right = cell_margin_lr
                cell.margin_top = cell_margin_tb
                cell.margin_bottom = cell_margin_tb

                # ── Vertical anchor ──
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # ── Text formatting ──
                text_frame = cell.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.NONE

                # Determine alignment: header row centers, numeric cols right-align
                if is_header:
                    align = PP_ALIGN.CENTER
                elif col_idx in numeric_cols:
                    align = PP_ALIGN.RIGHT
                elif col_idx == 0 and col_count >= 2:
                    align = PP_ALIGN.LEFT  # first column = label
                else:
                    align = PP_ALIGN.LEFT

                font_size = header_font_size if is_header else body_font_size
                text_color = header_text_color if is_header else body_text_color

                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = align
                    paragraph.line_spacing = 1.08
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        run.font.name = font_name
                        run.font.bold = is_header
                        run.font.color.rgb = RGBColor.from_string(text_color)

        # ── Table-level border styling via XML ──
        try:
            from pptx.oxml.ns import qn
            tbl_xml = table._tbl
            tbl_pr = tbl_xml.find(qn('a:tblPr'))
            if tbl_pr is None:
                from lxml import etree
                tbl_pr = etree.SubElement(tbl_xml, qn('a:tblPr'))
            # Set clean border style
            tbl_pr.set('bandRow', '1')  # enable banded rows
            tbl_pr.set('firstRow', '1')  # special first row formatting
            tbl_pr.set('bandCol', '0')
            tbl_pr.set('lastRow', '0')
            tbl_pr.set('lastCol', '0')
            tbl_pr.set('firstCol', '0')
        except Exception:
            pass  # XML styling is optional; table is usable without it

    _add_footer(ctx, slide, spec, data.get('footer_text', ''), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get('notes', ''))


def _tile_images(ctx: RenderContext, slide, image_paths: list[str],
                 rect: RectSpec | None) -> None:
    if not image_paths or rect is None:
        return
    if len(image_paths) == 1:
        ctx.safe_add_picture(
            slide.shapes, image_paths[0],
            Inches(rect.x), Inches(rect.y),
            width=Inches(rect.w), height=Inches(rect.h),
        )
        return

    cols = 2 if len(image_paths) > 1 else 1
    rows = int(math.ceil(len(image_paths) / cols))
    gap_x = min(rect.w * 0.03, 0.16)
    gap_y = min(rect.h * 0.04, 0.16)
    cell_w = (rect.w - gap_x * (cols - 1)) / cols
    cell_h = (rect.h - gap_y * (rows - 1)) / rows
    for idx, img_path in enumerate(image_paths):
        col = idx % cols
        row = idx // cols
        x = rect.x + col * (cell_w + gap_x)
        y = rect.y + row * (cell_h + gap_y)
        ctx.safe_add_picture(
            slide.shapes, img_path,
            Inches(x), Inches(y),
            width=Inches(cell_w), height=Inches(cell_h),
        )


# ── Icon helpers ─────────────────────────────────────────────────────

def _place_slide_icon(ctx: RenderContext, slide, spec: LayoutSpec,
                      slide_index: int, accent_hex: str) -> None:
    if not ctx.show_slide_icons:
        return
    icon_name = ctx.slide_icon_name(slide_index)
    if not icon_name:
        return
    icon_collection = ctx.slide_icon_collection(slide_index)
    icon_path = ctx.fetch_icon(icon_name, color_hex=accent_hex, required_collection=icon_collection)
    if not icon_path:
        return
    if spec.icon_rect:
        ctx.safe_add_picture(
            slide.shapes, icon_path,
            Inches(spec.icon_rect.x), Inches(spec.icon_rect.y),
            width=Inches(spec.icon_rect.w), height=Inches(spec.icon_rect.h),
        )
    else:
        ref = spec.title_rect or spec.key_message_rect or spec.content_rect
        if ref is not None:
            size = max(min(ref.h * 0.75, 0.72), 0.42)
            x = min(ref.x + ref.w - size, SLIDE_WIDTH_IN - size - 0.08)
            y = max(ref.y - size * 0.05, 0.10)
            ctx.safe_add_picture(
                slide.shapes, icon_path,
                Inches(x), Inches(y),
                width=Inches(size), height=Inches(size),
            )


def _place_panel_icon(ctx: RenderContext, slide, icon_id: str, accent_hex: str,
                      *, size_in: float, x: float, y: float,
                      required_collection: str | None = None) -> None:
    icon_path = ctx.fetch_icon(icon_id, color_hex=accent_hex, required_collection=required_collection)
    if icon_path:
        ctx.safe_add_design_picture(
            slide.shapes, icon_path,
            Inches(x), Inches(y),
            width=Inches(size_in), height=Inches(size_in),
        )


# ── Slide-level functions ────────────────────────────────────────────

def _set_speaker_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text or ""


def _add_footer(ctx: RenderContext, slide, spec: LayoutSpec,
                footer_text: str, colors: dict[str, str]) -> None:
    if not footer_text or not spec.footer_rect:
        return
    _add_textbox(ctx, slide, spec.footer_rect, footer_text,
                 8.5, colors["TEXT"], name="footer_citation")


def _add_title_and_key_message(
    ctx: RenderContext, slide, spec: LayoutSpec, data: dict,
    colors: dict[str, str], accent_a: str, accent_b: str, *,
    title_base_pt: float = 24,
) -> None:
    """Render title + optional key-message text for any content slide."""
    title_pt = _adjust_title_font(
        spec.title_rect, data["title_text"], title_base_pt,
        scale=ctx.style.title_font_scale,
    )
    title_align = PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT
    _add_textbox(ctx, slide, spec.title_rect, data["title_text"],
                 title_pt, ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                 bold=True, name="title", align=title_align)

    key_fill = _add_key_message_band(ctx, slide, spec, accent_a, accent_b, colors)
    if data["key_message_text"] and spec.key_message_rect:
        km_color = key_fill or colors["TEXT"]
        key_pt = _adjust_body_font(
            spec.key_message_rect.w, spec.key_message_rect.h,
            data["key_message_text"], 15.5,
        )
        _add_textbox(ctx, slide, spec.key_message_rect, data["key_message_text"],
                     key_pt, ctx.ensure_contrast(km_color, colors["BG"]),
                     name="key_message", align=title_align)


# ── Layout-type renderers ────────────────────────────────────────────

def _render_title_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                        data: dict, slide_index: int,
                        accent_a: str, accent_b: str,
                        colors: dict[str, str]) -> None:
    images = ctx.slide_image_paths(slide_index)
    if images:
        target_rect = spec.hero_rect or spec.sidebar_rect or spec.content_rect
        _tile_images(ctx, slide, images, target_rect)

    title_pt = _adjust_title_font(
        spec.title_rect, data["title_text"], 30,
        scale=ctx.style.title_font_scale,
    )
    _add_textbox(ctx, slide, spec.title_rect,
                 data["title_text"], title_pt,
                 ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                 bold=True, name="title",
                 align=PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT)

    key_fill = _add_key_message_band(ctx, slide, spec, accent_a, accent_b, colors)
    if data["key_message_text"] and spec.key_message_rect:
        km_color = key_fill or colors["TEXT"]
        key_pt = _adjust_body_font(
            spec.key_message_rect.w, spec.key_message_rect.h,
            data["key_message_text"], 17,
        )
        _add_textbox(ctx, slide, spec.key_message_rect, data["key_message_text"],
                     key_pt, ctx.ensure_contrast(km_color, colors["BG"]),
                     name="key_message",
                     align=PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT)

    # Chip labels
    chip_texts = data.get("chip_labels") or data.get("bullets") or []
    if spec.chips_rect and chip_texts:
        visible = chip_texts[:min(len(chip_texts), 3)]
        gap = min(spec.chips_rect.w * 0.025, 0.18)
        chip_w = (spec.chips_rect.w - gap * (len(visible) - 1)) / len(visible)
        for idx, chip_text in enumerate(visible):
            chip_rect = RectSpec(
                spec.chips_rect.x + idx * (chip_w + gap),
                spec.chips_rect.y, chip_w, spec.chips_rect.h,
            )
            fill_hex = ctx.accent_cycle[(slide_index + idx + 2) % len(ctx.accent_cycle)]
            chip = _add_panel(ctx, slide, chip_rect, fill_hex, fill_hex,
                              name=f"chip_{idx}")
            tf = chip.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            left_margin = 0.08
            right_margin = 0.08
            top_margin = 0.04
            bottom_margin = 0.04
            tf.margin_left = Inches(left_margin)
            tf.margin_right = Inches(right_margin)
            tf.margin_top = Inches(top_margin)
            tf.margin_bottom = Inches(bottom_margin)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.16
            chip_pt = _adjust_body_font(
                chip_rect.w - (left_margin + right_margin),
                chip_rect.h - (top_margin + bottom_margin),
                chip_text,
                11.2,
            )
            bg_for_contrast = _effective_panel_bg(ctx, fill_hex, colors)
            _write_run(p, chip_text, chip_pt,
                       ctx.ensure_contrast(colors["TEXT"], bg_for_contrast),
                       bold=False, font_name=ctx.font_family)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_bullets_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                          data: dict, slide_index: int,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> None:
    body_rect = spec.content_rect
    images = ctx.slide_image_paths(slide_index)
    if images and body_rect and not (spec.hero_rect or spec.sidebar_rect):
        body_rect, image_rect = _split_content_for_images(body_rect)
        _tile_images(ctx, slide, images, image_rect)
    elif images:
        _tile_images(ctx, slide, images, spec.hero_rect or spec.sidebar_rect or body_rect)

    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    bullet_rects = _grid_rects(body_rect, len(data["bullets"]), ctx.style.content_density,
                               ctx.style.layout_policy.whitespace_bias,
                               ctx.style.text_box_corner_style) if body_rect else []
    show_icons = ctx.text_box_style in ("with-icons", "mixed")

    for idx, bullet in enumerate(data["bullets"][:len(bullet_rects)]):
        rect = bullet_rects[idx]
        fill_hex = ctx.accent_cycle[(slide_index + idx) % len(ctx.accent_cycle)]
        panel = _add_panel(ctx, slide, rect, fill_hex, fill_hex,
                           accent_b_hex=accent_b, name=f"bullet_{idx}")
        _add_panel_stripe(ctx, slide, rect, fill_hex, idx)

        title_text, body_text, _ = _split_bullet_text(bullet)
        top_reserve = 0.0

        if show_icons:
            slide_icon = ctx.slide_icon_name(slide_index)
            icon_id = slide_icon or "fluent:checkmark-24-regular"
            icon_collection = ctx.slide_icon_collection(slide_index)
            icon_size = max(min(rect.h * 0.18, 0.34), 0.24)
            _place_panel_icon(ctx, slide, icon_id, fill_hex,
                              size_in=icon_size,
                              x=rect.x + rect.w - icon_size - 0.10,
                              y=rect.y + 0.10,
                              required_collection=icon_collection)
            top_reserve = 0.06

        stripe_w = (max(min(rect.w * 0.02, 0.08), 0.04) + 0.02) if ctx.style.panel_stripe else 0.0
        _write_panel_text(ctx, panel, title_text, body_text, fill_hex, colors,
                          title_pt=15 if body_text else 14,
                          body_pt=13 if body_text else 12.5,
                          top_reserve=top_reserve, left_reserve=stripe_w)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_cards_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                        data: dict, slide_index: int,
                        accent_a: str, accent_b: str,
                        colors: dict[str, str]) -> None:
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    if spec.cards:
        n_items = min(len(data["bullets"]), data.get("item_count") or len(data["bullets"]))
        card_rects = [spec.cards.card_rect(i) for i in range(n_items)]
        pattern = spec.cards.pattern or "standard"
        icon_size = spec.cards.icon_size or 0.46
        header_band_h = spec.cards.header_band_h or 0.34
        header_icon_count = max(spec.cards.header_icon_count, 1)
    else:
        card_rects = _grid_rects(spec.content_rect, len(data["bullets"]), ctx.style.content_density,
                                 ctx.style.layout_policy.whitespace_bias,
                                 ctx.style.text_box_corner_style) if spec.content_rect else []
        pattern = "standard"
        icon_size = 0.46
        header_band_h = 0.34
        header_icon_count = 1

    show_icons = ctx.text_box_style in ("with-icons", "mixed")

    for idx, bullet in enumerate(data["bullets"][:len(card_rects)]):
        rect = card_rects[idx]
        fill_hex = colors["PANEL_BASE"] if idx % 2 == 0 else colors.get("PANEL_ALT", colors["PANEL_BASE"])
        if ctx.style.panel_fill == "solid":
            fill_hex = ctx.accent_cycle[(slide_index + idx + 1) % len(ctx.accent_cycle)]
        panel = _add_panel(ctx, slide, rect, fill_hex, fill_hex,
                           accent_b_hex=accent_b, name=f"card_{idx}")
        title_text, body_text, _ = _split_bullet_text(bullet)
        slide_icon = ctx.slide_icon_name(slide_index)
        card_icon = slide_icon or "fluent:sparkle-24-regular"
        card_icon_collection = ctx.slide_icon_collection(slide_index)
        top_reserve = 0.0

        if pattern == "icon_card":
            ico = max(min(icon_size, rect.h * 0.22), 0.34)
            _place_panel_icon(ctx, slide, card_icon, fill_hex,
                              size_in=ico, x=rect.x + 0.14, y=rect.y + 0.10,
                              required_collection=card_icon_collection)
            top_reserve = ico + 0.08
        elif pattern == "header_icon_card":
            band_h = min(header_band_h, rect.h * 0.22)
            band = ctx.add_design_shape(
                slide.shapes, MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(rect.x), Inches(rect.y),
                Inches(rect.w), Inches(band_h),
                name=f"card_header_{idx}",
            )
            if ctx.style.color_treatment == "gradient":
                ctx.apply_gradient_fill(band, [fill_hex, accent_b],
                                        angle_degrees=ctx.style.gradient_angle)
                ctx.set_fill_transparency(band, 0.03)
            else:
                band.fill.solid()
                band.fill.fore_color.rgb = ctx.rgb_color(fill_hex)
                ctx.set_fill_transparency(band, 0.03)
            band.line.fill.background()

            ico = max(min(band_h * 0.58, 0.22), 0.16)
            n_icons = max(1, min(header_icon_count, 3))
            gap = ico * 0.35
            total_w = ico * n_icons + gap * (n_icons - 1)
            start_x = rect.x + max((rect.w - total_w) / 2, 0.10)
            for icon_idx in range(n_icons):
                bg_for_icon = _effective_panel_bg(ctx, fill_hex, colors)
                _place_panel_icon(ctx, slide, card_icon,
                                  ctx.ensure_contrast(colors["TEXT"], bg_for_icon),
                                  size_in=ico,
                                  x=start_x + icon_idx * (ico + gap),
                                  y=rect.y + (band_h - ico) / 2,
                                  required_collection=card_icon_collection)
            top_reserve = band_h + 0.05
        elif show_icons:
            ico = max(min(rect.h * 0.15, 0.28), 0.18)
            _place_panel_icon(ctx, slide, card_icon, fill_hex,
                              size_in=ico,
                              x=rect.x + rect.w - ico - 0.10,
                              y=rect.y + 0.10,
                              required_collection=card_icon_collection)
            top_reserve = 0.06

        _write_panel_text(ctx, panel, title_text, body_text, fill_hex, colors,
                          title_pt=15.5 if body_text else 14.5,
                          body_pt=13, top_reserve=top_reserve)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_comparison_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                             data: dict, slide_index: int,
                             accent_a: str, accent_b: str,
                             colors: dict[str, str]) -> None:
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    if spec.comparison:
        left_rect = spec.comparison.left
        right_rect = spec.comparison.right
    elif spec.content_rect:
        left_rect, right_rect = _split_content_for_images(spec.content_rect)
    else:
        return

    split_at = int(math.ceil(len(data["bullets"]) / 2))
    left_items = data["bullets"][:split_at]
    right_items = data["bullets"][split_at:]

    left_panel = _add_panel(ctx, slide, left_rect, accent_a, accent_a,
                            accent_b_hex=accent_b, name="comparison_left")
    right_panel = _add_panel(ctx, slide, right_rect, accent_b, accent_b,
                             accent_b_hex=accent_a, name="comparison_right")

    show_icons = ctx.text_box_style in ("with-icons", "mixed")
    top_left = top_right = 0.0
    if show_icons:
        comparison_icon_collection = ctx.slide_icon_collection(slide_index)
        left_ico = max(min(left_rect.h * 0.12, 0.26), 0.18)
        right_ico = max(min(right_rect.h * 0.12, 0.26), 0.18)
        _place_panel_icon(ctx, slide, "fluent:warning-24-regular", accent_a,
                          size_in=left_ico,
                          x=left_rect.x + left_rect.w - left_ico - 0.10,
                          y=left_rect.y + 0.10,
                          required_collection=comparison_icon_collection)
        _place_panel_icon(ctx, slide, "fluent:checkmark-24-regular", accent_b,
                          size_in=right_ico,
                          x=right_rect.x + right_rect.w - right_ico - 0.10,
                          y=right_rect.y + 0.10,
                          required_collection=comparison_icon_collection)
        top_left = top_right = 0.05

    left_heading = data.get("comparison_left_label", "Before")
    right_heading = data.get("comparison_right_label", "After")
    # If data doesn't have explicit labels, use generic ones based on language
    title_text = data.get("title_text", "")
    if left_heading == "Before" and any(ord(c) > 0x3000 for c in title_text):
        left_heading = "\u30a2\u30f3\u30c1\u30d1\u30bf\u30fc\u30f3"    # アンチパターン
        right_heading = "\u63a8\u5968\u30d1\u30bf\u30fc\u30f3"          # 推奨パターン

    _write_bullets_panel(ctx, left_panel, left_heading, left_items,
                         accent_a, colors,
                         heading_pt=15.5, bullet_pt=13, top_reserve=top_left)
    _write_bullets_panel(ctx, right_panel, right_heading, right_items,
                         accent_b, colors,
                         heading_pt=15.5, bullet_pt=13, top_reserve=top_right)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_timeline_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                           data: dict, slide_index: int,
                           accent_a: str, accent_b: str,
                           colors: dict[str, str]) -> None:
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    if spec.timeline:
        tl = spec.timeline
        # Vertical spine line
        spine = ctx.add_design_shape(
            slide.shapes, MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(tl.line_x), Inches(tl.line_y),
            Inches(0.04), Inches(tl.line_h),
            name="timeline_spine",
        )
        spine.fill.solid()
        spine.fill.fore_color.rgb = ctx.rgb_color(accent_a)
        spine.line.fill.background()

        for idx, bullet in enumerate(data["bullets"][:spec.max_items or 6]):
            node = tl.node_rect(idx)
            fill_hex = ctx.accent_cycle[(slide_index + idx) % len(ctx.accent_cycle)]

            # Dot
            dot = ctx.add_design_shape(
                slide.shapes, MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(tl.dot_x), Inches(node.y + node.h * 0.3),
                Inches(tl.dot_size), Inches(tl.dot_size),
                name=f"dot_{idx}",
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = ctx.rgb_color(fill_hex)
            dot.line.fill.background()

            # Text panel
            panel = _add_panel(ctx, slide, node, fill_hex, fill_hex,
                               name=f"timeline_{idx}")
            title_text, body_text, _ = _split_bullet_text(bullet)
            _write_panel_text(ctx, panel, title_text, body_text, fill_hex, colors,
                              title_pt=14, body_pt=12)
    else:
        # Fallback: render as bullets
        _render_bullets_body(ctx, slide, spec, data, slide_index, accent_a, accent_b, colors)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_stats_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                        data: dict, slide_index: int,
                        accent_a: str, accent_b: str,
                        colors: dict[str, str]) -> None:
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    if spec.stats:
        for idx, bullet in enumerate(data["bullets"][:spec.max_items or 3]):
            box = spec.stats.box_rect(idx)
            fill_hex = ctx.accent_cycle[(slide_index + idx) % len(ctx.accent_cycle)]
            panel = _add_panel(ctx, slide, box, fill_hex, fill_hex,
                               name=f"stat_{idx}")
            title_text, body_text, _ = _split_bullet_text(bullet)
            title_pt, body_pt = _stats_panel_fonts(title_text, body_text)
            _write_panel_text(ctx, panel, title_text, body_text, fill_hex, colors,
                              title_pt=title_pt, body_pt=body_pt, align=PP_ALIGN.CENTER)
    else:
        _render_bullets_body(ctx, slide, spec, data, slide_index, accent_a, accent_b, colors)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_summary_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                          data: dict, slide_index: int,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> None:
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)

    rects = _grid_rects(spec.content_rect, len(data["bullets"]), ctx.style.content_density,
                        ctx.style.layout_policy.whitespace_bias,
                        ctx.style.text_box_corner_style) if spec.content_rect else []
    show_icons = ctx.text_box_style in ("with-icons", "mixed")

    for idx, bullet in enumerate(data["bullets"][:len(rects)]):
        rect = rects[idx]
        fill_hex = ctx.accent_cycle[(slide_index + idx) % len(ctx.accent_cycle)]
        panel = _add_panel(ctx, slide, rect, fill_hex, fill_hex,
                           name=f"summary_{idx}")
        top_reserve = 0.0
        if show_icons:
            summary_icon_collection = ctx.slide_icon_collection(slide_index)
            ico = max(min(rect.h * 0.15, 0.26), 0.18)
            _place_panel_icon(ctx, slide, "fluent:task-list-24-regular", fill_hex,
                              size_in=ico,
                              x=rect.x + rect.w - ico - 0.10,
                              y=rect.y + 0.10,
                              required_collection=summary_icon_collection)
            top_reserve = 0.05
        _write_panel_text(ctx, panel, bullet, "", fill_hex, colors,
                          title_pt=14, body_pt=12, top_reserve=top_reserve)

    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_section_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                          data: dict, slide_index: int,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> None:
    """Section divider — large title, optional key message, no body panels."""
    title_pt = _adjust_title_font(
        spec.title_rect, data["title_text"], 36,
        scale=ctx.style.title_font_scale,
    )
    _add_textbox(ctx, slide, spec.title_rect, data["title_text"],
                 title_pt, ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                 bold=True, name="title",
                 align=PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT)

    _add_key_message_band(ctx, slide, spec, accent_a, accent_b, colors)
    if data["key_message_text"] and spec.key_message_rect:
        key_pt = _adjust_body_font(
            spec.key_message_rect.w, spec.key_message_rect.h,
            data["key_message_text"], 18,
        )
        _add_textbox(ctx, slide, spec.key_message_rect, data["key_message_text"],
                     key_pt, ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                     name="key_message")

    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_closing_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                          data: dict, slide_index: int,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> None:
    """Closing slide — similar to section but with footer."""
    _render_section_slide(ctx, slide, spec, data, slide_index,
                          accent_a, accent_b, colors)
    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)


def _render_quote_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                        data: dict, slide_index: int,
                        accent_a: str, accent_b: str,
                        colors: dict[str, str]) -> None:
    """Large quote text with optional attribution."""
    quote_text = data["bullets"][0] if data["bullets"] else data["title_text"]
    if spec.content_rect:
        _add_textbox(ctx, slide, spec.content_rect, quote_text,
                     28, ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                     bold=False, name="quote",
                     align=PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT,
                     line_spacing=1.5)
    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_big_number_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                             data: dict, slide_index: int,
                             accent_a: str, accent_b: str,
                             colors: dict[str, str]) -> None:
    """Hero number + supporting text."""
    metric_like_title = _is_metric_like_title(data["title_text"])
    title_base_pt = 72 if metric_like_title else 42
    title_min_pt = 28 if metric_like_title else 24
    title_pt = title_base_pt * ctx.style.title_font_scale
    if spec.title_rect:
        title_pt = _adjust_title_font(
            spec.title_rect,
            data["title_text"],
            title_base_pt,
            scale=ctx.style.title_font_scale,
            min_pt=title_min_pt,
        )

    key_pt = 18.0
    if data["key_message_text"] and spec.key_message_rect:
        key_pt = _adjust_body_font(
            spec.key_message_rect.w,
            max(spec.key_message_rect.h, 0.4),
            data["key_message_text"],
            18,
        )

    render_spec = _reflow_big_number_spec(
        spec,
        data["title_text"],
        data.get("key_message_text", ""),
        title_pt=title_pt,
        key_pt=key_pt,
    )

    if render_spec.title_rect:
        _add_textbox(ctx, slide, render_spec.title_rect, data["title_text"],
                     title_pt,
                     ctx.ensure_contrast(accent_a, colors["BG"]),
                     bold=True, name="big_number",
                     align=PP_ALIGN.CENTER)
    if data["key_message_text"] and render_spec.key_message_rect:
        _add_textbox(ctx, slide, render_spec.key_message_rect, data["key_message_text"],
                     key_pt, ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                     name="key_message", align=PP_ALIGN.CENTER)
    if data["bullets"] and render_spec.content_rect:
        _render_bullets_body(ctx, slide, render_spec, data, slide_index, accent_a, accent_b, colors)
    _place_slide_icon(ctx, slide, render_spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_photo_fullbleed_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                                  data: dict, slide_index: int,
                                  accent_a: str, accent_b: str,
                                  colors: dict[str, str]) -> None:
    """Full-bleed photo with title overlay."""
    images = ctx.slide_image_paths(slide_index)
    if images and spec.hero_rect:
        _tile_images(ctx, slide, images, spec.hero_rect)

    if spec.title_rect:
        _add_textbox(ctx, slide, spec.title_rect, data["title_text"],
                     28 * ctx.style.title_font_scale,
                     ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                     bold=True, name="title",
                     align=PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT)
    if data["key_message_text"] and spec.key_message_rect:
        _add_textbox(ctx, slide, spec.key_message_rect, data["key_message_text"],
                     16, ctx.ensure_contrast(colors["TEXT"], colors["BG"]),
                     name="key_message",
                     align=PP_ALIGN.CENTER if ctx.style.title_centered else PP_ALIGN.LEFT)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_chart_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                        data: dict, slide_index: int,
                        accent_a: str, accent_b: str,
                        colors: dict[str, str]) -> None:
    """Chart placeholder — title + empty content rect for chart data."""
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)
    # Chart data rendering is handled separately; place a placeholder border
    if spec.content_rect:
        placeholder = ctx.add_managed_shape(
            slide.shapes, MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(spec.content_rect.x), Inches(spec.content_rect.y),
            Inches(spec.content_rect.w), Inches(spec.content_rect.h),
            name="chart_area",
        )
        placeholder.fill.background()
        placeholder.line.fill.solid()
        placeholder.line.color.rgb = ctx.rgb_color(colors["BORDER"])
        placeholder.line.width = Pt(0.5)
    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_diagram_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                          data: dict, slide_index: int,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> None:
    """Diagram/agenda with optional sidebar."""
    _add_title_and_key_message(ctx, slide, spec, data, colors, accent_a, accent_b)
    images = ctx.slide_image_paths(slide_index)
    if images and spec.sidebar_rect:
        _tile_images(ctx, slide, images, spec.sidebar_rect)
    _render_bullets_body(ctx, slide, spec, data, slide_index, accent_a, accent_b, colors)
    _add_footer(ctx, slide, spec, data.get("footer_text", ""), colors)
    _place_slide_icon(ctx, slide, spec, slide_index, accent_a)
    _set_speaker_notes(slide, data.get("notes", ""))


def _render_process_slide(ctx: RenderContext, slide, spec: LayoutSpec,
                          data: dict, slide_index: int,
                          accent_a: str, accent_b: str,
                          colors: dict[str, str]) -> None:
    """Process / multi-column cards (uses cards spec with multi columns)."""
    _render_cards_slide(ctx, slide, spec, data, slide_index, accent_a, accent_b, colors)


def _render_bullets_body(ctx: RenderContext, slide, spec: LayoutSpec,
                         data: dict, slide_index: int,
                         accent_a: str, accent_b: str,
                         colors: dict[str, str]) -> None:
    """Shared bullet body rendering (no title/key-message) for reuse."""
    body_rect = spec.content_rect
    if not body_rect or not data["bullets"]:
        return
    bullet_rects = _grid_rects(body_rect, len(data["bullets"]), ctx.style.content_density,
                               ctx.style.layout_policy.whitespace_bias,
                               ctx.style.text_box_corner_style)
    for idx, bullet in enumerate(data["bullets"][:len(bullet_rects)]):
        rect = bullet_rects[idx]
        fill_hex = ctx.accent_cycle[(slide_index + idx) % len(ctx.accent_cycle)]
        panel = _add_panel(ctx, slide, rect, fill_hex, fill_hex, name=f"body_{idx}")
        _add_panel_stripe(ctx, slide, rect, fill_hex, idx)
        title_text, body_text, _ = _split_bullet_text(bullet)
        stripe_w = (max(min(rect.w * 0.02, 0.08), 0.04) + 0.02) if ctx.style.panel_stripe else 0.0
        _write_panel_text(ctx, panel, title_text, body_text, fill_hex, colors,
                          title_pt=14, body_pt=12, left_reserve=stripe_w)


# ── Renderer dispatch table ──────────────────────────────────────────

_RENDERERS: dict[str, object] = {
    "title": _render_title_slide,
    "section": _render_section_slide,
    "closing": _render_closing_slide,
    "bullets": _render_bullets_slide,
    "cards": _render_cards_slide,
    "comparison": _render_comparison_slide,
    "timeline": _render_timeline_slide,
    "stats": _render_stats_slide,
    "summary": _render_summary_slide,
    "diagram": _render_diagram_slide,
    "chart": _render_chart_slide,
    "table": _render_table_slide,
    "quote": _render_quote_slide,
    "big_number": _render_big_number_slide,
    "photo_fullbleed": _render_photo_fullbleed_slide,
    "process": _render_process_slide,
    "multi_column": _render_cards_slide,
    "agenda": _render_diagram_slide,
    "content_caption": _render_bullets_slide,
    "picture_caption": _render_photo_fullbleed_slide,
    "two_content": _render_comparison_slide,
    "title_only": _render_bullets_slide,
    "pyramid": _render_timeline_slide,
}


# ── Public API ───────────────────────────────────────────────────────

def render_presentation(
    layout_input: list[dict],
    layout_specs: list,
    theme: dict[str, str],
    style: StyleConfig,
    slide_assets: list[dict],
    output_path: str,
    *,
    template_path: str | None = None,
    template_meta: dict | None = None,
    font_family: str = "Calibri",
    title: str = "Presentation",
    # Utility functions injected from pptx-python-runner.py
    rgb_color=None,
    ensure_contrast=None,
    set_fill_transparency=None,
    apply_gradient_fill=None,
    fetch_icon=None,
    safe_add_picture=None,
    safe_add_design_picture=None,
    add_design_shape=None,
    add_managed_textbox=None,
    add_managed_shape=None,
    tag_as_design=None,
    resolve_font=None,
    get_blank_layout=None,
    apply_widescreen=None,
    slide_image_paths=None,
    slide_icon_name=None,
    slide_icon_collection=None,
    ensure_parent_dir=None,
    safe_image_path=None,
    workspace_dir: str = "",
    theme_explicit: bool = False,
    text_box_style: str = "plain",
    show_slide_icons: bool = True,
) -> str:
    """Render a PPTX file from structured slide data.

    Returns the output path (as string).
    """
    # ── Validate required utilities ──
    required_fns = {
        "rgb_color": rgb_color,
        "ensure_contrast": ensure_contrast,
        "set_fill_transparency": set_fill_transparency,
        "apply_gradient_fill": apply_gradient_fill,
        "fetch_icon": fetch_icon,
        "safe_add_picture": safe_add_picture,
        "add_design_shape": add_design_shape,
        "add_managed_textbox": add_managed_textbox,
        "add_managed_shape": add_managed_shape,
        "get_blank_layout": get_blank_layout,
        "apply_widescreen": apply_widescreen,
        "slide_image_paths": slide_image_paths,
        "slide_icon_name": slide_icon_name,
        "slide_icon_collection": slide_icon_collection,
    }
    missing = [k for k, v in required_fns.items() if v is None]
    if missing:
        raise ValueError(f"Missing required utility functions: {', '.join(missing)}")

    # ── Create presentation ──
    if template_path:
        prs = apply_widescreen(Presentation(template_path))
        blank_layout = get_blank_layout(prs)
    else:
        prs = apply_widescreen(Presentation())
        blank_layout = prs.slide_layouts[6]

    prs.core_properties.title = title

    colors = _build_colors(theme)
    accent_cycle = _build_accent_cycle(theme)

    # When the style is dark-mode but the theme's BG is light (e.g. no palette
    # generated or the TS side didn't swap BG/TEXT), swap them here so dark
    # styles always get a dark background and light text.
    if style.dark_mode:
        bg_lum = _luminance(colors["BG"])
        text_lum = _luminance(colors["TEXT"])
        if bg_lum > text_lum:
            # BG is lighter than TEXT — swap them
            colors["BG"], colors["TEXT"] = colors["TEXT"], colors["BG"]
            colors["LIGHT"], colors["DARK"] = colors["DARK"], colors["LIGHT"]
            colors["LIGHT2"], colors["DARK2"] = colors["DARK2"], colors["LIGHT2"]
            # Also update the raw theme dict so downstream _build_colors() stays
            # consistent (e.g. _add_panel frosted uses _build_colors(ctx.theme)).
            theme = dict(theme)
            theme["BG"] = colors["BG"]
            theme["TEXT"] = colors["TEXT"]
            theme["LIGHT"] = colors["LIGHT"]
            theme["DARK"] = colors["DARK"]
            theme["LIGHT2"] = colors["LIGHT2"]
            theme["DARK2"] = colors["DARK2"]

    print(f"[renderer] BG={colors['BG']} TEXT={colors['TEXT']} ACCENT1={colors['ACCENT1']} "
          f"style={style.panel_fill}/{style.color_treatment} dark={style.dark_mode}", file=sys.stderr)

    ctx = RenderContext(
        prs=prs,
        theme=theme,
        style=style,
        font_family=font_family,
        slide_assets=slide_assets,
        accent_cycle=accent_cycle,
        template_path=template_path,
        workspace_dir=workspace_dir,
        theme_explicit=theme_explicit,
        text_box_style=text_box_style,
        show_slide_icons=show_slide_icons,
        rgb_color=rgb_color,
        ensure_contrast=ensure_contrast,
        set_fill_transparency=set_fill_transparency,
        apply_gradient_fill=apply_gradient_fill,
        fetch_icon=fetch_icon,
        safe_add_picture=safe_add_picture,
        safe_add_design_picture=safe_add_design_picture or safe_add_picture,
        add_design_shape=add_design_shape,
        add_managed_textbox=add_managed_textbox,
        add_managed_shape=add_managed_shape,
        tag_as_design=tag_as_design or (lambda s, n="": None),
        resolve_font=resolve_font or (lambda t, b="Calibri": b),
        get_blank_layout=get_blank_layout,
        apply_widescreen=apply_widescreen,
        slide_image_paths=slide_image_paths,
        slide_icon_name=slide_icon_name,
        slide_icon_collection=slide_icon_collection,
        ensure_parent_dir=ensure_parent_dir or (lambda p: os.makedirs(os.path.dirname(p), exist_ok=True)),
        safe_image_path=safe_image_path or (lambda p: p),
    )

    # ── Render each slide ──
    for slide_index, raw_data in enumerate(layout_input):
        raw_bullets = raw_data.get("bullets") or []
        raw_chip_labels = raw_data.get("chip_labels") or []
        data = {
            "layout_type": str(raw_data.get("layout_type") or raw_data.get("layout") or "bullets").strip().lower(),
            "title_text": str(raw_data.get("title_text") or raw_data.get("title") or "").strip(),
            "key_message_text": str(raw_data.get("key_message_text") or raw_data.get("keyMessage") or "").strip(),
            "bullets": [str(item).strip() for item in raw_bullets if str(item).strip()],
            "chip_labels": [str(item).strip() for item in raw_chip_labels if str(item).strip()],
            "footer_text": str(raw_data.get("footer_text") or "").strip(),
            "notes": str(raw_data.get("notes") or "").strip(),
            "item_count": int(raw_data.get("item_count") or 0),
        }

        if slide_index >= len(layout_specs):
            print(f"[renderer] Warning: no layout spec for slide {slide_index + 1}, skipping", file=sys.stderr)
            continue

        spec = layout_specs[slide_index]
        accent_a = accent_cycle[slide_index % len(accent_cycle)]
        accent_b = accent_cycle[(slide_index + 1) % len(accent_cycle)]
        slide_colors = _resolve_slide_colors(ctx, colors, accent_a, accent_b, slide_index)

        # Add slide
        slide = prs.slides.add_slide(blank_layout)
        if not template_path:
            # Apply style-specific gradient background when available.
            # bg_colors is a signature element of the design style and
            # should be honoured even when the user has set an explicit
            # theme palette (theme_explicit only gates content colours).
            use_style_bg = (
                style.bg_colors
                and len(style.bg_colors) >= 2
            )
            if use_style_bg:
                _apply_slide_bg_gradient(slide, style.bg_colors)
            else:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = rgb_color(slide_colors["BG"])

        # Decorative accents (controlled by style)
        _add_design_language(ctx, slide, spec, accent_a, accent_b, slide_colors)

        # Dispatch to layout-specific renderer
        layout_type = data["layout_type"]
        renderer_fn = _RENDERERS.get(layout_type, _render_bullets_slide)
        renderer_fn(ctx, slide, spec, data, slide_index, accent_a, accent_b, slide_colors)

    # ── Save ──
    ctx.ensure_parent_dir(output_path)
    prs.save(output_path)
    print(f"[renderer] Saved {len(layout_input)} slides to {output_path}", file=sys.stderr)
    return output_path


# ── CLI entry point for standalone testing ───────────────────────────

def _cli_main() -> int:
    """Run the renderer from the command line for testing."""
    import argparse

    parser = argparse.ArgumentParser(description="Deterministic PPTX slide renderer")
    parser.add_argument("--layout-input", required=True, help="Path to layout-input.json")
    parser.add_argument("--layout-specs", required=True, help="Path to layout-specs.json")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument("--workspace-dir", default="", help="Workspace directory")
    parser.add_argument("--design-style", default="Blank White", help="Design style name")
    parser.add_argument("--theme-json", default="", help="Theme JSON string")
    parser.add_argument("--font-family", default="Calibri", help="Font family")
    parser.add_argument("--title", default="Presentation", help="Presentation title")
    args = parser.parse_args()

    # Load inputs
    with open(args.layout_input, "r", encoding="utf-8") as f:
        layout_input = json.load(f)

    with open(args.layout_specs, "r", encoding="utf-8") as f:
        specs_raw = json.load(f)

    # Deserialize layout specs
    from scripts.layout.hybrid_layout import deserialize_specs  # noqa: E402
    layout_specs = deserialize_specs(json.dumps(specs_raw))

    # Theme
    if args.theme_json:
        theme = json.loads(args.theme_json)
    else:
        theme = {"BG": "FFFFFF", "TEXT": "000000", "ACCENT1": "4472C4", "ACCENT2": "ED7D31",
                 "ACCENT3": "70AD47", "ACCENT4": "FFC000", "ACCENT5": "5B9BD5", "ACCENT6": "A5A5A5"}

    # Style
    from style_config import resolve_style_config  # noqa: E402
    style = resolve_style_config(args.design_style)

    # Slide assets (empty for CLI test)
    slide_assets_path = os.path.join(os.path.dirname(args.layout_input), "slide-assets.json")
    slide_assets = []
    if os.path.exists(slide_assets_path):
        with open(slide_assets_path, "r", encoding="utf-8") as f:
            slide_assets = json.load(f)

    # Import utility functions from pptx-python-runner (when run as CLI)
    runner_dir = str(Path(__file__).resolve().parent)
    sys.path.insert(0, runner_dir)

    # Minimal stubs for utilities that the renderer needs
    def _stub_rgb_color(hex_val, fallback="000000"):
        h = (hex_val or fallback or "000000").strip().lstrip("#")
        if len(h) < 6:
            h = h.ljust(6, "0")
        return RGBColor.from_string(h[:6])

    def _stub_ensure_contrast(fg, bg, min_ratio=4.5):
        return fg  # simplified for CLI testing

    def _stub_set_transparency(shape, val):
        pass

    def _stub_gradient(shape, stops, angle_degrees=35):
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(stops[0] if stops else "000000")

    def _stub_fetch_icon(name, color_hex="000000", size=256, required_collection=None):
        return None

    def _make_shape_adder(shapes, shape_type, left, top, w, h, name=""):
        shape = shapes.add_shape(shape_type, left, top, w, h)
        if name:
            shape.name = name
        return shape

    def _make_textbox_adder(shapes, left, top, w, h, name=""):
        tb = shapes.add_textbox(left, top, w, h)
        if name:
            tb.name = name
        return tb

    def _stub_image_paths(idx):
        if idx < len(slide_assets):
            sa = slide_assets[idx]
            p = sa.get("primaryImagePath")
            if p and os.path.isfile(p):
                return [p]
        return []

    def _stub_icon_name(idx):
        if idx < len(slide_assets):
            return slide_assets[idx].get("iconName")
        return None

    def _stub_icon_collection(idx):
        if idx < len(slide_assets):
            return slide_assets[idx].get("iconCollection")
        return None

    render_presentation(
        layout_input=layout_input,
        layout_specs=layout_specs,
        theme=theme,
        style=style,
        slide_assets=slide_assets,
        output_path=args.output,
        font_family=args.font_family,
        title=args.title,
        workspace_dir=args.workspace_dir,
        rgb_color=_stub_rgb_color,
        ensure_contrast=_stub_ensure_contrast,
        set_fill_transparency=_stub_set_transparency,
        apply_gradient_fill=_stub_gradient,
        fetch_icon=_stub_fetch_icon,
        safe_add_picture=lambda *a, **kw: None,
        safe_add_design_picture=lambda *a, **kw: None,
        add_design_shape=_make_shape_adder,
        add_managed_textbox=_make_textbox_adder,
        add_managed_shape=_make_shape_adder,
        tag_as_design=lambda s, n="": None,
        get_blank_layout=lambda prs: prs.slide_layouts[6],
        apply_widescreen=lambda prs: prs,
        slide_image_paths=_stub_image_paths,
        slide_icon_name=_stub_icon_name,
        slide_icon_collection=_stub_icon_collection,
        ensure_parent_dir=lambda p: os.makedirs(os.path.dirname(p), exist_ok=True) if os.path.dirname(p) else None,
        safe_image_path=lambda p: p if os.path.isfile(p) else None,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(_cli_main())
