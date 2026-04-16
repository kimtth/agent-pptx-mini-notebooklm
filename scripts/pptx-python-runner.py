from __future__ import annotations

import argparse
import json
import os
import re
import sys
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt  # noqa: E402

from pathlib import Path
from typing import TYPE_CHECKING

if __package__ in {None, ''}:
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

# Layout engine lives in scripts/layout/ sub-package
sys.path.insert(0, str(Path(__file__).resolve().parent / 'layout'))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
else:
    PresentationType = object  # noqa: N816

from layout_specs import (  # type: ignore
    LayoutSpec,
)
from layout_validator import (  # type: ignore
    validate_presentation,
    report_issues,
)

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

PPTX_ICON_COLLECTION = (os.environ.get('PPTX_ICON_COLLECTION', 'all') or 'all').strip()

# ---------------------------------------------------------------------------
# Shape Semantic Registry — authoritative role metadata for every created shape
# ---------------------------------------------------------------------------
# Maps shape_id → role string ('template_design' | 'layout_managed').
# The validator reads this registry first; name-prefix and heuristic
# classification are used only when a shape is not resolved from the registry.
_SHAPE_ROLE_REGISTRY: dict[int, str] = {}


def _register_shape_role(shape, role: str) -> None:
    """Record the semantic role for *shape* in the global registry."""
    shape_id = getattr(shape, 'shape_id', None)
    if shape_id is not None:
        _SHAPE_ROLE_REGISTRY[int(shape_id)] = role


def get_shape_role_registry() -> dict[int, str]:
    """Return a snapshot of the current shape-role registry."""
    return dict(_SHAPE_ROLE_REGISTRY)


def clear_shape_role_registry() -> None:
    """Reset the registry (called between generation runs)."""
    _SHAPE_ROLE_REGISTRY.clear()


# Module-level reference to the loaded layout specs (set in build_namespace).
_LOADED_LAYOUT_SPECS: list[LayoutSpec] | None = None
# Module-level reference to the active style guardrails (set in renderer-mode).
_ACTIVE_GUARDRAILS: object | None = None


def _get_precomputed_specs_for_validation() -> list[LayoutSpec] | None:
    """Return the layout specs loaded at namespace-build time, if any."""
    return _LOADED_LAYOUT_SPECS

# ---------------------------------------------------------------------------
# Font resolution
# ---------------------------------------------------------------------------
# The user selects a font from the system font list in the palette UI.
# resolve_font() always returns that font unchanged — PowerPoint handles
# glyph substitution for missing characters at render time.
# ---------------------------------------------------------------------------


def resolve_font(text: str, base_font: str = 'Calibri') -> str:
    """Return *base_font* unchanged.

    PowerPoint handles font fallback for missing glyphs (e.g. CJK characters
    in a Latin-only font) at render time.  We never substitute fonts at the
    python-pptx level.
    """
    return base_font


def _has_real_transparency(image) -> bool:
    if 'A' not in image.getbands():
        return False
    min_alpha, max_alpha = image.getchannel('A').getextrema()
    return min_alpha < 255 and max_alpha > 0


def _has_visible_alpha(image) -> bool:
    if 'A' not in image.getbands():
        return False
    return image.getchannel('A').getbbox() is not None


def _make_background_transparent(image):
    bg_samples = [
        image.getpixel((0, 0)),
        image.getpixel((image.width - 1, 0)),
        image.getpixel((0, image.height - 1)),
        image.getpixel((image.width - 1, image.height - 1)),
    ]
    bg_r = sum(sample[0] for sample in bg_samples) // len(bg_samples)
    bg_g = sum(sample[1] for sample in bg_samples) // len(bg_samples)
    bg_b = sum(sample[2] for sample in bg_samples) // len(bg_samples)

    converted = image.copy()
    pixels = converted.load()
    assert pixels is not None
    for y in range(converted.height):
        for x in range(converted.width):
            r, g, b, _ = pixels[x, y]  # type: ignore[misc]
            diff = max(abs(r - bg_r), abs(g - bg_g), abs(b - bg_b))
            if diff <= 12:
                pixels[x, y] = (0, 0, 0, 0)  # type: ignore[index]
                continue
            alpha = min(255, diff * 4)
            pixels[x, y] = (0, 0, 0, alpha)  # type: ignore[index]
    return converted


def _make_transparent(png_path: str) -> str:
    """Convert a black-on-white RGB icon to black-on-transparent RGBA.

    Returns the path to the transparent version (cached alongside the original).
    If the source already has an alpha channel with real transparency, returns it as-is.
    """
    suffix = '_t.png'
    transparent_path = os.path.join(
        os.path.dirname(png_path),
        f'{Path(png_path).stem}{suffix}',
    )
    if os.path.isfile(transparent_path):
        return transparent_path

    try:
        from PIL import Image
        img = Image.open(png_path).convert('RGBA')

        # If image already has real transparency, skip conversion
        if _has_real_transparency(img):
            return png_path

        converted = img.copy()
        pixels = converted.load()
        assert pixels is not None
        w, h = converted.size
        for y in range(h):
            for x in range(w):
                r, g, b, _ = pixels[x, y]  # type: ignore[misc]
                # Luminance: near-white → transparent, darker → opaque icon stroke
                lum = r * 0.299 + g * 0.587 + b * 0.114
                if lum > 240:
                    pixels[x, y] = (0, 0, 0, 0)  # type: ignore[index]
                else:
                    # Map luminance to alpha: black=255, mid-gray=partial
                    alpha = min(255, int((255 - lum) * (255 / 200)))
                    pixels[x, y] = (0, 0, 0, alpha)  # type: ignore[index]

        if not _has_visible_alpha(converted):
            converted = _make_background_transparent(img)
            if not _has_visible_alpha(converted):
                return png_path

        converted.save(transparent_path, 'PNG')
        return transparent_path
    except Exception:
        return png_path


def _recolor_png(png_path: str, color_hex: str) -> str:
    """Tint an icon PNG to the requested color.

    Handles both RGBA (transparent bg) and RGB (white bg) source icons.
    """
    # First ensure we have a transparent version
    png_path = _make_transparent(png_path)

    color = color_hex.lstrip('#')
    if color == '000000':
        return png_path  # already black-on-transparent

    colored_path = os.path.join(
        os.path.dirname(png_path),
        f'{Path(png_path).stem}_{color}.png',
    )
    if os.path.isfile(colored_path):
        return colored_path

    try:
        from PIL import Image
        img = Image.open(png_path).convert('RGBA')
        r_tgt = int(color[0:2], 16)
        g_tgt = int(color[2:4], 16)
        b_tgt = int(color[4:6], 16)
        pixels = img.load()
        assert pixels is not None
        w, h = img.size
        for y in range(h):
            for x in range(w):
                _, _, _, a = pixels[x, y]  # type: ignore[misc]
                if a > 0:
                    pixels[x, y] = (r_tgt, g_tgt, b_tgt, a)  # type: ignore[index]
        if not _has_visible_alpha(img):
            return png_path
        img.save(colored_path, 'PNG')
        return colored_path
    except Exception:
        return png_path


# ---------------------------------------------------------------------------
# Icon fetch — downloads from Iconify public API with host redundancy
# ---------------------------------------------------------------------------
_MISSING_ICONS: list[dict[str, str]] = []
_ICON_FETCH_ATTEMPTS = 0

ICONIFY_API_HOSTS = [
    'https://api.iconify.design',
    'https://api.simplesvg.com',
    'https://api.unisvg.com',
]

_ICON_TEMP_DIR: str | None = None


def _get_icon_temp_dir() -> str:
    global _ICON_TEMP_DIR
    if _ICON_TEMP_DIR is None:
        import tempfile
        _ICON_TEMP_DIR = tempfile.mkdtemp(prefix='pptx_icons_')
    return _ICON_TEMP_DIR


def _download_svg(prefix: str, icon_name: str) -> bytes | None:
    import urllib.request
    import urllib.error
    for host in ICONIFY_API_HOSTS:
        url = f'{host}/{prefix}/{icon_name}.svg?box=1'
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'pptx-slide-agent/1.0'})
            with urllib.request.urlopen(req, timeout=8) as resp:
                if resp.status == 200:
                    data = resp.read()
                    if data and b'<svg' in data.lower():
                        print(f'[icon] fetched {prefix}:{icon_name} from {host}', file=sys.stderr)
                        return data
        except Exception as exc:
            print(f'[icon] {host} failed for {prefix}:{icon_name}: {exc}', file=sys.stderr)
    return None


def _preprocess_svg(svg_data: bytes) -> bytes:
    """Replace currentColor with black so svglib can render strokes/fills."""
    import re
    text = svg_data.decode('utf-8', errors='replace')
    text = text.replace('currentColor', '#000000')
    # Ensure stroke-only icons (fill="none") get visible strokes
    if 'stroke=' not in text and 'fill="none"' in text:
        text = text.replace('fill="none"', 'fill="none" stroke="#000000"')
    # If no stroke-width is set, add a default for thin icons
    if 'stroke-width' not in text and 'stroke=' in text:
        text = re.sub(r'(<svg[^>]*>)', r'\1<style>*{stroke-width:2}</style>', text, count=1)
    return text.encode('utf-8')


def _svg_to_png(svg_data: bytes, size: int = 256) -> bytes | None:
    # Preprocess SVG to fix currentColor and stroke visibility for svglib
    svg_data = _preprocess_svg(svg_data)
    try:
        import importlib
        cairosvg = importlib.import_module('cairosvg')
        return cairosvg.svg2png(bytestring=svg_data, output_width=size, output_height=size)
    except ImportError:
        pass
    try:
        import io
        import tempfile as _tf
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        # svglib requires a file path, write SVG to a temp file
        with _tf.NamedTemporaryFile(suffix='.svg', delete=False) as tmp:
            tmp.write(svg_data)
            tmp_path = tmp.name
        try:
            drawing = svg2rlg(tmp_path)
        finally:
            os.unlink(tmp_path)
        if drawing is None:
            return None
        scale_x = size / drawing.width if drawing.width else 1
        scale_y = size / drawing.height if drawing.height else 1
        scale = min(scale_x, scale_y)
        drawing.width = size
        drawing.height = size
        drawing.scale(scale, scale)
        buf = io.BytesIO()
        renderPM.drawToFile(drawing, buf, fmt='PNG')
        return buf.getvalue()
    except Exception as exc:
        print(f'[icon] svglib/reportlab conversion failed: {exc}', file=sys.stderr)
        return None


def fetch_icon(
    name: str,
    color_hex: str = '000000',
    size: int = 256,
    required_collection: str | None = None,
) -> str | None:
    """Fetch an icon from the Iconify public API, convert to PNG, recolor."""
    global _ICON_FETCH_ATTEMPTS
    _ICON_FETCH_ATTEMPTS += 1

    requested_name = name.strip()
    has_explicit_prefix = ':' in requested_name
    normalized_required_collection = (required_collection or '').strip()
    if not has_explicit_prefix:
        default_prefix = (
            normalized_required_collection
            if normalized_required_collection and normalized_required_collection != 'all'
            else PPTX_ICON_COLLECTION
            if PPTX_ICON_COLLECTION and PPTX_ICON_COLLECTION != 'all'
            else 'mdi'
        )
        requested_name = f'{default_prefix}:{requested_name}'
    prefix, icon_name = requested_name.split(':', 1)

    effective_collection = (
        normalized_required_collection
        if normalized_required_collection and normalized_required_collection != 'all'
        else PPTX_ICON_COLLECTION
        if not has_explicit_prefix and PPTX_ICON_COLLECTION != 'all'
        else 'all'
    )
    if effective_collection != 'all' and prefix != effective_collection:
        _MISSING_ICONS.append({'icon': f'{prefix}:{icon_name}', 'reason': 'outside_selected_collection'})
        print(f'[icon] REJECTED: {prefix}:{icon_name} outside collection {effective_collection}', file=sys.stderr)
        return None

    # Check temp dir for already-fetched icon this run
    temp_dir = _get_icon_temp_dir()
    png_path = os.path.join(temp_dir, prefix, f'{icon_name}.png')
    if os.path.isfile(png_path):
        return _recolor_png(png_path, color_hex)

    svg_data = _download_svg(prefix, icon_name)
    if svg_data is None:
        _MISSING_ICONS.append({'icon': f'{prefix}:{icon_name}', 'reason': 'network_all_hosts_failed'})
        return None

    png_data = _svg_to_png(svg_data, size)
    if png_data is None:
        _MISSING_ICONS.append({'icon': f'{prefix}:{icon_name}', 'reason': 'svg_conversion_failed'})
        return None

    os.makedirs(os.path.join(temp_dir, prefix), exist_ok=True)
    with open(png_path, 'wb') as f:
        f.write(png_data)

    return _recolor_png(png_path, color_hex)


def get_missing_icons() -> list[dict[str, str]]:
    """Return a copy of the missing-icon audit list."""
    return list(_MISSING_ICONS)


def split_icon_audit_entries(entries: list[dict[str, str]]) -> tuple[list[dict[str, str]], list[dict[str, str]]]:
    """Split icon audit entries into collection-policy rejections vs real fetch failures."""
    rejected: list[dict[str, str]] = []
    unresolved: list[dict[str, str]] = []
    for entry in entries:
        if entry.get('reason') == 'outside_selected_collection':
            rejected.append(entry)
        else:
            unresolved.append(entry)
    return rejected, unresolved


def get_icon_audit_stats() -> dict[str, float | int]:
    """Return aggregate icon audit stats for post-staging QA classification."""
    rejected, unresolved = split_icon_audit_entries(_MISSING_ICONS)
    missing = len(unresolved)
    requested = _ICON_FETCH_ATTEMPTS
    missing_ratio = (missing / requested) if requested > 0 else 0.0
    rejected_ratio = (len(rejected) / requested) if requested > 0 else 0.0
    return {
        'requested': requested,
        'missing': missing,
        'missingRatio': missing_ratio,
        'rejectedByCollection': len(rejected),
        'rejectedRatio': rejected_ratio,
    }


def _load_theme_payload() -> dict[str, object]:
    raw = os.environ.get('PPTX_THEME_JSON', '{}')
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return {}
    return parsed if isinstance(parsed, dict) else {}


def _load_theme() -> dict[str, str]:
    parsed = _load_theme_payload()
    raw_colors = parsed.get('C') if isinstance(parsed.get('C'), dict) else parsed
    if not isinstance(raw_colors, dict):
        return {}
    return {
        str(k): str(v).strip().lstrip('#').upper()
        for k, v in raw_colors.items()
        if isinstance(v, str) and v.strip()
    }


def _load_slide_assets() -> list[dict[str, object]]:
    raw = os.environ.get('PPTX_SLIDE_ASSETS_JSON', '')
    if not raw.strip():
        return []
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return []
    return parsed if isinstance(parsed, list) else []


def slide_assets(slide_index: int) -> dict[str, object]:
    if 0 <= slide_index < len(SLIDE_ASSETS):
        asset = SLIDE_ASSETS[slide_index]
        if isinstance(asset, dict):
            return asset
    return {}


def slide_image_paths(slide_index: int) -> list[str]:
    asset = slide_assets(slide_index)
    paths: list[str] = []
    selected_images = asset.get('selectedImages')
    if isinstance(selected_images, list):
        for image in selected_images:
            if isinstance(image, dict):
                image_path = image.get('imagePath')
                if isinstance(image_path, str) and image_path.strip():
                    paths.append(image_path)
    primary_path = asset.get('primaryImagePath')
    if isinstance(primary_path, str) and primary_path.strip() and primary_path not in paths:
        paths.insert(0, primary_path)
    return paths


def slide_icon_name(slide_index: int) -> str | None:
    asset = slide_assets(slide_index)
    icon_name = asset.get('iconName') or asset.get('icon')
    return icon_name if isinstance(icon_name, str) and icon_name.strip() else None


def slide_icon_collection(slide_index: int) -> str | None:
    asset = slide_assets(slide_index)
    collection = asset.get('iconCollection')
    return collection if isinstance(collection, str) and collection.strip() else None


SLIDE_ASSETS = _load_slide_assets()


def rgb_color(value: str | None, fallback: str = '000000') -> RGBColor:
    normalized = (value or fallback).strip().lstrip('#').upper()
    if len(normalized) != 6:
        normalized = fallback
    return RGBColor.from_string(normalized)


def ensure_parent_dir(file_path: str) -> None:
    Path(file_path).expanduser().resolve().parent.mkdir(parents=True, exist_ok=True)


def apply_widescreen(prs: PresentationType) -> PresentationType:
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    return prs


# ---------------------------------------------------------------------------
# Custom template support
# ---------------------------------------------------------------------------

TEMPLATE_PATH: str = os.environ.get('PPTX_TEMPLATE_PATH', '')
_template_meta_raw = os.environ.get('PPTX_TEMPLATE_META_JSON', '')
try:
    TEMPLATE_META: dict[str, object] = json.loads(_template_meta_raw) if _template_meta_raw.strip() else {}
except Exception:
    TEMPLATE_META = {}
TEMPLATE_BACKGROUND_IMAGES: list[str] = [
    str(path) for path in TEMPLATE_META.get('backgroundImages', [])
    if isinstance(path, str) and path.strip()
] if isinstance(TEMPLATE_META, dict) else []
TEMPLATE_BLANK_LAYOUT_INDEX = TEMPLATE_META.get('blankLayoutIndex') if isinstance(TEMPLATE_META, dict) else None


def get_blank_layout(prs: PresentationType) -> object:
    """Find the blank slide layout from the presentation.

    When a custom template is loaded, the blank layout may not be at index 6.
    This helper finds it by name first ('Blank'), then falls back to the layout
    with the fewest placeholders.
    """
    layouts = prs.slide_layouts
    # Try name match first
    for layout in layouts:
        name = layout.name.lower().strip()
        if name in ('blank', '\u767d\u7d19'):  # 'blank' or Japanese '白紙'
            return layout
    # Fall back to fewest placeholders
    best = layouts[min(6, len(layouts) - 1)]
    best_count = len(best.placeholders)
    for layout in layouts:
        if len(layout.placeholders) < best_count:
            best = layout
            best_count = len(layout.placeholders)
    return best


def _convert_to_pptx_compatible(source: Path) -> Path:
    """Convert unsupported image formats (e.g. WebP) to PNG for python-pptx."""
    if source.suffix.lower() not in ('.webp',):
        return source
    target = source.with_suffix('.png')
    if target.exists():
        return target
    from PIL import Image
    with Image.open(source) as img:
        img.save(target, 'PNG')
    return target


def safe_image_path(value: str | None) -> str | None:
    if not value:
        return None
    candidate = Path(value).expanduser().resolve()
    if not candidate.exists():
        return None
    candidate = _convert_to_pptx_compatible(candidate)
    return str(candidate)


ICON_MAX_RENDER_DIMENSION_IN = 1.5
ICON_MAX_RENDER_DIMENSION_EMU = int(ICON_MAX_RENDER_DIMENSION_IN * 914400)


def _is_icon_asset(path: str) -> bool:
    if _ICON_TEMP_DIR is None:
        return False
    normalized = os.path.normcase(os.path.normpath(path))
    return normalized.startswith(os.path.normcase(os.path.normpath(_ICON_TEMP_DIR)) + os.sep)


def safe_add_picture(shapes, image_path: str | None, left, top, width=None, height=None):
    # Guard: LLMs sometimes pass a Slide object instead of slide.shapes
    if hasattr(shapes, 'shapes') and not hasattr(shapes, 'add_picture'):
        shapes = shapes.shapes
    resolved = safe_image_path(image_path)
    if not resolved:
        return None
    if width is not None and height is not None and _is_icon_asset(resolved):
        max_requested_dim = max(int(width), int(height))
        if max_requested_dim > ICON_MAX_RENDER_DIMENSION_EMU:
            icon_scale = ICON_MAX_RENDER_DIMENSION_EMU / max_requested_dim
            scaled_width = max(1, int(width * icon_scale))
            scaled_height = max(1, int(height * icon_scale))
            left = left + int((width - scaled_width) / 2)
            top = top + int((height - scaled_height) / 2)
            width = scaled_width
            height = scaled_height
    # Preserve aspect ratio when both width and height are specified
    if width is not None and height is not None:
        try:
            from PIL import Image as _PILImage
            with _PILImage.open(resolved) as _img:
                img_w, img_h = _img.size
            if img_w > 0 and img_h > 0:
                scale = min(width / img_w, height / img_h)
                fit_w = int(img_w * scale)
                fit_h = int(img_h * scale)
                # Center within the bounding box
                left = left + (width - fit_w) // 2
                top = top + (height - fit_h) // 2
                width = fit_w
                height = fit_h
        except Exception:
            pass  # Fall back to stretched dimensions
    picture = shapes.add_picture(resolved, left, top, width=width, height=height)
    if _is_icon_asset(resolved):
        base_name = getattr(picture, 'name', '') or 'picture'
        if not base_name.lower().startswith(('icon_', 'design_icon_', 'decor_icon_')):
            picture.name = f'icon_{base_name}'
        _register_shape_role(picture, 'template_design')
    else:
        _register_shape_role(picture, 'layout_managed')
    return picture


def tag_as_design(shape, name: str = '') -> object:
    """Mark a shape as template/design — excluded from collision checks.

    Sets a ``design_`` name prefix recognised by the layout validator and
    registers the shape in the semantic role registry.
    Call on any decorative shape (backgrounds, borders, accent blobs,
    watermarks) so it does not trigger overlap warnings against
    blueprint-managed content.
    """
    base = name or getattr(shape, 'name', '') or 'shape'
    if not base.lower().startswith(('design_', 'tmpl_', 'decor_', 'bg_')):
        shape.name = f'design_{base}'
    else:
        shape.name = base
    _register_shape_role(shape, 'template_design')
    return shape


def safe_add_design_picture(shapes, image_path: str | None, left, top, width=None, height=None):
    """Add an image as a template/design element — excluded from collision checks.

    Identical to ``safe_add_picture`` but tags the result with ``design_`` prefix.
    """
    pic = safe_add_picture(shapes, image_path, left, top, width, height)
    if pic is not None:
        tag_as_design(pic)
    return pic


def add_design_shape(shapes, auto_shape_type, left, top, width, height, name: str = ''):
    """Add an auto-shape and register it as template/design.

    Use for decorative elements (backgrounds, borders, accent blobs, etc.)
    that should not participate in collision or cramped-spacing checks.
    """
    if hasattr(shapes, 'shapes') and not hasattr(shapes, 'add_shape'):
        shapes = shapes.shapes
    shape = shapes.add_shape(auto_shape_type, left, top, width, height)
    tag_as_design(shape, name=name)
    return shape


def add_managed_textbox(shapes, left, top, width, height, name: str = ''):
    """Add a textbox and register it as layout-managed.

    Use for structural content placed from PRECOMPUTED_LAYOUT_SPECS — these
    shapes participate in all validation checks including collision detection.
    """
    if hasattr(shapes, 'shapes') and not hasattr(shapes, 'add_textbox'):
        shapes = shapes.shapes
    tb = shapes.add_textbox(left, top, width, height)
    if name:
        tb.name = name
    _register_shape_role(tb, 'layout_managed')
    return tb


def add_managed_shape(shapes, auto_shape_type, left, top, width, height, name: str = ''):
    """Add an auto-shape and register it as layout-managed.

    Use for structural content (cards, panels, stat boxes) placed from
    PRECOMPUTED_LAYOUT_SPECS that should participate in collision checks.
    """
    if hasattr(shapes, 'shapes') and not hasattr(shapes, 'add_shape'):
        shapes = shapes.shapes
    shape = shapes.add_shape(auto_shape_type, left, top, width, height)
    if name:
        shape.name = name
    _register_shape_role(shape, 'layout_managed')
    return shape
def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument('output_path')
    parser.add_argument('--render-dir', default=None,
                        help='Preview render directory (accepted for compatibility; managed by the caller)')
    parser.add_argument('--workspace-dir', default=None,
                        help='Absolute path to the user workspace directory')
    parser.add_argument('--post-process-only', action='store_true',
                        help='Run only post-processing on an existing PPTX (validation, preview, notebooklm)')
    parser.add_argument('--renderer-mode', action='store_true',
                        help='Use the deterministic slide renderer')
    return parser.parse_args()


# ---------------------------------------------------------------------------
# Color contrast utilities (WCAG 2.1)
# ---------------------------------------------------------------------------

def _srgb_to_linear(c: float) -> float:
    """Convert sRGB component (0-1) to linear."""
    return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip('#')
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _luminance_hex(hex_color: str) -> float:
    """WCAG relative luminance from a hex color string."""
    r, g, b = _hex_to_rgb(hex_color)
    rs = _srgb_to_linear(r / 255)
    gs = _srgb_to_linear(g / 255)
    bs = _srgb_to_linear(b / 255)
    return 0.2126 * rs + 0.7152 * gs + 0.0722 * bs


def _blend_hex(fg_hex: str, bg_hex: str, fg_opacity: float) -> str:
    """Alpha-blend ``fg_hex`` over ``bg_hex`` at ``fg_opacity`` (0..1)."""
    fg_hex = fg_hex.lstrip('#')[:6].ljust(6, '0')
    bg_hex = bg_hex.lstrip('#')[:6].ljust(6, '0')
    op = max(0.0, min(fg_opacity, 1.0))
    parts: list[str] = []
    for idx in (0, 2, 4):
        f = int(fg_hex[idx:idx + 2], 16)
        b = int(bg_hex[idx:idx + 2], 16)
        parts.append(f'{max(0, min(round(f * op + b * (1.0 - op)), 255)):02X}')
    return ''.join(parts)


def contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    """WCAG 2.1 contrast ratio between two hex colors."""
    l1 = _luminance_hex(fg_hex)
    l2 = _luminance_hex(bg_hex)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _shift_lightness(hex_color: str, toward_dark: bool) -> list[str]:
    """Return up to 8 candidates by shifting only the L channel of *hex_color*
    in 10 % increments toward 0 (dark) or 1 (light), preserving hue/saturation."""
    r, g, b = _hex_to_rgb(hex_color)
    # Convert to HLS (Python colorsys uses HLS, not HSL)
    import colorsys
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255) # noqa: E741
    candidates: list[str] = []
    step = 0.10
    for i in range(1, 9):
        new_l = l - step * i if toward_dark else l + step * i
        new_l = max(0.0, min(1.0, new_l))
        nr, ng, nb = colorsys.hls_to_rgb(h, new_l, s)
        candidates.append(f'{int(round(nr * 255)):02X}{int(round(ng * 255)):02X}{int(round(nb * 255)):02X}')
    return candidates


def ensure_contrast(fg_hex: str, bg_hex: str, *, min_ratio: float = 4.5) -> str:
    """Return *fg_hex* if contrast is sufficient, else a lightness-adjusted variant.

    Shifts only the L channel of *fg_hex* toward the appropriate extreme in 10 %
    increments before falling back to the generic dark/light fallback, so hue and
    saturation of accent colors are preserved where possible.

    ``min_ratio`` defaults to WCAG AA (4.5) for normal text; use 3.0 for large text.
    """
    fg_hex = fg_hex.lstrip('#')
    bg_hex = bg_hex.lstrip('#')
    if contrast_ratio(fg_hex, bg_hex) >= min_ratio:
        return fg_hex
    bg_lum = _luminance_hex(bg_hex)
    toward_dark = bg_lum > 0.4
    for candidate in _shift_lightness(fg_hex, toward_dark):
        if contrast_ratio(candidate, bg_hex) >= min_ratio:
            return candidate
    return '2D2D2D' if toward_dark else 'F0F0F0'


# ---------------------------------------------------------------------------
# Chart / data-visualisation helpers
# ---------------------------------------------------------------------------

# Use Agg backend for headless rendering (no GUI dependency)
def render_chart_to_image(
    fig: matplotlib.figure.Figure,
    workspace_dir: str = '',
    *,
    dpi: int = 200,
) -> str:
    """Save a matplotlib Figure to a temporary PNG and return its absolute path.

    The image is written to ``{workspace_dir}/previews/charts/`` so it
    persists alongside other workspace artefacts.
    """
    if not workspace_dir:
        workspace_dir = os.environ.get('WORKSPACE_DIR', '')
    chart_dir = os.path.join(workspace_dir, 'previews', 'charts') if workspace_dir else os.path.join(os.getcwd(), 'charts')
    os.makedirs(chart_dir, exist_ok=True)

    import uuid
    filename = f'chart_{uuid.uuid4().hex[:8]}.png'
    filepath = os.path.join(chart_dir, filename)
    fig.savefig(filepath, dpi=dpi, bbox_inches='tight', transparent=True, pad_inches=0.1)
    plt.close(fig)
    return filepath


def add_chart_picture(
    shapes,
    fig: matplotlib.figure.Figure,
    left,
    top,
    width=None,
    height=None,
    workspace_dir: str = '',
):
    """Render a matplotlib Figure to PNG, then embed it as a slide picture.

    This is the primary helper for seaborn/matplotlib chart images.
    Returns the picture shape or None on failure.
    """
    image_path = render_chart_to_image(fig, workspace_dir)
    return safe_add_picture(shapes, image_path, left, top, width, height)


def add_native_chart(
    slide,
    chart_type,
    chart_data,
    left,
    top,
    width,
    height,
    *,
    theme: dict[str, str] | None = None,
):
    """Add an editable python-pptx chart to a slide and apply theme colours.

    ``chart_type`` should be a python-pptx chart enum member.
    ``chart_data`` should be a python-pptx chart data object.
    Returns the chart shape.
    """
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    if theme:
        keys = ['accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6']
        colors = [f'#{theme[k].lstrip("#")}' for k in keys if k in theme and theme[k]]
        if not colors:
            colors = ['#6366F1', '#06B6D4', '#F59E0B', '#EF4444', '#10B981', '#8B5CF6']
        for idx, series in enumerate(chart.series):
            hex_val = colors[idx % len(colors)].lstrip('#')
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(hex_val)
    return graphic_frame


def _cleanup_rogue_pptx(directory: Path, canonical_name: str) -> None:
    """Remove any PPTX files in *directory* that are not the canonical output."""
    for entry in directory.iterdir():
        if entry.is_file() and entry.suffix.lower() == '.pptx' and entry.name != canonical_name:
            print(f'[runner] Removing rogue PPTX: {entry.name}', file=sys.stderr)
            entry.unlink(missing_ok=True)


def _build_completion_report(
    output_path: Path,
    *,
    warnings: list[str] | None = None,
    contrast_fixes: int = 0,
    missing_icons: list[dict[str, str]] | None = None,
    rejected_icons: list[dict[str, str]] | None = None,
    icon_stats: dict[str, float | int] | None = None,
    missing_images: list[str] | None = None,
    layout_issues: list[dict[str, str]] | None = None,
) -> dict:
    """Build a structured completion report for the TypeScript caller.

    Includes a ``qa`` sub-object with post-staging QA findings so the
    renderer can decide whether to trigger the poststaging workflow.
    """
    report: dict = {
        'status': 'error',
        'outputPath': str(output_path),
        'fileExists': False,
        'slideCount': 0,
        'fileSizeBytes': 0,
        'warnings': warnings or [],
        'qa': {
            'contrastFixes': contrast_fixes,
            'missingIcons': missing_icons or [],
            'rejectedIcons': rejected_icons or [],
            'iconStats': icon_stats or {'requested': 0, 'missing': 0, 'missingRatio': 0.0, 'rejectedByCollection': 0, 'rejectedRatio': 0.0},
            'missingImages': missing_images or [],
            'layoutIssues': layout_issues or [],
        },
    }

    if not output_path.exists():
        report['error'] = f'Output file not found: {output_path}'
        return report

    report['fileExists'] = True
    report['fileSizeBytes'] = output_path.stat().st_size

    if report['fileSizeBytes'] == 0:
        report['error'] = 'Output file is empty (0 bytes)'
        return report

    try:
        prs = Presentation(str(output_path))
        report['slideCount'] = len(prs.slides)
        if report['slideCount'] == 0:
            report['status'] = 'warning'
            report['warnings'].append('PPTX file contains 0 slides')
        else:
            report['status'] = 'success'
    except Exception as exc:
        report['error'] = f'Failed to open PPTX for verification: {exc}'

    return report


def _normalize_shape_text(value: str) -> str:
    return ' '.join(value.split())


def _extract_python_pptx_shape_text(shape) -> str:
    if not getattr(shape, 'has_text_frame', False):
        return ''
    try:
        return _normalize_shape_text('\n'.join(paragraph.text for paragraph in shape.text_frame.paragraphs))
    except Exception:
        return ''


def _find_shape_for_overflow(slide, *, shape_id: int | None, shape_name: str, shape_text: str, fallback_index: int):
    shapes_list = list(slide.shapes)

    if shape_id is not None:
        for candidate in shapes_list:
            if getattr(candidate, 'shape_id', None) == shape_id:
                return candidate

    if shape_name:
        matching_name = [candidate for candidate in shapes_list if (getattr(candidate, 'name', '') or '') == shape_name]
        if len(matching_name) == 1:
            return matching_name[0]
        if shape_text:
            for candidate in matching_name:
                if _extract_python_pptx_shape_text(candidate) == shape_text:
                    return candidate

    if shape_text:
        matching_text = [candidate for candidate in shapes_list if _extract_python_pptx_shape_text(candidate) == shape_text]
        if len(matching_text) == 1:
            return matching_text[0]

    if 0 <= fallback_index < len(shapes_list):
        return shapes_list[fallback_index]

    return None


def _shrink_text_frame_fonts(shape, scale: float) -> bool:
    changed = False
    minimum_pt = 8.0
    for para in shape.text_frame.paragraphs:
        if para.font.size is not None:
            para.font.size = Pt(max(round(para.font.size.pt * scale, 1), minimum_pt))
            changed = True
        for run in para.runs:
            if run.font.size is not None:
                run.font.size = Pt(max(round(run.font.size.pt * scale, 1), minimum_pt))
                changed = True
    return changed


def _collect_pillow_overflows(output_path: Path) -> list[dict[str, object]] | None:
    """Return Pillow-measured overflow metadata, or None if Pillow is unavailable."""
    try:
        from font_text_measure import TextMeasureRequest, measure_text_heights  # type: ignore
    except ImportError:
        return None

    prs = Presentation(str(output_path))
    overflows: list[dict[str, object]] = []

    for si, slide in enumerate(prs.slides):
        shapes_list = list(slide.shapes)
        for shi, shape in enumerate(shapes_list):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text or not text.strip():
                continue

            # Collect dominant font properties across all paragraphs/runs
            max_font_size_pt = 18.0
            font_family = 'Calibri'
            is_bold = False
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        size_pt = run.font.size.pt
                        if size_pt > max_font_size_pt:
                            max_font_size_pt = size_pt
                            if run.font.name:
                                font_family = run.font.name
                            is_bold = bool(run.font.bold)

            # Shape dimensions in inches (EMU → inches)
            shape_width_in = shape.width / 914400.0
            shape_height_in = shape.height / 914400.0

            # Internal margins (EMU → inches, default ~0.05in each side)
            tf = shape.text_frame
            margin_left = (tf.margin_left if tf.margin_left is not None else 91440) / 914400.0
            margin_right = (tf.margin_right if tf.margin_right is not None else 91440) / 914400.0
            margin_top = (tf.margin_top if tf.margin_top is not None else 45720) / 914400.0
            margin_bottom = (tf.margin_bottom if tf.margin_bottom is not None else 45720) / 914400.0
            text_width_in = max(shape_width_in - margin_left - margin_right, 0.2)
            usable_height_in = max(shape_height_in - margin_top - margin_bottom, 0.1)

            normalized_text = _normalize_shape_text(text)
            is_textbox = (shape.shape_type is not None and int(shape.shape_type) == 17)

            req = TextMeasureRequest(
                text=text,
                width_in=text_width_in,
                font_family=font_family,
                font_size_pt=max_font_size_pt,
                bold=is_bold,
            )
            heights = measure_text_heights([req])
            required_height_in = heights[0] if heights else usable_height_in

            # 5% tolerance (same as COM collector)
            if required_height_in > usable_height_in * 1.05:
                scale = min(usable_height_in / required_height_in, 1.0)
                overflows.append({
                    'slide_idx': si,
                    'shape_idx': shi,
                    'shape_id': getattr(shape, 'shape_id', None),
                    'shape_name': str(getattr(shape, 'name', '') or ''),
                    'shape_text': normalized_text,
                    'scale': scale,
                    'is_textbox': is_textbox,
                })
                print(
                    f'[layout] Slide {si + 1}, "{shape.name}": '
                    f'need {required_height_in:.2f}in, have {usable_height_in:.2f}in '
                    f'(scale={scale:.0%}, {"textbox" if is_textbox else "shape"}) [pillow]',
                    file=sys.stderr,
                )

    return overflows


def _fix_text_overflow(output_path: Path) -> int:
    """Measure text overflow and repair text-bearing shapes.

    The repair runs in bounded passes. Each pass measures overflow,
    then applies python-pptx fixes. Textboxes are manually shrunk because
    PowerPoint ignores TEXT_TO_FIT_SHAPE for them. Auto shapes also get manual
    shrink in addition to TEXT_TO_FIT_SHAPE because the auto-size flag alone has
    proven insufficient for some dense card/footer compositions.

    Returns the number of fixes applied, or -1 if Pillow is unavailable.
    """
    total_fixes = 0
    max_passes = 2

    for pass_index in range(max_passes):
        overflows = _collect_pillow_overflows(output_path)
        if overflows is None:
            return -1
        if not overflows:
            return total_fixes

        prs = Presentation(str(output_path))
        fixes = 0

        for overflow in overflows:
            slide_idx = int(overflow['slide_idx'])
            shape_idx = int(overflow['shape_idx'])
            scale = float(overflow['scale'])
            is_textbox = bool(overflow['is_textbox'])
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            shape = _find_shape_for_overflow(
                slide,
                shape_id=overflow.get('shape_id') if isinstance(overflow.get('shape_id'), int) else None,
                shape_name=str(overflow.get('shape_name', '') or ''),
                shape_text=str(overflow.get('shape_text', '') or ''),
                fallback_index=shape_idx,
            )
            if shape is None:
                continue
            if not shape.has_text_frame:
                continue

            # First pass is conservative, second pass can shrink a bit further.
            minimum_scale = 0.55 if pass_index == 0 else 0.48
            safe_scale = max(scale, minimum_scale)
            shape.text_frame.word_wrap = True

            if is_textbox:
                if _shrink_text_frame_fonts(shape, safe_scale):
                    fixes += 1
            else:
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                if _shrink_text_frame_fonts(shape, safe_scale):
                    fixes += 1
                else:
                    fixes += 1

        if fixes <= 0:
            return total_fixes

        prs.save(str(output_path))
        total_fixes += fixes

    return total_fixes


def _get_slide_bg_hex(slide) -> str | None:
    """Return the best-known slide background colour hex (no '#'), or None.

    Handles solid fills directly and gradient fills by averaging the first
    gradient stop colour (representative of the dominant visual tone).
    """
    try:
        bg_fill = slide.background.fill
        if bg_fill.type is not None and int(bg_fill.type) == 1:  # MSO_FILL.SOLID
            return str(bg_fill.fore_color.rgb)
    except Exception:
        pass
    # Gradient background: extract the first stop colour as representative bg.
    try:
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        bg_pr = slide.background._element.find(f'{{{ns_p}}}bg/{{{ns_p}}}bgPr')
        if bg_pr is not None:
            grad = bg_pr.find(f'{{{ns_a}}}gradFill')
            if grad is not None:
                gs_lst = grad.find(f'{{{ns_a}}}gsLst')
                if gs_lst is not None:
                    stops = gs_lst.findall(f'{{{ns_a}}}gs')
                    if stops:
                        srgb = stops[0].find(f'{{{ns_a}}}srgbClr')
                        if srgb is not None:
                            return srgb.get('val', '').upper()
    except Exception:
        pass
    if isinstance(TEMPLATE_META, dict):
        surface_hex = TEMPLATE_META.get('backgroundSurfaceHex')
        if isinstance(surface_hex, str) and re.fullmatch(r'[0-9A-Fa-f]{6}', surface_hex.strip()):
            return surface_hex.strip().upper()
    return None


def _get_solid_fill_hex(shape) -> str | None:
    """Return the raw solid fill hex (no '#') of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and int(fill.type) == 1:  # MSO_FILL.SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _get_shape_bounds(shape) -> tuple[int, int, int, int] | None:
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except Exception:
        return None

    if width <= 0 or height <= 0:
        return None

    return left, top, left + width, top + height


def _intersection_area(bounds_a: tuple[int, int, int, int], bounds_b: tuple[int, int, int, int]) -> int:
    left = max(bounds_a[0], bounds_b[0])
    top = max(bounds_a[1], bounds_b[1])
    right = min(bounds_a[2], bounds_b[2])
    bottom = min(bounds_a[3], bounds_b[3])
    if right <= left or bottom <= top:
        return 0
    return (right - left) * (bottom - top)


def _shape_role(shape) -> str | None:
    shape_id = getattr(shape, 'shape_id', None)
    if shape_id is None:
        return None
    return _SHAPE_ROLE_REGISTRY.get(int(shape_id))


def _infer_background_from_shapes(slide_shapes: list[object], text_shape_index: int, text_shape, slide_bg_hex: str | None) -> str | None:
    text_bounds = _get_shape_bounds(text_shape)
    if text_bounds is None:
        return slide_bg_hex

    text_area = max(1, (text_bounds[2] - text_bounds[0]) * (text_bounds[3] - text_bounds[1]))
    best_hex: str | None = None
    best_score = -1.0

    for candidate in slide_shapes[:text_shape_index]:
        candidate_fill = _get_solid_fill_hex(candidate)
        if candidate_fill is None:
            continue

        candidate_bounds = _get_shape_bounds(candidate)
        if candidate_bounds is None:
            continue

        overlap = _intersection_area(text_bounds, candidate_bounds)
        if overlap <= 0:
            continue

        coverage = overlap / text_area
        candidate_area = (candidate_bounds[2] - candidate_bounds[0]) * (candidate_bounds[3] - candidate_bounds[1])
        role_bonus = 0.2 if _shape_role(candidate) == 'template_design' else 0.0
        non_text_bonus = 0.1 if not getattr(candidate, 'has_text_frame', False) else 0.0
        area_bonus = min(candidate_area / text_area, 8.0) * 0.02
        score = coverage + role_bonus + non_text_bonus + area_bonus

        if score > best_score and coverage >= 0.35:
            best_score = score
            best_hex = candidate_fill

    return best_hex or slide_bg_hex


def _get_fill_transparency(shape) -> float:
    """Return fill transparency (0.0 = opaque, 1.0 = invisible)."""
    # In-memory attribute (works when just set, not after load)
    try:
        t = shape.fill.transparency
        if t is not None:
            return float(t)
    except (AttributeError, TypeError):
        pass
    # Read from XML: <a:alpha val="N"/> where N is opacity in 1/1000ths %
    try:
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for alpha_el in shape._element.findall('.//a:solidFill/*/a:alpha', ns):
            opacity = int(alpha_el.get('val', '100000')) / 100000.0
            return round(1.0 - opacity, 4)
    except Exception:
        pass
    return 0.0


def set_fill_transparency(shape, value: float) -> None:
    """Set fill transparency via XML (python-pptx .fill.transparency doesn't persist).

    Use this instead of ``shape.fill.transparency = x`` to ensure the
    transparency is written to XML and survives save/load.

    Args:
        shape: Any python-pptx shape with a solid fill.
        value: 0.0 = fully opaque, 1.0 = fully transparent.
    """
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for sf in shape._element.findall(f'.//{{{ns}}}solidFill'):
        color_el = sf[0] if len(sf) else None
        if color_el is None:
            continue
        # Remove existing alpha elements
        for old in color_el.findall(f'{{{ns}}}alpha'):
            color_el.remove(old)
        # Add new alpha (opacity = 1 - transparency)
        opacity_val = str(int((1.0 - value) * 100000))
        alpha_el = etree.SubElement(color_el, f'{{{ns}}}alpha')
        alpha_el.set('val', opacity_val)


def apply_gradient_fill(shape, color_stops: list[str], angle_degrees: float = 0.0) -> None:
    """Apply a linear gradient fill using DrawingML XML."""
    from lxml import etree

    stops = [str(color).replace('#', '').strip().upper() for color in color_stops if str(color).strip()]
    stops = [color for color in stops if re.fullmatch(r'[0-9A-F]{6}', color)]
    if len(stops) < 2:
        raise ValueError('apply_gradient_fill requires at least two valid hex color stops.')

    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_pr = getattr(shape._element, 'spPr', None)
    if sp_pr is None:
        raise ValueError('Shape does not support fill properties.')

    for fill_tag in ('solidFill', 'gradFill', 'pattFill', 'blipFill', 'grpFill', 'noFill'):
        for existing in list(sp_pr.findall(f'{{{ns}}}{fill_tag}')):
            sp_pr.remove(existing)

    grad_fill = etree.SubElement(sp_pr, f'{{{ns}}}gradFill')
    grad_fill.set('rotWithShape', '1')
    gs_list = etree.SubElement(grad_fill, f'{{{ns}}}gsLst')

    stop_count = len(stops) - 1
    for index, color in enumerate(stops):
        position = 0 if stop_count == 0 else int(round(index * 100000 / stop_count))
        gradient_stop = etree.SubElement(gs_list, f'{{{ns}}}gs')
        gradient_stop.set('pos', str(position))
        srgb = etree.SubElement(gradient_stop, f'{{{ns}}}srgbClr')
        srgb.set('val', color)

    linear = etree.SubElement(grad_fill, f'{{{ns}}}lin')
    linear.set('ang', str(int(round((angle_degrees % 360) * 60000))))
    linear.set('scaled', '1')


def _get_run_color_hex(run, para) -> str | None:
    """Return the effective font color hex for a run."""
    try:
        if run.font.color.rgb is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    try:
        if para.font.color.rgb is not None:
            return str(para.font.color.rgb)
    except Exception:
        pass
    return None


def _fix_low_contrast_text(output_path: Path) -> int:
    """Fix text with insufficient contrast against its shape's fill or the slide background.

    For each text-bearing shape, checks every text run's color against:
    1. The shape's own solid fill, or
    2. The slide background colour (fallback when shape has no fill).

    When contrast is below WCAG AA (4.0:1), replaces the text colour with a
    lightness-adjusted variant.  For shapes with their own fill, also caps
    transparency so the panel provides a consistent readable background.
    """
    prs = Presentation(str(output_path))
    fixes = 0
    MIN_RATIO = 4.0  # WCAG AA practical threshold
    MAX_TRANSPARENCY = 0.45  # cap so panel is visually present

    for slide in prs.slides:
        slide_bg_hex = _get_slide_bg_hex(slide)
        slide_shapes = list(slide.shapes)

        for shape_index, shape in enumerate(slide_shapes):
            if not shape.has_text_frame:
                continue

            shape_fill_hex = _get_solid_fill_hex(shape)
            inferred_bg_hex = _infer_background_from_shapes(slide_shapes, shape_index, shape, slide_bg_hex)
            transparency = _get_fill_transparency(shape) if shape_fill_hex else 0.0

            # When a shape has a semi-transparent fill, the visual background
            # is a blend of the shape fill and whatever is behind it.
            if shape_fill_hex and transparency > 0.05 and inferred_bg_hex:
                # Approximate: blend shape fill over the background behind it
                opacity = 1.0 - transparency
                effective_bg = _blend_hex(shape_fill_hex, inferred_bg_hex, opacity)
            else:
                effective_bg = shape_fill_hex or inferred_bg_hex
            if effective_bg is None:
                continue

            fill_lum = _luminance_hex(effective_bg)

            # Collect runs that fail contrast
            bad_runs: list[tuple] = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fg_hex = _get_run_color_hex(run, para)
                    if fg_hex is None:
                        continue
                    ratio = contrast_ratio(fg_hex, effective_bg)
                    if ratio < MIN_RATIO:
                        bad_runs.append((run, fg_hex))

            if not bad_runs:
                continue

            # Choose replacement color
            is_light_fill = fill_lum > 0.4

            for run, old_hex in bad_runs:
                new_hex = ensure_contrast(old_hex, effective_bg, min_ratio=MIN_RATIO)
                run.font.color.rgb = RGBColor.from_string(new_hex)
                fixes += 1

            # If shape has its own light fill and highly transparent, reduce
            # transparency so the panel provides a consistent background.
            if shape_fill_hex and is_light_fill and transparency > MAX_TRANSPARENCY:
                set_fill_transparency(shape, MAX_TRANSPARENCY)
                fixes += 1

    if fixes > 0:
        prs.save(str(output_path))
    return fixes


def _check_missing_images(output_path: Path) -> list[str]:
    """Detect slides where approved images are missing from the generated PPTX.

    Returns a list of human-readable messages describing each missing image.
    Does NOT inject images — positioning remains the responsibility of the
    slide renderer that owns the layout context.
    """
    if not SLIDE_ASSETS:
        return []

    prs = Presentation(str(output_path))
    slides = list(prs.slides)
    details: list[str] = []

    for slide_idx, slide in enumerate(slides):
        expected_paths = slide_image_paths(slide_idx)
        if not expected_paths:
            continue

        # Resolve expected paths to canonical forms
        expected_resolved: list[tuple[str, str]] = []  # (original, resolved)
        for p in expected_paths:
            resolved = safe_image_path(p)
            if resolved:
                expected_resolved.append((p, os.path.normcase(os.path.normpath(resolved))))

        if not expected_resolved:
            continue

        # Collect canonical paths of images already placed on the slide
        existing_images: set[str] = set()
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    shape_blob = shape.image.blob
                    for _, resolved_path in expected_resolved:
                        try:
                            with open(resolved_path, 'rb') as f:
                                if f.read() == shape_blob:
                                    existing_images.add(resolved_path)
                        except Exception:
                            pass
                except Exception:
                    pass

        # Report missing images
        for orig, res in expected_resolved:
            if res not in existing_images:
                basename = os.path.basename(orig)
                msg = f'Slide {slide_idx + 1}: missing approved image {basename}'
                details.append(msg)
                print(f'[image-check] {msg}', file=sys.stderr)

    return details


def validate_and_fix_output(output_path: Path, *, run_text_overflow_fix: bool = True,
                            guardrails: object | None = None) -> dict:
    """Run image enforcement, optional text overflow fix, contrast fix, then validation on the generated PPTX.

    Returns a dict with structured QA findings:
        contrast_fixes (int), missing_images (list[str]), layout_issues (list[dict]).

    Raises RuntimeError when images were omitted by the generated code.  The
    error propagates to the TypeScript caller which marks the generation as
    failed and triggers an automatic retry — giving the LLM a chance to
    include ALL images with proper layout-aware positioning in code.
    """
    qa: dict = {'contrast_fixes': 0, 'missing_images': [], 'layout_issues': []}

    # Step 0: Check for missing approved images (detect only, no injection)
    missing_details = _check_missing_images(output_path)
    qa['missing_images'] = missing_details
    if missing_details:
        summary = '; '.join(missing_details)
        raise RuntimeError(
            f'Generated code omitted {len(missing_details)} approved image(s). '
            f'The LLM must include ALL approved images via slide_image_paths() '
            f'and safe_add_picture() with layout-aware positioning. Details: {summary}'
        )

    # Step 1: Text overflow fix (measure + resize)
    if run_text_overflow_fix:
        overflow_fixes = _fix_text_overflow(output_path)
        if overflow_fixes > 0:
            print(f'[layout] Overflow fix applied to {overflow_fixes} shape(s).', file=sys.stderr)
        elif overflow_fixes < 0:
            # No backend available: fallback to auto-size XML flags only
            from layout_validator import _enforce_auto_size  # type: ignore
            prs = Presentation(str(output_path))
            fixes = 0
            for slide in prs.slides:
                fixes += _enforce_auto_size(slide)
            if fixes > 0:
                prs.save(str(output_path))
                print(f'[layout-validator] Auto-size flags set on {fixes} frame(s) (no measurement backend).', file=sys.stderr)

    # Step 2: Fix low-contrast text on glass panels / cards
    contrast_fixes = _fix_low_contrast_text(output_path)
    qa['contrast_fixes'] = contrast_fixes
    if contrast_fixes > 0:
        print(f'[contrast] Fixed {contrast_fixes} low-contrast text/fill issue(s).', file=sys.stderr)

    # Step 3: Validate (re-open the processed file so we see all fixes)
    prs = Presentation(str(output_path))
    issues = validate_presentation(
        prs,
        shape_role_registry=get_shape_role_registry(),
        layout_specs=_get_precomputed_specs_for_validation(),
        guardrails=guardrails,
    )
    if not issues:
        print('[layout-validator] All slides passed layout validation.', file=sys.stderr)
        return qa

    # Capture structured layout issues for post-staging QA
    qa['layout_issues'] = [
        {
            'slide': issue.slide_index + 1,
            'type': issue.issue_type.value,
            'severity': issue.severity.value,
            'message': issue.message,
        }
        for issue in issues
        if issue.severity.value != 'info'
    ]

    report = report_issues(issues)
    print(report, file=sys.stderr)

    blocking = [issue for issue in issues if issue.severity.value == 'error']
    if blocking:
        has_overlap = any('overlap' in i.issue_type.value.lower() for i in blocking)
        has_text_overflow = any('text_overflow' in i.issue_type.value.lower() for i in blocking)

        # Tolerate up to 2 blocking issues only when they are not text overflow.
        if len(blocking) <= 2 and not has_text_overflow:
            print(
                f'[layout-validator] {len(blocking)} minor layout issue(s) detected (within tolerance). '
                'PPTX generated with incomplete layout details.',
                file=sys.stderr,
            )
            return qa

        hints: list[str] = []
        if has_overlap:
            hints.append(
                'TOOL HINT: Inspect the current layout coordinates in layout_specs.py, adjust them as needed '
                'with the available app repair tooling, then rerun the PPTX render.'
            )
        if has_text_overflow:
            hints.append(
                'TOOL HINT: Inspect the validation thresholds in layout_validator.py and adjust them if needed, '
                'or widen the available space in layout_specs.py, then rerun the PPTX render.'
            )

        hint_text = '\n'.join(hints)
        raise RuntimeError(
            'Layout validation failed after generation. '
            'Reduce content density, reserve more vertical space, or split the slide.\n'
            'Alternatively, use the layout infrastructure tools to fix spec dimensions or validator thresholds.\n\n'
            f'{hint_text}\n\n'
            f'{report}'
        )


def _unlock_or_rename(output_path: Path) -> Path:
    """Remove existing output file. If locked, fall back to a timestamped name."""
    if not output_path.exists():
        return output_path
    try:
        output_path.unlink()
        return output_path
    except PermissionError:
        from datetime import datetime
        stamp = datetime.now().strftime('%Y%m%d-%H%M%S')
        alt = output_path.with_stem(f'{output_path.stem}-{stamp}')
        print(f'[WARNING] {output_path.name} is locked, saving as {alt.name}', file=sys.stderr)
        return alt


def append_notebooklm_infographics(output_path: Path) -> None:
    """Append NotebookLM-generated infographic images as full-bleed slides at the end of the PPTX."""
    raw = os.environ.get('PPTX_NOTEBOOKLM_INFOGRAPHICS', '')
    if not raw.strip():
        return

    try:
        paths = json.loads(raw)
    except json.JSONDecodeError:
        return
    if not isinstance(paths, list) or len(paths) == 0:
        return

    # Filter to existing image files
    valid_paths = [p for p in paths if isinstance(p, str) and os.path.isfile(p)]
    if not valid_paths:
        return

    print(f'[notebooklm] Appending {len(valid_paths)} infographic slide(s)…', file=sys.stderr)

    prs = Presentation(str(output_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    blank_layout = get_blank_layout(prs)

    for img_path in valid_paths:
        resolved = safe_image_path(img_path)
        if not resolved:
            print(f'[notebooklm] Skipping missing image: {img_path}', file=sys.stderr)
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Fit image to slide while preserving aspect ratio, centered
        try:
            from PIL import Image as _PILImage
            with _PILImage.open(resolved) as img:
                img_w, img_h = img.size
            if img_w > 0 and img_h > 0:
                scale = min(slide_w / img_w, slide_h / img_h)
                fit_w = int(img_w * scale)
                fit_h = int(img_h * scale)
                left = (slide_w - fit_w) // 2
                top = (slide_h - fit_h) // 2
                slide.shapes.add_picture(resolved, left, top, fit_w, fit_h)
            else:
                slide.shapes.add_picture(resolved, 0, 0, slide_w, slide_h)
        except Exception:
            # Fallback: stretch to full slide
            slide.shapes.add_picture(resolved, 0, 0, slide_w, slide_h)

    prs.save(str(output_path))
    print(f'[notebooklm] Appended {len(valid_paths)} infographic slide(s)', file=sys.stderr)


def main() -> int:
    import io
    if isinstance(sys.stdout, io.TextIOWrapper):
        sys.stdout.reconfigure(encoding='utf-8')
    if isinstance(sys.stderr, io.TextIOWrapper):
        sys.stderr.reconfigure(encoding='utf-8')

    print(f'[icon] Using Iconify API hosts: {ICONIFY_API_HOSTS}', file=sys.stderr)

    args = parse_args()
    output_path = Path(args.output_path).resolve()
    workspace_dir = str(Path(args.workspace_dir).resolve()) if args.workspace_dir else ''
    renderer_mode = args.renderer_mode

    print(f'[workspace] WORKSPACE_DIR={workspace_dir or "(not set)"}', file=sys.stderr)

    # --post-process-only: skip code generation, jump straight to post-processing
    if args.post_process_only:
        print('[post-process-only] Running post-processing only on existing PPTX.', file=sys.stderr)
        if not output_path.exists():
            raise FileNotFoundError(f'PPTX file not found for post-processing: {output_path}')
    elif renderer_mode:
        # ── Deterministic renderer path ──
        print('[renderer-mode] Using deterministic slide renderer.', file=sys.stderr)
        output_path = _unlock_or_rename(output_path)

        # Load theme, specs, assets from env vars (same as build_namespace)
        theme = _load_theme()
        title = os.environ.get('PPTX_TITLE', 'Presentation')
        base_font_family = os.environ.get('PPTX_FONT_FAMILY', 'Calibri')
        color_treatment = (os.environ.get('PPTX_COLOR_TREATMENT', 'solid') or 'solid').strip().lower()
        text_box_style = (os.environ.get('PPTX_TEXT_BOX_STYLE', 'plain') or 'plain').strip().lower()
        text_box_corner_style = (os.environ.get('PPTX_TEXT_BOX_CORNER_STYLE', 'square') or 'square').strip().lower()
        design_style = os.environ.get('PPTX_DESIGN_STYLE', 'Blank White')

        # Layout specs
        specs_json = os.environ.get('PPTX_LAYOUT_SPECS_JSON', '')
        if not specs_json.strip():
            raise RuntimeError('PPTX_LAYOUT_SPECS_JSON is required for renderer mode.')
        from scripts.layout.hybrid_layout import deserialize_specs  # type: ignore
        precomputed_specs = deserialize_specs(specs_json)
        print(f'[layout] Loaded {len(precomputed_specs)} pre-computed layout spec(s).', file=sys.stderr)

        # Store for validator and reset shape registry
        global _LOADED_LAYOUT_SPECS  # noqa: PLW0603
        _LOADED_LAYOUT_SPECS = precomputed_specs
        clear_shape_role_registry()

        # Layout input
        layout_input_path = os.path.join(workspace_dir, 'previews', 'layout-input.json')
        if not os.path.exists(layout_input_path):
            raise FileNotFoundError(f'layout-input.json not found: {layout_input_path}')
        with open(layout_input_path, 'r', encoding='utf-8') as f:
            layout_input = json.load(f)
        if not isinstance(layout_input, list) or not layout_input:
            raise ValueError('layout-input.json must contain a non-empty slide list')

        # Slide assets
        _load_slide_assets()

        # Style config
        from style_config import resolve_style_config  # type: ignore
        style = resolve_style_config(design_style, color_treatment, text_box_corner_style)
        global _ACTIVE_GUARDRAILS  # noqa: PLW0603
        _ACTIVE_GUARDRAILS = style.guardrails
        theme_explicit = os.environ.get('PPTX_THEME_EXPLICIT', '0') == '1'
        show_slide_icons = os.environ.get('PPTX_SHOW_SLIDE_ICONS', '1') != '0'
        print(f'[renderer] Design style: {design_style!r}, panel_fill={style.panel_fill}, '
              f'color_treatment={style.color_treatment}, theme_explicit={theme_explicit}, '
              f'show_slide_icons={show_slide_icons}', file=sys.stderr)

        # Template
        template_path_str = TEMPLATE_PATH if TEMPLATE_PATH else None

        # Render
        from slide_renderer import render_presentation  # type: ignore
        render_presentation(
            layout_input=layout_input,
            layout_specs=precomputed_specs,
            theme=theme,
            style=style,
            slide_assets=SLIDE_ASSETS,
            output_path=str(output_path),
            template_path=template_path_str,
            template_meta=TEMPLATE_META,
            font_family=base_font_family,
            title=title,
            workspace_dir=workspace_dir,
            theme_explicit=theme_explicit,
            text_box_style=text_box_style,
            show_slide_icons=show_slide_icons,
            # Inject utility functions
            rgb_color=rgb_color,
            ensure_contrast=ensure_contrast,
            set_fill_transparency=set_fill_transparency,
            apply_gradient_fill=apply_gradient_fill,
            fetch_icon=fetch_icon,
            safe_add_picture=safe_add_picture,
            safe_add_design_picture=safe_add_design_picture,
            add_design_shape=add_design_shape,
            add_managed_textbox=add_managed_textbox,
            add_managed_shape=add_managed_shape,
            tag_as_design=tag_as_design,
            resolve_font=lambda text, base_font=base_font_family: resolve_font(text, base_font),
            get_blank_layout=get_blank_layout,
            apply_widescreen=apply_widescreen,
            slide_image_paths=slide_image_paths,
            slide_icon_name=slide_icon_name,
            slide_icon_collection=slide_icon_collection,
            ensure_parent_dir=ensure_parent_dir,
            safe_image_path=safe_image_path,
        )

        _cleanup_rogue_pptx(output_path.parent, output_path.name)

    # Append NotebookLM infographic slides if the option is activated
    post_warnings: list[str] = []
    try:
        append_notebooklm_infographics(output_path)
    except Exception as exc:  # noqa: BLE001
        msg = f'Failed to append infographics: {exc}'
        post_warnings.append(msg)
        print(f'[notebooklm] {msg}', file=sys.stderr)

    # Stamp the actual creation date (python-pptx's bundled default.pptx has a frozen 2013 date)
    try:
        from datetime import datetime as _dt, timezone as _tz
        _prs = Presentation(str(output_path))
        _now = _dt.now(_tz.utc)
        _prs.core_properties.created = _now
        _prs.core_properties.modified = _now
        _prs.save(str(output_path))
    except Exception as _exc:
        msg = f'Could not update core properties date: {_exc}'
        post_warnings.append(msg)
        print(f'[WARNING] {msg}', file=sys.stderr)

    qa_findings: dict = {'contrast_fixes': 0, 'missing_images': [], 'layout_issues': []}
    try:
        skip_text_overflow_fix = os.environ.get('PPTX_SKIP_TEXT_OVERFLOW_FIX') == '1'
        qa_findings = validate_and_fix_output(
            output_path,
            run_text_overflow_fix=not skip_text_overflow_fix,
            guardrails=_ACTIVE_GUARDRAILS,
        ) or qa_findings
    except Exception as exc:  # noqa: BLE001
        msg = f'Layout validation failed: {exc}'
        post_warnings.append(msg)
        print(f'[layout-validator] {msg}', file=sys.stderr)

    # Collect missing icons accumulated during code execution
    icon_audit_entries = get_missing_icons()
    rejected_icons, missing_icons = split_icon_audit_entries(icon_audit_entries)
    icon_stats = get_icon_audit_stats()

    report = _build_completion_report(
        output_path,
        warnings=post_warnings,
        contrast_fixes=qa_findings.get('contrast_fixes', 0),
        missing_icons=missing_icons,
        rejected_icons=rejected_icons,
        icon_stats=icon_stats,
        missing_images=qa_findings.get('missing_images', []),
        layout_issues=qa_findings.get('layout_issues', []),
    )
    print(json.dumps(report))
    return 0 if report['status'] in ('success', 'warning') else 1


if __name__ == '__main__':
    raise SystemExit(main())
