from __future__ import annotations

from pathlib import Path
import sys


REPO_ROOT = Path(__file__).resolve().parents[2]
LAYOUT_DIR = REPO_ROOT / 'scripts' / 'layout'
if str(LAYOUT_DIR) not in sys.path:
    sys.path.insert(0, str(LAYOUT_DIR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from layout_blueprint import ZoneRole, get_blueprint, list_layout_types
from layout_specs import (
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    LayoutSpec,
    RectSpec,
    flow_layout_spec,
)
from constraint_solver import solve_layout


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


BG_SPEC = rgb('F7F4EC')
BG_BLUEPRINT = rgb('EFF5F1')
INK = rgb('20302C')
MUTED = rgb('60706B')
LINE = rgb('6E7F79')
CARD_FILL = rgb('DDD2B0')
CARD_FILL_ALT = rgb('C7DCCF')
CARD_FILL_LIGHT = rgb('E8E1CD')
ACCENT = rgb('B8643C')
TIMELINE = rgb('6588A6')


def _set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_text(shape, text: str, font_size: int, *, bold: bool = False) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = 'Aptos'
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = INK
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.5)


def _add_box(
    slide,
    rect: RectSpec,
    label: str,
    *,
    fill_color: RGBColor,
    font_size: int = 18,
    rounded: bool = False,
) -> None:
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(rect.x),
        Inches(rect.y),
        Inches(rect.w),
        Inches(rect.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    _set_text(shape, label, font_size, bold=font_size >= 20)


def _add_label_tag(slide, x: float, y: float, label: str, *, fill_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(1.22),
        Inches(0.28),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    _set_text(shape, label, 10)


def _add_accent(slide, rect: RectSpec) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(rect.x),
        Inches(rect.y),
        Inches(rect.w),
        Inches(max(rect.h, 0.04)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    _add_label_tag(slide, rect.x, max(rect.y - 0.18, 0.1), 'accent', fill_color=CARD_FILL_LIGHT)


def _add_icon(slide, rect: RectSpec) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(rect.x),
        Inches(rect.y),
        Inches(rect.w),
        Inches(rect.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_FILL_ALT
    _set_text(shape, 'icon', 16, bold=True)


def _add_timeline(slide, spec, *, label: str) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(spec.line_x),
        Inches(spec.line_y),
        Inches(spec.line_x),
        Inches(spec.line_y + spec.line_h),
    )
    line.line.color.rgb = TIMELINE
    line.line.width = Pt(2.5)
    _add_label_tag(slide, spec.line_x + 0.14, spec.line_y, label, fill_color=CARD_FILL_LIGHT)

    for index in range(5):
        node = spec.node_rect(index)
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(spec.dot_x),
            Inches(node.y + 0.08),
            Inches(spec.dot_size),
            Inches(spec.dot_size),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = TIMELINE
        dot.line.fill.background()
        _add_box(slide, node, label, fill_color=CARD_FILL_LIGHT, font_size=13, rounded=True)


def _render_spec_slide(slide, spec: LayoutSpec) -> None:
    _set_slide_background(slide, BG_SPEC)

    if spec.title_rect is not None:
        _add_box(slide, spec.title_rect, 'title', fill_color=CARD_FILL, font_size=22, rounded=True)
    if spec.key_message_rect is not None:
        _add_box(slide, spec.key_message_rect, 'key_message', fill_color=CARD_FILL_LIGHT, font_size=16)
    if spec.accent_rect is not None:
        _add_accent(slide, spec.accent_rect)
    if spec.icon_rect is not None:
        _add_icon(slide, spec.icon_rect)
    if spec.summary_box is not None:
        _add_box(slide, spec.summary_box, 'summary_box', fill_color=CARD_FILL_ALT, font_size=16, rounded=True)
    if spec.hero_rect is not None:
        _add_box(slide, spec.hero_rect, 'hero', fill_color=CARD_FILL_ALT, font_size=22, rounded=True)
    if spec.sidebar_rect is not None:
        _add_box(slide, spec.sidebar_rect, 'sidebar', fill_color=CARD_FILL_ALT, font_size=16)
    if spec.chips_rect is not None:
        _add_box(slide, spec.chips_rect, 'chips', fill_color=CARD_FILL_LIGHT, font_size=13, rounded=True)
    if spec.notes_rect is not None:
        _add_box(slide, spec.notes_rect, 'notes', fill_color=rgb('F0E8D8'), font_size=11)

    if spec.content_rect is not None:
        has_variant = spec.cards or spec.stats or spec.comparison or spec.timeline
        if not has_variant:
            _add_box(slide, spec.content_rect, 'content', fill_color=rgb('F2EAD6'), font_size=18)

    if spec.cards is not None:
        for index in range(min(spec.max_items, 4)):
            _add_box(
                slide,
                spec.cards.card_rect(index),
                'cards',
                fill_color=rgb('D9E8E1'),
                font_size=16,
                rounded=True,
            )

    if spec.stats is not None:
        for index in range(3):
            _add_box(
                slide,
                spec.stats.box_rect(index),
                'stats',
                fill_color=rgb('D9E8E1'),
                font_size=16,
                rounded=True,
            )

    if spec.comparison is not None:
        _add_box(slide, spec.comparison.left, 'comparison', fill_color=rgb('D9E8E1'), font_size=16)
        _add_box(slide, spec.comparison.right, 'comparison', fill_color=rgb('D9E8E1'), font_size=16)

    if spec.timeline is not None:
        _add_timeline(slide, spec.timeline, label='timeline')

    if spec.footer_rect is not None:
        _add_box(slide, spec.footer_rect, 'footer', fill_color=CARD_FILL_LIGHT, font_size=13)


def _blueprint_stack_width(layout_type: str, icon_size: float, has_hero: bool, has_sidebar: bool, margin_x: float, gap_x: float) -> float:
    available = SLIDE_WIDTH_IN - (margin_x * 2)
    if has_hero:
        return min(7.85, available - (3.65 + gap_x))
    if has_sidebar:
        sidebar_w = 2.88 if layout_type == 'agenda' else 3.33
        return min(8.9, available - (sidebar_w + gap_x))
    if icon_size > 0:
        return min(9.3, available - (icon_size + 0.9))
    return available


def _build_blueprint_rects(layout_type: str) -> tuple[dict[ZoneRole, RectSpec], dict[str, RectSpec]]:
    blueprint = get_blueprint(layout_type)
    tokens = blueprint.tokens
    stack_width = _blueprint_stack_width(
        layout_type,
        blueprint.icon_size,
        blueprint.has_hero,
        blueprint.has_sidebar,
        tokens.margin_x,
        tokens.gap_x,
    )
    fixed_total = 0.0
    stretch_zones = []
    stack_zones = [zone for zone in blueprint.zones if zone.role is not ZoneRole.NOTES]

    gap_total = 0.0
    for index, zone in enumerate(stack_zones):
        if zone.stretch:
            stretch_zones.append(zone)
        else:
            fixed_total += zone.fixed_h or zone.preferred_h
        if index < len(stack_zones) - 1:
            gap_total += tokens.accent_gap if zone.role is ZoneRole.ACCENT else tokens.gap_y

    content_bottom = tokens.notes_y - 0.12
    remaining = max(content_bottom - tokens.margin_top - fixed_total - gap_total, 0.8)
    stretch_height = remaining / len(stretch_zones) if stretch_zones else 0.0

    rects: dict[ZoneRole, RectSpec] = {}
    extras: dict[str, RectSpec] = {}
    current_y = tokens.margin_top

    for zone in stack_zones:
        height = zone.fixed_h or (stretch_height if zone.stretch else zone.preferred_h)
        if zone.role in {ZoneRole.TITLE, ZoneRole.KEY_MESSAGE}:
            width = round(stack_width * tokens.header_w_ratio * zone.width_fraction, 2)
            x = round(tokens.margin_x + (stack_width - width) / 2, 2)
        elif zone.role is ZoneRole.ACCENT:
            width = 0.9 if layout_type in {'title', 'section'} else 1.5
            x = tokens.margin_x
        else:
            width = round(stack_width * zone.width_fraction, 2)
            x = tokens.margin_x

        rects[zone.role] = RectSpec(round(x, 2), round(current_y, 2), width, round(height, 2))
        current_y += height + (tokens.accent_gap if zone.role is ZoneRole.ACCENT else tokens.gap_y)

    rects[ZoneRole.NOTES] = RectSpec(
        tokens.margin_x,
        tokens.notes_y,
        round(SLIDE_WIDTH_IN - (tokens.margin_x * 2), 2),
        tokens.notes_h,
    )

    if blueprint.icon_size > 0:
        extras['icon'] = RectSpec(
            round(SLIDE_WIDTH_IN - tokens.icon_corner_margin_x - blueprint.icon_size, 2),
            tokens.icon_corner_margin_y,
            blueprint.icon_size,
            blueprint.icon_size,
        )

    content_rect = rects.get(ZoneRole.CONTENT)
    if blueprint.has_hero and content_rect is not None:
        extras['hero'] = RectSpec(8.85, content_rect.y, 3.65, min(content_rect.h, 3.65))
    if blueprint.has_sidebar and content_rect is not None:
        sidebar_w = 2.88 if layout_type == 'agenda' else 3.33
        extras['sidebar'] = RectSpec(
            round(SLIDE_WIDTH_IN - tokens.margin_x - sidebar_w, 2),
            content_rect.y,
            sidebar_w,
            content_rect.h,
        )

    return rects, extras


def _render_blueprint_slide(slide, layout_type: str) -> None:
    _set_slide_background(slide, BG_BLUEPRINT)
    blueprint = get_blueprint(layout_type)
    rects, extras = _build_blueprint_rects(layout_type)

    if ZoneRole.TITLE in rects:
        _add_box(slide, rects[ZoneRole.TITLE], 'title', fill_color=rgb('CFE0D2'), font_size=22, rounded=True)
    if ZoneRole.KEY_MESSAGE in rects:
        _add_box(slide, rects[ZoneRole.KEY_MESSAGE], 'key_message', fill_color=rgb('E0ECE4'), font_size=16)
    if ZoneRole.ACCENT in rects:
        _add_accent(slide, rects[ZoneRole.ACCENT])
    if ZoneRole.SUMMARY_BOX in rects:
        _add_box(slide, rects[ZoneRole.SUMMARY_BOX], 'summary_box', fill_color=rgb('CFE0D2'), font_size=16, rounded=True)
    if ZoneRole.CONTENT in rects:
        has_variant = blueprint.cards or blueprint.stats or blueprint.comparison or blueprint.timeline
        if not has_variant:
            _add_box(slide, rects[ZoneRole.CONTENT], 'content', fill_color=rgb('E4EFE7'), font_size=18)
    if ZoneRole.CHIPS in rects:
        _add_box(slide, rects[ZoneRole.CHIPS], 'chips', fill_color=rgb('E0ECE4'), font_size=13, rounded=True)
    if ZoneRole.NOTES in rects:
        _add_box(slide, rects[ZoneRole.NOTES], 'notes', fill_color=rgb('E8F0EA'), font_size=11)

    if 'icon' in extras:
        _add_icon(slide, extras['icon'])
    if 'hero' in extras:
        _add_box(slide, extras['hero'], 'hero', fill_color=rgb('CFE0D2'), font_size=22, rounded=True)
    if 'sidebar' in extras:
        _add_box(slide, extras['sidebar'], 'sidebar', fill_color=rgb('CFE0D2'), font_size=16)

    content_rect = rects.get(ZoneRole.CONTENT)
    if content_rect is None:
        return

    if blueprint.cards is not None:
        card_w = (content_rect.w - blueprint.cards.gap_x) / blueprint.cards.columns
        card_h = (content_rect.h - blueprint.cards.gap_y) / 2
        for row in range(2):
            for col in range(blueprint.cards.columns):
                rect = RectSpec(
                    x=content_rect.x + col * (card_w + blueprint.cards.gap_x),
                    y=content_rect.y + row * (card_h + blueprint.cards.gap_y),
                    w=card_w,
                    h=card_h,
                )
                _add_box(slide, rect, 'cards', fill_color=rgb('D6E6DE'), font_size=15, rounded=True)

    if blueprint.stats is not None:
        box_w = (content_rect.w - (blueprint.stats.gap_x * (blueprint.stats.columns - 1))) / blueprint.stats.columns
        for index in range(blueprint.stats.columns):
            rect = RectSpec(
                x=content_rect.x + index * (box_w + blueprint.stats.gap_x),
                y=content_rect.y,
                w=box_w,
                h=content_rect.h * 0.56,
            )
            _add_box(slide, rect, 'stats', fill_color=rgb('D6E6DE'), font_size=15, rounded=True)

    if blueprint.comparison is not None:
        left_w = (content_rect.w - blueprint.comparison.gap_x) / 2
        left = RectSpec(content_rect.x, content_rect.y, left_w, content_rect.h)
        right = RectSpec(content_rect.x + left_w + blueprint.comparison.gap_x, content_rect.y, left_w, content_rect.h)
        _add_box(slide, left, 'comparison', fill_color=rgb('D6E6DE'), font_size=16)
        _add_box(slide, right, 'comparison', fill_color=rgb('D6E6DE'), font_size=16)

    if blueprint.timeline is not None:
        class _BlueprintTimeline:
            line_x = blueprint.timeline.line_x
            line_y = content_rect.y
            line_h = content_rect.h
            dot_x = blueprint.timeline.dot_x
            dot_size = blueprint.timeline.dot_size
            start_y = content_rect.y
            step_y = content_rect.h / 5
            text_x = blueprint.timeline.text_x
            text_w = max(content_rect.right - blueprint.timeline.text_x, 2.4)

            @staticmethod
            def node_rect(index: int) -> RectSpec:
                return RectSpec(
                    x=_BlueprintTimeline.text_x,
                    y=_BlueprintTimeline.start_y + index * _BlueprintTimeline.step_y,
                    w=_BlueprintTimeline.text_w,
                    h=_BlueprintTimeline.step_y * 0.82,
                )

        _add_timeline(slide, _BlueprintTimeline, label='timeline')

    if ZoneRole.FOOTER in rects:
        _add_box(slide, rects[ZoneRole.FOOTER], 'footer', fill_color=rgb('E0ECE4'), font_size=13)


def build_presentation(output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    _SAMPLE_ITEM_COUNTS = {
        'title': 0, 'section': 0, 'agenda': 5, 'bullets': 6,
        'cards': 4, 'stats': 3, 'comparison': 6, 'timeline': 5,
        'summary': 3, 'diagram': 5, 'chart': 1,
        'closing': 0, 'photo_fullbleed': 0, 'multi_column': 5,
    }

    for layout_type in list_layout_types():
        bp = get_blueprint(layout_type)
        has_icon = bp.icon_size > 0
        item_count = _SAMPLE_ITEM_COUNTS.get(layout_type, 0)
        base_spec = solve_layout(bp, {}, has_icon=has_icon, item_count=item_count)
        spec = flow_layout_spec(
            base_spec,
            title_text='Sample Slide Title Text',
            key_message_text='Key message subtitle for this slide layout',
            chip_texts=['chips', 'chips', 'chips'],
        )

        spec_slide = prs.slides.add_slide(blank_layout)
        _render_spec_slide(spec_slide, spec)

        blueprint_slide = prs.slides.add_slide(blank_layout)
        _render_blueprint_slide(blueprint_slide, layout_type)

    prs.save(output_path)
    return output_path


def main() -> None:
    output_path = REPO_ROOT / 'samples' / 'layout_blueprint.pptx'
    saved_path = build_presentation(output_path)
    print(saved_path)


if __name__ == '__main__':
    main()