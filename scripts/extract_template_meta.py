"""Extract design metadata from a PPTX template file.

Outputs JSON to stdout with:
  - themeColors: 12 OOXML color slots (dk1, lt1, dk2, lt2, accent1-6, hlink, folHlink)
  - backgroundImages: list of extracted background image paths
  - blankLayoutIndex: index of the blank slide layout
  - fonts: { major, minor } from the theme font scheme
  - originalDimensions: { widthIn, heightIn }

Usage:
    python extract_template_meta.py <pptx_path> <assets_dir>
"""
from __future__ import annotations

import json
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.util import Emu


NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def _hex_from_clr_element(el: etree._Element) -> str | None:
    """Return 6-digit hex from an <a:srgbClr> or <a:sysClr> child."""
    srgb = el.find('a:srgbClr', NS)
    if srgb is not None:
        return srgb.get('val', '')
    sys_clr = el.find('a:sysClr', NS)
    if sys_clr is not None:
        return sys_clr.get('lastClr', '') or sys_clr.get('val', '')
    return None


def extract_theme_colors(prs: Presentation) -> dict[str, str]:
    """Extract 12 OOXML theme color slots from the presentation's theme XML."""
    slots = {
        'dk1': '000000', 'lt1': 'FFFFFF', 'dk2': '44546A', 'lt2': 'E7E6E6',
        'accent1': '4472C4', 'accent2': 'ED7D31', 'accent3': 'A5A5A5',
        'accent4': 'FFC000', 'accent5': '5B9BD5', 'accent6': '70AD47',
        'hlink': '0563C1', 'folHlink': '954F72',
    }
    try:
        # Navigate to the actual theme XML part
        theme_el = None
        for part in prs.part.package.iter_parts():
            if hasattr(part, 'content_type') and 'theme' in part.content_type:
                theme_el = etree.fromstring(part.blob)
                break
        if theme_el is None:
            return slots

        clr_scheme = theme_el.find('.//a:clrScheme', NS)
        if clr_scheme is None:
            return slots

        for slot_name in slots:
            el = clr_scheme.find(f'a:{slot_name}', NS)
            if el is not None:
                hex_val = _hex_from_clr_element(el)
                if hex_val:
                    slots[slot_name] = hex_val.upper()
    except Exception as exc:
        print(f'[extract] Warning: could not extract theme colors: {exc}', file=sys.stderr)
    return slots


def extract_fonts(prs: Presentation) -> dict[str, str]:
    """Extract major/minor font families from theme."""
    fonts = {'major': 'Calibri', 'minor': 'Calibri'}
    try:
        for part in prs.part.package.iter_parts():
            if hasattr(part, 'content_type') and 'theme' in part.content_type:
                theme_el = etree.fromstring(part.blob)
                font_scheme = theme_el.find('.//a:fontScheme', NS)
                if font_scheme is None:
                    break
                major_font = font_scheme.find('a:majorFont/a:latin', NS)
                minor_font = font_scheme.find('a:minorFont/a:latin', NS)
                if major_font is not None:
                    fonts['major'] = major_font.get('typeface', 'Calibri')
                if minor_font is not None:
                    fonts['minor'] = minor_font.get('typeface', 'Calibri')
                break
    except Exception as exc:
        print(f'[extract] Warning: could not extract fonts: {exc}', file=sys.stderr)
    return fonts


def find_blank_layout_index(prs: Presentation) -> int:
    """Find the index of the blank slide layout (no placeholders)."""
    best_index = 6  # default fallback
    best_ph_count = 999

    for i, layout in enumerate(prs.slide_layouts):
        name_lower = layout.name.lower().strip()
        ph_count = len(layout.placeholders)

        # Exact name match: prefer "Blank"
        if name_lower == 'blank' or name_lower == '白紙':
            return i

        # Fewest placeholders wins
        if ph_count < best_ph_count:
            best_ph_count = ph_count
            best_index = i

    return best_index


def extract_background_images(prs: Presentation, assets_dir: str) -> list[str]:
    """Extract background images from slide masters and save to assets_dir."""
    os.makedirs(assets_dir, exist_ok=True)
    extracted: list[str] = []

    try:
        for master_idx, master in enumerate(prs.slide_masters):
            bg = master.background
            if bg is None:
                continue

            bg_el = bg._element
            blip_fill = bg_el.find('.//p:bgPr/a:blipFill', NS) or bg_el.find('.//a:blipFill', NS)
            if blip_fill is None:
                continue

            blip = blip_fill.find('a:blip', NS)
            if blip is None:
                continue

            r_embed = blip.get(f'{{{NS["r"]}}}embed')
            if not r_embed:
                continue

            try:
                rel = master.part.rels[r_embed]
                image_blob = rel.target_part.blob
                ext = os.path.splitext(rel.target_ref)[1] or '.png'
                filename = f'master_{master_idx}_bg{ext}'
                out_path = os.path.join(assets_dir, filename)
                with open(out_path, 'wb') as f:
                    f.write(image_blob)
                extracted.append(out_path)
            except Exception as exc:
                print(f'[extract] Warning: could not extract bg image from master {master_idx}: {exc}', file=sys.stderr)
    except Exception as exc:
        print(f'[extract] Warning: background image extraction failed: {exc}', file=sys.stderr)

    return extracted


def main() -> None:
    if len(sys.argv) < 3:
        print('Usage: extract_template_meta.py <pptx_path> <assets_dir>', file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    assets_dir = sys.argv[2]

    if not os.path.isfile(pptx_path):
        print(json.dumps({'error': f'File not found: {pptx_path}'}))
        sys.exit(1)

    prs = Presentation(pptx_path)

    # Dimensions
    width_in = round(prs.slide_width / Emu(914400), 3)
    height_in = round(prs.slide_height / Emu(914400), 3)

    meta = {
        'themeColors': extract_theme_colors(prs),
        'backgroundImages': extract_background_images(prs, assets_dir),
        'blankLayoutIndex': find_blank_layout_index(prs),
        'fonts': extract_fonts(prs),
        'originalDimensions': {
            'widthIn': width_in,
            'heightIn': height_in,
        },
    }

    print(json.dumps(meta, ensure_ascii=False))


if __name__ == '__main__':
    main()
