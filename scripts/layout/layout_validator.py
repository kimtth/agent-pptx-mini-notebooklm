"""Post-processing validator for python-pptx slide layouts.

Detects overlapping shapes, out-of-bounds elements, and cramped spacing,
while keeping geometry changes limited to auto-size enforcement and boundary clamping.

Usage:

    from pptx import Presentation
    from layout_validator import validate_presentation, fix_presentation

    prs = Presentation('deck.pptx')
    issues = validate_presentation(prs)
    if issues:
        fixed = fix_presentation(prs, issues)
        print(f'Applied {fixed} non-layout geometry fixes')
    prs.save('deck.pptx')
"""

from __future__ import annotations

import json
import os
from dataclasses import dataclass
from enum import Enum

from pptx.util import Inches
from layout_specs import estimate_text_height_in

SLIDE_WIDTH_EMU = Inches(13.333)
SLIDE_HEIGHT_EMU = Inches(7.5)
SAFE_MARGIN_EMU = Inches(0.3)
MIN_GAP_EMU = Inches(0.15)
OVERLAP_TOLERANCE_EMU = Inches(0.05)

# Auto-size constant (avoids import issues when MSO_AUTO_SIZE isn't available)
_TEXT_TO_FIT_SHAPE = 2  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _should_include_images_in_layout() -> bool:
    """Check whether Picture shapes should participate in collision detection.

    Reads from PPTX_INCLUDE_IMAGES_IN_LAYOUT env var (set by pptx-handler),
    or falls back to layout-meta.json in the workspace previews directory.
    """
    if os.environ.get('PPTX_INCLUDE_IMAGES_IN_LAYOUT', '') == '1':
        return True
    workspace = os.environ.get('WORKSPACE_DIR', '')
    if workspace:
        meta_path = os.path.join(workspace, 'previews', 'layout-meta.json')
        try:
            with open(meta_path, encoding='utf-8') as f:
                meta = json.load(f)
            return bool(meta.get('includeImagesInLayout', False))
        except (OSError, json.JSONDecodeError, TypeError):
            pass
    return False


class IssueSeverity(Enum):
    INFO = 'info'
    WARNING = 'warning'
    ERROR = 'error'


class IssueType(Enum):
    OVERLAP = 'overlap'
    OUT_OF_BOUNDS = 'out_of_bounds'
    CRAMPED = 'cramped'
    TEXT_OVERFLOW = 'text_overflow'


class ShapeRole(Enum):
    """Semantic role for validation check eligibility.

    TEMPLATE_DESIGN — decorative / template shapes: checked for bounds and
    text overflow but excluded from overlap and cramped-spacing detection.
    LAYOUT_MANAGED — blueprint-driven content: participates in all checks.
    """
    TEMPLATE_DESIGN = 'template_design'
    LAYOUT_MANAGED = 'layout_managed'


# Shape-name prefixes that mark a shape as template/design.
# Shapes whose lowercased name starts with any of these are automatically
# classified as TEMPLATE_DESIGN and excluded from collision checks.
_DESIGN_NAME_PREFIXES = ('tmpl_', 'design_', 'bg_blob', 'bg_', 'decor_')
_ICON_NAME_PREFIXES = ('icon_', 'design_icon_', 'decor_icon_')


def _classify_shape_role(
    shape,
    *,
    include_images: bool,
    shape_role_registry: dict[int, str] | None = None,
    layout_specs: list | None = None,
    slide_index: int = -1,
) -> ShapeRole:
    """Classify a shape's semantic role for validation purposes.

    Priority:
      1. Explicit registry entry (authoritative, set at creation time)
      2. Blueprint geometry match (shape aligns with a known LayoutSpec rect)
      3. Explicit name prefix (backward compat)
      4. Heuristics (background fill, decorative frame, glow oval)
      5. Picture shape + include_images flag
      6. Default → LAYOUT_MANAGED
    """
    # 1. Registry lookup — single source of truth when available
    if shape_role_registry:
        shape_id = getattr(shape, 'shape_id', None)
        if shape_id is not None and int(shape_id) in shape_role_registry:
            role_str = shape_role_registry[int(shape_id)]
            if role_str == 'template_design':
                return ShapeRole.TEMPLATE_DESIGN
            return ShapeRole.LAYOUT_MANAGED

    # 2. Blueprint geometry match
    if layout_specs and slide_index >= 0:
        role = _match_blueprint_geometry(shape, layout_specs, slide_index)
        if role is not None:
            return role

    name = (getattr(shape, 'name', '') or '').lower()

    # 3. Explicit name-based classification (backward compat)
    for prefix in _ICON_NAME_PREFIXES:
        if name.startswith(prefix):
            return ShapeRole.TEMPLATE_DESIGN

    for prefix in _DESIGN_NAME_PREFIXES:
        if name.startswith(prefix):
            return ShapeRole.TEMPLATE_DESIGN

    # 4. Heuristic classification
    if _is_background_fill(shape):
        return ShapeRole.TEMPLATE_DESIGN
    if _is_decorative_frame(shape):
        return ShapeRole.TEMPLATE_DESIGN
    if _is_decorative_glow(shape):
        return ShapeRole.TEMPLATE_DESIGN

    # 5. Picture shapes: controlled by include_images flag
    if _is_picture_shape(shape):
        if not include_images:
            return ShapeRole.TEMPLATE_DESIGN

    # 6. Default
    return ShapeRole.LAYOUT_MANAGED


@dataclass
class ShapeBox:
    """Bounding box for a shape in EMU."""
    left: int
    top: int
    width: int
    height: int
    shape_name: str = ''
    shape_id: int = 0

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height


@dataclass
class LayoutIssue:
    slide_index: int
    issue_type: IssueType
    severity: IssueSeverity
    message: str
    shape_a: str = ''
    shape_b: str = ''


def _get_shape_box(shape) -> ShapeBox | None:
    """Extract bounding box from a python-pptx shape, or None if not positioned."""
    try:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
    except (AttributeError, TypeError):
        return None
    if left is None or top is None or width is None or height is None:
        return None
    if width <= 0 or height <= 0:
        return None
    name = getattr(shape, 'name', '') or ''
    shape_id = getattr(shape, 'shape_id', 0) or 0
    return ShapeBox(left=left, top=top, width=width, height=height,
                    shape_name=name, shape_id=shape_id)


def _boxes_overlap(a: ShapeBox, b: ShapeBox) -> bool:
    """AABB intersection test with a small tolerance to ignore trivial edge touches."""
    if a.right - OVERLAP_TOLERANCE_EMU <= b.left:
        return False
    if b.right - OVERLAP_TOLERANCE_EMU <= a.left:
        return False
    if a.bottom - OVERLAP_TOLERANCE_EMU <= b.top:
        return False
    if b.bottom - OVERLAP_TOLERANCE_EMU <= a.top:
        return False
    return True


def _overlap_area(a: ShapeBox, b: ShapeBox) -> int:
    """Return the overlapping area in EMU² (0 if no overlap)."""
    ox = max(0, min(a.right, b.right) - max(a.left, b.left))
    oy = max(0, min(a.bottom, b.bottom) - max(a.top, b.top))
    return ox * oy


def _is_contained(inner: ShapeBox, outer: ShapeBox, tolerance: int = 2) -> bool:
    """Return True if *inner* is fully contained within *outer* (within EMU tolerance)."""
    return (inner.left >= outer.left - tolerance and
            inner.top >= outer.top - tolerance and
            inner.right <= outer.right + tolerance and
            inner.bottom <= outer.bottom + tolerance)


def _is_background_fill(shape) -> bool:
    """Heuristic: shapes covering the full slide are background fills, not content."""
    box = _get_shape_box(shape)
    if box is None:
        return False
    covers_w = box.width >= SLIDE_WIDTH_EMU * 0.95
    covers_h = box.height >= SLIDE_HEIGHT_EMU * 0.95
    return covers_w and covers_h


def _is_decorative_frame(shape) -> bool:
    """Heuristic: unfilled rectangles that span a large portion of the slide
    are decorative borders / frames (e.g. double-frame outlines) and should
    be excluded from overlap detection.

    Criteria:
      - Rectangle or Rounded Rectangle auto-shape (type 1 or 5)
      - No solid fill (background fill or no fill)
      - No meaningful text content
      - Covers at least 50% of slide width AND 50% of slide height
    """
    box = _get_shape_box(shape)
    if box is None:
        return False

    # Must be large enough to act as a border frame
    is_large = (box.width >= SLIDE_WIDTH_EMU * 0.50 and
                box.height >= SLIDE_HEIGHT_EMU * 0.50)
    if not is_large:
        return False

    # Must have no text content
    text = _shape_text(shape)
    if text:
        return False

    # Must have no solid fill (background or no fill)
    try:
        fill = shape.fill
        # fill.type: None=no fill, 0=background, 1=solid, ...
        fill_type = fill.type
        if fill_type is not None and int(fill_type) == 1:  # solid fill
            return False
    except Exception:
        pass

    return True


def _shape_text(shape) -> str:
    if not getattr(shape, 'has_text_frame', False):
        return ''
    try:
        return '\n'.join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
    except Exception:
        return ''


def _is_decorative_glow(shape) -> bool:
    """Heuristic: text-free oval shapes are background decoration (glow blobs)."""
    name = getattr(shape, 'name', '') or ''
    if not name.startswith('Oval'):
        return False
    if _shape_text(shape):
        return False
    # Text-free ovals are decorative by design
    return True


def _is_picture_shape(shape) -> bool:
    """Check if a shape is a Picture (image) shape.

    Uses both ``isinstance`` and ``shape_type == 13`` (MSO_SHAPE_TYPE.PICTURE)
    so that placeholder-picture shapes and other image variants are caught.
    """
    from pptx.shapes.picture import Picture
    if isinstance(shape, Picture):
        return True
    shape_type = getattr(shape, 'shape_type', None)
    return shape_type is not None and int(shape_type) == 13


# ---------------------------------------------------------------------------
# Blueprint geometry matching
# ---------------------------------------------------------------------------

# Tolerance in EMU (≈0.12″) for matching a shape's position to a LayoutSpec rect.
_BLUEPRINT_MATCH_TOLERANCE_EMU = int(0.12 * 914400)


def _rect_matches(shape_box: 'ShapeBox', rect_x: float, rect_y: float, rect_w: float, rect_h: float) -> bool:
    """Return True if *shape_box* aligns with a blueprint rect within tolerance."""
    tol = _BLUEPRINT_MATCH_TOLERANCE_EMU
    rx = int(rect_x * 914400)
    ry = int(rect_y * 914400)
    rw = int(rect_w * 914400)
    rh = int(rect_h * 914400)
    return (abs(shape_box.left - rx) <= tol
            and abs(shape_box.top - ry) <= tol
            and abs(shape_box.width - rw) <= tol
            and abs(shape_box.height - rh) <= tol)


def _match_blueprint_geometry(shape, layout_specs: list, slide_index: int) -> ShapeRole | None:
    """Try to match *shape* against known LayoutSpec rects for *slide_index*.

    Returns ``ShapeRole.LAYOUT_MANAGED`` if the shape's bounding box matches
    a blueprint rect within tolerance, ``None`` otherwise (so the caller falls
    through to further classification).  Shapes matching accent rects are not
    positively classified here because they can be either structural or decorative.
    """
    if slide_index >= len(layout_specs):
        return None
    spec = layout_specs[slide_index]

    box = _get_shape_box(shape)
    if box is None:
        return None

    # Check each simple RectSpec field on the spec
    _RECT_FIELDS = (
        'title_rect', 'key_message_rect', 'icon_rect', 'content_rect',
        'notes_rect', 'summary_box', 'hero_rect', 'chips_rect',
        'footer_rect', 'sidebar_rect',
    )
    for field_name in _RECT_FIELDS:
        rect = getattr(spec, field_name, None)
        if rect is None:
            continue
        if _rect_matches(box, rect.x, rect.y, rect.w, rect.h):
            return ShapeRole.LAYOUT_MANAGED

    # Check card rects
    cards = getattr(spec, 'cards', None)
    if cards:
        max_items = getattr(spec, 'max_items', 0) or 12
        for idx in range(max_items):
            try:
                card_rect = cards.card_rect(idx)
                if _rect_matches(box, card_rect.x, card_rect.y, card_rect.w, card_rect.h):
                    return ShapeRole.LAYOUT_MANAGED
            except (IndexError, AttributeError):
                break

    # Check stats boxes
    stats = getattr(spec, 'stats', None)
    if stats:
        for idx in range(6):  # practical upper bound
            try:
                stat_rect = stats.box_rect(idx)
                if _rect_matches(box, stat_rect.x, stat_rect.y, stat_rect.w, stat_rect.h):
                    return ShapeRole.LAYOUT_MANAGED
            except (IndexError, AttributeError):
                break

    # Check timeline nodes
    timeline = getattr(spec, 'timeline', None)
    if timeline:
        max_items = getattr(spec, 'max_items', 0) or 8
        for idx in range(max_items):
            try:
                node_rect = timeline.node_rect(idx)
                if _rect_matches(box, node_rect.x, node_rect.y, node_rect.w, node_rect.h):
                    return ShapeRole.LAYOUT_MANAGED
            except (IndexError, AttributeError):
                break

    # Check comparison left/right
    comparison = getattr(spec, 'comparison', None)
    if comparison:
        for side_field in ('left', 'right'):
            side = getattr(comparison, side_field, None)
            if side is not None and _rect_matches(box, side.x, side.y, side.w, side.h):
                return ShapeRole.LAYOUT_MANAGED

    return None


def _max_font_size_pt(shape, fallback: float = 18.0) -> float:
    if not getattr(shape, 'has_text_frame', False):
        return fallback
    sizes: list[float] = []
    try:
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.font.size is not None:
                sizes.append(paragraph.font.size.pt)
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
    except Exception:
        return fallback
    return max(sizes) if sizes else fallback


def _dominant_font_props(shape) -> tuple[str, bool]:
    font_family = 'Calibri'
    is_bold = False
    max_font_size = 0.0
    if not getattr(shape, 'has_text_frame', False):
        return font_family, is_bold
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is None:
                    continue
                size_pt = run.font.size.pt
                if size_pt >= max_font_size:
                    max_font_size = size_pt
                    if run.font.name:
                        font_family = run.font.name
                    if run.font.bold is not None:
                        is_bold = bool(run.font.bold)
    except Exception:
        return font_family, is_bold
    return font_family, is_bold


def _estimate_required_text_height_in(shape, box: ShapeBox) -> float | None:
    text = _shape_text(shape)
    if not text:
        return None

    margin_x_emu = 0
    margin_y_emu = 0
    if getattr(shape, 'has_text_frame', False):
        try:
            tf = shape.text_frame
            margin_x_emu += int(getattr(tf, 'margin_left', 0) or 0)
            margin_x_emu += int(getattr(tf, 'margin_right', 0) or 0)
            margin_y_emu += int(getattr(tf, 'margin_top', 0) or 0)
            margin_y_emu += int(getattr(tf, 'margin_bottom', 0) or 0)
        except Exception:
            margin_x_emu = 0
            margin_y_emu = 0

    usable_width_emu = max(box.width - margin_x_emu, int(0.15 * 914400))
    width_in = usable_width_emu / 914400
    if width_in <= 0.15:
        return None

    font_pt = _max_font_size_pt(shape)
    font_family, is_bold = _dominant_font_props(shape)
    required = None
    try:
        from font_text_measure import TextMeasureRequest, measure_text_heights  # type: ignore

        req = TextMeasureRequest(
            text=text,
            width_in=width_in,
            font_family=font_family,
            font_size_pt=font_pt,
            bold=is_bold,
        )
        heights = measure_text_heights([req])
        if heights:
            required = heights[0]
    except ImportError:
        required = None

    if required is None:
        required = estimate_text_height_in(text, width_in, font_pt)

    # Add a small pad for text frame margins and bullet indentation.
    if any(line.lstrip().startswith(('-', '*', '\u2022', '\u25cf', '\u25aa', '•')) for line in text.splitlines()):
        required += 0.08

    margin_y_in = margin_y_emu / 914400
    return required + margin_y_in + 0.04


def _text_content_box(shape, outer_box: ShapeBox) -> ShapeBox | None:
    """Approximate the usable text region inside a text-bearing shape.

    This excludes text-frame margins so decorative inline icons can be allowed
    inside a panel while still being flagged if they intrude into the area
    where actual text is expected to render.
    """
    text = _shape_text(shape)
    if not text or not getattr(shape, 'has_text_frame', False):
        return None

    try:
        tf = shape.text_frame
        margin_left = int(getattr(tf, 'margin_left', 0) or 0)
        margin_right = int(getattr(tf, 'margin_right', 0) or 0)
        margin_top = int(getattr(tf, 'margin_top', 0) or 0)
        margin_bottom = int(getattr(tf, 'margin_bottom', 0) or 0)
    except Exception:
        margin_left = 0
        margin_right = 0
        margin_top = 0
        margin_bottom = 0

    content_left = outer_box.left + margin_left
    content_top = outer_box.top + margin_top
    content_width = outer_box.width - margin_left - margin_right
    content_height = outer_box.height - margin_top - margin_bottom

    if content_width <= 0 or content_height <= 0:
        return None

    return ShapeBox(
        left=content_left,
        top=content_top,
        width=content_width,
        height=content_height,
        shape_name=outer_box.shape_name,
        shape_id=outer_box.shape_id,
    )


def _is_icon_shape(shape) -> bool:
    name = (getattr(shape, 'name', '') or '').lower()
    return any(name.startswith(prefix) for prefix in _ICON_NAME_PREFIXES)


def validate_slide(
    slide,
    slide_index: int,
    *,
    include_images: bool | None = None,
    shape_role_registry: dict[int, str] | None = None,
    layout_specs: list | None = None,
) -> list[LayoutIssue]:
    """Validate a single slide for layout issues.

    Shapes are classified into two semantic roles:

    - **TEMPLATE_DESIGN** — decorative / background / template shapes.
      Checked for out-of-bounds and text overflow, but excluded from
      overlap and cramped-spacing detection.
    - **LAYOUT_MANAGED** — blueprint-driven content shapes.
      Participates in all validation checks including collision detection.

    Classification priority (first match wins):
      1. ``shape_role_registry`` — explicit role set at creation time
      2. ``layout_specs`` geometry match — shape aligns with a blueprint rect
      3. Name-prefix conventions
      4. Heuristic detection
      5. Picture + include_images flag
      6. Default → LAYOUT_MANAGED
    """
    if include_images is None:
        include_images = _should_include_images_in_layout()

    issues: list[LayoutIssue] = []

    # Collect ALL shapes with their semantic role and bounding box.
    all_boxes: list[tuple[ShapeBox, object, ShapeRole]] = []

    for shape in slide.shapes:
        box = _get_shape_box(shape)
        if box is None:
            continue
        role = _classify_shape_role(
            shape,
            include_images=include_images,
            shape_role_registry=shape_role_registry,
            layout_specs=layout_specs,
            slide_index=slide_index,
        )
        all_boxes.append((box, shape, role))

    # Only layout-managed shapes participate in collision checks.
    collision_boxes: list[tuple[ShapeBox, object]] = [
        (box, shape) for box, shape, role in all_boxes
        if role == ShapeRole.LAYOUT_MANAGED
    ]

    # 0. Text overflow / dense text check  (all shapes)
    for box, shape, _role in all_boxes:
        required_h_in = _estimate_required_text_height_in(shape, box)
        if required_h_in is None:
            continue
        available_h_in = box.height / 914400
        ratio = required_h_in / max(available_h_in, 0.01)
        if ratio > 1.12:
            overflow_sev = IssueSeverity.ERROR
            shape_name_lower = (box.shape_name or '').strip().lower()
            # Auto-sized panel shapes may still render with visible clipping, so
            # only downgrade mild overflows. Severe overflow remains blocking.
            try:
                if (getattr(shape, 'has_text_frame', False)
                        and shape.text_frame.auto_size == _TEXT_TO_FIT_SHAPE
                        and ratio <= 2.50):
                    overflow_sev = IssueSeverity.WARNING
            except Exception:
                pass
            # Flow-managed top textboxes are pre-sized upstream and can read as
            # under-height to static estimators even when the rendered result is
            # acceptable. Keep them actionable without blocking export.
            if shape_name_lower in {'title', 'key_message'}:
                overflow_sev = IssueSeverity.WARNING
            issues.append(LayoutIssue(
                slide_index=slide_index,
                issue_type=IssueType.TEXT_OVERFLOW,
                severity=overflow_sev,
                message=(
                    f'Shape "{box.shape_name}" text needs about {required_h_in:.2f}" '
                    f'but only {available_h_in:.2f}" is available'
                ),
                shape_a=box.shape_name,
            ))
        elif ratio > 1.06:
            issues.append(LayoutIssue(
                slide_index=slide_index,
                issue_type=IssueType.TEXT_OVERFLOW,
                severity=IssueSeverity.WARNING,
                message=(
                    f'Shape "{box.shape_name}" is text-dense '
                    f'({required_h_in:.2f}" needed vs {available_h_in:.2f}" available)'
                ),
                shape_a=box.shape_name,
            ))

    # 1. Out-of-bounds check  (all shapes)
    for box, _, _role in all_boxes:
        oob_parts = []
        if box.left < 0:
            oob_parts.append('extends left of slide')
        if box.top < 0:
            oob_parts.append('extends above slide')
        if box.right > SLIDE_WIDTH_EMU + OVERLAP_TOLERANCE_EMU:
            oob_parts.append('extends right of slide')
        if box.bottom > SLIDE_HEIGHT_EMU + OVERLAP_TOLERANCE_EMU:
            oob_parts.append('extends below slide')
        if oob_parts:
            # Promote right-edge overflow to ERROR so it's actionable
            severity = IssueSeverity.ERROR if 'extends right of slide' in oob_parts else IssueSeverity.WARNING
            issues.append(LayoutIssue(
                slide_index=slide_index,
                issue_type=IssueType.OUT_OF_BOUNDS,
                severity=severity,
                message=f'Shape "{box.shape_name}" {", ".join(oob_parts)}',
                shape_a=box.shape_name,
            ))

    # 1.5 Decorative icon vs text-content check.
    # Icons embedded inside text-box compositions are decorative for layout,
    # but they still must not intrude into the textbox's usable text region.
    icon_boxes = [
        (box, shape) for box, shape, _role in all_boxes
        if _is_icon_shape(shape)
    ]
    text_content_boxes = [
        (content_box, shape) for outer_box, shape, _role in all_boxes
        for content_box in [_text_content_box(shape, outer_box)]
        if content_box is not None
    ]

    for icon_box, icon_shape in icon_boxes:
        for text_box, text_shape in text_content_boxes:
            if getattr(icon_shape, 'shape_id', None) == getattr(text_shape, 'shape_id', None):
                continue
            if not _boxes_overlap(icon_box, text_box):
                continue

            overlap_area = _overlap_area(icon_box, text_box)
            if overlap_area <= 0:
                continue

            icon_area = max(icon_box.width * icon_box.height, 1)
            overlap_ratio = overlap_area / icon_area
            severity = IssueSeverity.ERROR if overlap_ratio > 0.08 else IssueSeverity.WARNING
            issues.append(LayoutIssue(
                slide_index=slide_index,
                issue_type=IssueType.OVERLAP,
                severity=severity,
                message=(
                    f'Decorative icon "{icon_box.shape_name}" overlaps the text area '
                    f'of "{text_box.shape_name}" ({overlap_ratio:.0%} of icon area)'
                ),
                shape_a=icon_box.shape_name,
                shape_b=text_box.shape_name,
            ))

    # 2. Overlap check  (layout-managed shapes only)
    for i in range(len(collision_boxes)):
        for j in range(i + 1, len(collision_boxes)):
            a, shape_a = collision_boxes[i]
            b, shape_b = collision_boxes[j]
            if _boxes_overlap(a, b):
                area = _overlap_area(a, b)
                # Only flag significant overlaps (> 5% of smaller shape area)
                smaller_area = min(a.width * a.height, b.width * b.height)
                if smaller_area > 0 and area / smaller_area > 0.05:
                    severity = IssueSeverity.ERROR if area / smaller_area > 0.25 else IssueSeverity.WARNING

                    # Downgrade to WARNING when a small non-text shape (icon badge,
                    # accent rule) overlaps a larger content shape — these are
                    # intentional decorative overlays, not layout errors.
                    if severity == IssueSeverity.ERROR:
                        a_text = _shape_text(shape_a)
                        b_text = _shape_text(shape_b)
                        a_area = a.width * a.height
                        b_area = b.width * b.height
                        a_is_small_decor = (not a_text and a_area < b_area * 0.35)
                        b_is_small_decor = (not b_text and b_area < a_area * 0.35)
                        if a_is_small_decor or b_is_small_decor:
                            severity = IssueSeverity.WARNING

                    # Fully contained shapes are intentional parent-child
                    # overlays (e.g., a caption label inside an image frame).
                    if severity == IssueSeverity.ERROR:
                        if _is_contained(a, b) or _is_contained(b, a):
                            severity = IssueSeverity.WARNING

                    # Image caption panels are an intentional overlay pattern when
                    # they sit on top of a picture with a solid background.
                    if severity == IssueSeverity.ERROR:
                        names = {a.shape_name, b.shape_name}
                        if any(name.startswith('image_caption') for name in names) and any(name.startswith('Picture') for name in names):
                            severity = IssueSeverity.WARNING

                    # Clay-morphism drop shadows are intentional decorative overlays.
                    # The shadow shape is always slightly offset from its paired main
                    # shape, so they inherently overlap. Downgrade to WARNING.
                    if severity == IssueSeverity.ERROR:
                        names = {a.shape_name, b.shape_name}
                        if any(name.startswith('clay_shadow') for name in names):
                            severity = IssueSeverity.WARNING

                    issues.append(LayoutIssue(
                        slide_index=slide_index,
                        issue_type=IssueType.OVERLAP,
                        severity=severity,
                        message=f'Shapes "{a.shape_name}" and "{b.shape_name}" overlap '
                                f'({area / smaller_area:.0%} of smaller shape)',
                        shape_a=a.shape_name,
                        shape_b=b.shape_name,
                    ))

    # 3. Cramped spacing check  (layout-managed shapes only)
    for i in range(len(collision_boxes)):
        for j in range(i + 1, len(collision_boxes)):
            a, _ = collision_boxes[i]
            b, _ = collision_boxes[j]
            if _boxes_overlap(a, b):
                continue  # already flagged as overlap
            # Horizontal gap
            h_gap = max(b.left - a.right, a.left - b.right)
            # Vertical gap
            v_gap = max(b.top - a.bottom, a.top - b.bottom)
            # Shapes are adjacent if they share a horizontal or vertical band
            h_overlap = not (a.right <= b.left or b.right <= a.left)
            v_overlap = not (a.bottom <= b.top or b.bottom <= a.top)
            if h_overlap and 0 < v_gap < MIN_GAP_EMU:
                issues.append(LayoutIssue(
                    slide_index=slide_index,
                    issue_type=IssueType.CRAMPED,
                    severity=IssueSeverity.INFO,
                    message=f'Shapes "{a.shape_name}" and "{b.shape_name}" have only '
                            f'{v_gap / 914400:.2f}" vertical gap (min {MIN_GAP_EMU / 914400:.2f}")',
                    shape_a=a.shape_name,
                    shape_b=b.shape_name,
                ))
            if v_overlap and 0 < h_gap < MIN_GAP_EMU:
                issues.append(LayoutIssue(
                    slide_index=slide_index,
                    issue_type=IssueType.CRAMPED,
                    severity=IssueSeverity.INFO,
                    message=f'Shapes "{a.shape_name}" and "{b.shape_name}" have only '
                            f'{h_gap / 914400:.2f}" horizontal gap (min {MIN_GAP_EMU / 914400:.2f}")',
                    shape_a=a.shape_name,
                    shape_b=b.shape_name,
                ))

    return issues


def validate_presentation(
    prs,
    *,
    shape_role_registry: dict[int, str] | None = None,
    layout_specs: list | None = None,
) -> list[LayoutIssue]:
    """Validate all slides in a presentation."""
    include_images = _should_include_images_in_layout()
    issues: list[LayoutIssue] = []
    for idx, slide in enumerate(prs.slides):
        issues.extend(validate_slide(
            slide,
            idx,
            include_images=include_images,
            shape_role_registry=shape_role_registry,
            layout_specs=layout_specs,
        ))
    return issues


# ---------------------------------------------------------------------------
# Auto-correction (kept: auto-size enforcement and OOB clamping only)
# Overlap resolution is intentionally NOT handled here.
# ---------------------------------------------------------------------------

def _clamp_shape_to_slide(shape) -> bool:
    """Clamp a shape so it fits within slide boundaries. Returns True if adjusted."""
    box = _get_shape_box(shape)
    if box is None:
        return False
    changed = False

    if box.left < 0:
        shape.left = int(SAFE_MARGIN_EMU)
        changed = True
    if box.top < 0:
        shape.top = int(SAFE_MARGIN_EMU)
        changed = True
    if box.right > SLIDE_WIDTH_EMU:
        new_left = int(SLIDE_WIDTH_EMU - box.width - SAFE_MARGIN_EMU)
        if new_left >= int(SAFE_MARGIN_EMU):
            shape.left = new_left
        else:
            shape.left = int(SAFE_MARGIN_EMU)
            shape.width = int(SLIDE_WIDTH_EMU - 2 * SAFE_MARGIN_EMU)
        changed = True
    if box.bottom > SLIDE_HEIGHT_EMU:
        new_top = int(SLIDE_HEIGHT_EMU - box.height - SAFE_MARGIN_EMU)
        if new_top >= int(SAFE_MARGIN_EMU):
            shape.top = new_top
        else:
            shape.top = int(SAFE_MARGIN_EMU)
            shape.height = int(SLIDE_HEIGHT_EMU - 2 * SAFE_MARGIN_EMU)
        changed = True

    return changed


def _enforce_auto_size(slide) -> int:
    """Enable TEXT_TO_FIT_SHAPE auto-size on panel/card shapes only.

    Free textboxes (used for title, key_message, notes) should keep
    MSO_AUTO_SIZE.NONE because their height is pre-computed by flow_layout_spec.
    Returns count of fixes.
    """
    fixed = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip free textboxes — they use flow-computed heights
        # Textboxes have shape_type == MSO_SHAPE_TYPE.TEXT_BOX (17)
        shape_type = getattr(shape, 'shape_type', None)
        if shape_type is not None and int(shape_type) == 17:
            # Ensure word_wrap is on, but do NOT change auto_size
            shape.text_frame.word_wrap = True
            continue
        tf = shape.text_frame
        tf.word_wrap = True
        try:
            current = tf.auto_size
        except Exception:
            current = None
        if current != _TEXT_TO_FIT_SHAPE:
            tf.auto_size = _TEXT_TO_FIT_SHAPE
            fixed += 1
    return fixed


def fix_slide(slide, issues: list[LayoutIssue]) -> int:
    """Apply auto-corrections for a single slide's issues. Returns count of fixes.

    Only handles auto-size enforcement and OOB clamping.
    Overlap resolution is handled by layout_engine.relayout_presentation().
    """
    fixed = 0

    # Fix 0: Enforce auto-size on all text frames
    fixed += _enforce_auto_size(slide)

    # Fix out-of-bounds
    for issue in issues:
        if issue.issue_type == IssueType.OUT_OF_BOUNDS:
            for shape in slide.shapes:
                name = getattr(shape, 'name', '') or ''
                if name == issue.shape_a:
                    if _clamp_shape_to_slide(shape):
                        fixed += 1
                    break

    return fixed


def fix_presentation(prs, issues: list[LayoutIssue]) -> int:
    """Apply auto-corrections across all slides. Returns total fixes applied."""
    slides_list = list(prs.slides)
    total_fixed = 0

    by_slide: dict[int, list[LayoutIssue]] = {}
    for issue in issues:
        by_slide.setdefault(issue.slide_index, []).append(issue)

    for slide_idx, slide_issues in by_slide.items():
        if slide_idx < len(slides_list):
            total_fixed += fix_slide(slides_list[slide_idx], slide_issues)

    return total_fixed


def report_issues(issues: list[LayoutIssue]) -> str:
    """Format issues into a human-readable report."""
    filtered = [i for i in issues if i.severity.value.lower() != 'info']
    
    if not filtered:
        return 'Layout validation passed — no issues found.'
    
    lines = [f'Layout validation found {len(filtered)} issue(s):']
    for issue in filtered:
        severity = issue.severity.value.lower()
        prefix = {'error': 'ERROR', 'warning': 'WARN'}[severity]
        lines.append(f'  [{prefix}] Slide {issue.slide_index + 1}: {issue.message}')
    
    return '\n'.join(lines)


